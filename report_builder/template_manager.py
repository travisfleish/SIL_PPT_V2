# report_builder/template_manager.py
"""
Template manager for PowerPoint presentations
Handles slide templates, layouts, and styling configurations
"""

import logging
from pathlib import Path
from typing import Dict, Optional, Any, List
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
import yaml
import json

logger = logging.getLogger(__name__)


class TemplateManager:
    """Manages PowerPoint templates and styling configurations"""

    def __init__(self, template_path: Optional[Path] = None):
        """
        Initialize template manager

        Args:
            template_path: Path to PowerPoint template file (.pptx)
        """
        self.template_path = template_path
        self.template_config = self._load_template_config()

        # Default style configurations
        self.styles = {
            'fonts': {
                'title': {'name': 'Arial', 'size': 28, 'bold': True},
                'subtitle': {'name': 'Arial', 'size': 20, 'bold': False},
                'heading': {'name': 'Arial', 'size': 14, 'bold': True},
                'body': {'name': 'Arial', 'size': 11, 'bold': False},
                'table_header': {'name': 'Arial', 'size': 9, 'bold': True},
                'table_body': {'name': 'Arial', 'size': 10, 'bold': False}
            },
            'colors': {
                'primary': RGBColor(0, 34, 68),  # Default blue
                'secondary': RGBColor(255, 182, 18),  # Default yellow
                'accent': RGBColor(139, 139, 139),  # Default gray
                'positive': RGBColor(0, 176, 80),  # Green
                'negative': RGBColor(255, 0, 0),  # Red
                'neutral': RGBColor(0, 0, 0),  # Black
                'background': RGBColor(255, 255, 255),  # White
                'header_bg': RGBColor(240, 240, 240),  # Light gray
                'table_header_bg': RGBColor(217, 217, 217)  # Table header gray
            },
            'spacing': {
                'margin_top': Inches(0.5),
                'margin_bottom': Inches(0.5),
                'margin_left': Inches(0.5),
                'margin_right': Inches(0.5),
                'line_spacing': 1.2
            }
        }

        # Load custom template if provided
        if self.template_path and self.template_path.exists():
            self.template_presentation = Presentation(str(self.template_path))
            logger.info(f"Loaded template from: {self.template_path}")
        else:
            self.template_presentation = None
            logger.info("No template file loaded, using default settings")

    def _load_template_config(self) -> Dict[str, Any]:
        """Load template configuration from YAML file"""
        config_path = Path(__file__).parent.parent / 'config' / 'slide_templates.yaml'

        if config_path.exists():
            with open(config_path, 'r') as f:
                config = yaml.safe_load(f)
                logger.info(f"Loaded template config from: {config_path}")
                return config
        else:
            logger.warning(f"Template config not found at: {config_path}")
            return {}

    def create_presentation(self) -> Presentation:
        """
        Create a new presentation using the template

        Returns:
            New Presentation object
        """
        if self.template_presentation:
            # Create from template
            return Presentation(str(self.template_path))
        else:
            # Create blank presentation
            return Presentation()

    def apply_team_styling(self, team_config: Dict[str, Any]):
        """
        Apply team-specific styling to the template

        Args:
            team_config: Team configuration with colors
        """
        if 'colors' in team_config:
            team_colors = team_config['colors']

            # Update primary colors with team colors
            if 'primary' in team_colors:
                self.styles['colors']['primary'] = self._hex_to_rgb(team_colors['primary'])
            if 'secondary' in team_colors:
                self.styles['colors']['secondary'] = self._hex_to_rgb(team_colors['secondary'])
            if 'accent' in team_colors:
                self.styles['colors']['accent'] = self._hex_to_rgb(team_colors['accent'])

            logger.info(f"Applied team styling for: {team_config.get('team_name', 'Unknown')}")

    def get_slide_layout(self, layout_name: str) -> Any:
        """
        Get a specific slide layout

        Args:
            layout_name: Name of the layout to retrieve

        Returns:
            Slide layout object
        """
        layout_mapping = {
            'title': 0,
            'title_and_content': 1,
            'section_header': 2,
            'two_content': 3,
            'comparison': 4,
            'blank': 5,
            'content_with_caption': 6,
            'picture_with_caption': 7
        }

        if self.template_presentation:
            layouts = self.template_presentation.slide_layouts
        else:
            # Use default presentation layouts
            temp_pres = Presentation()
            layouts = temp_pres.slide_layouts

        layout_index = layout_mapping.get(layout_name, 5)  # Default to blank

        if layout_index < len(layouts):
            return layouts[layout_index]
        else:
            logger.warning(f"Layout '{layout_name}' not found, using blank layout")
            return layouts[5]  # Blank layout

    def format_text_frame(self, text_frame, style_name: str = 'body'):
        """
        Apply formatting to a text frame

        Args:
            text_frame: PowerPoint text frame object
            style_name: Name of the style to apply
        """
        style = self.styles['fonts'].get(style_name, self.styles['fonts']['body'])

        for paragraph in text_frame.paragraphs:
            paragraph.font.name = style['name']
            paragraph.font.size = Pt(style['size'])
            paragraph.font.bold = style['bold']
            paragraph.line_spacing = self.styles['spacing']['line_spacing']

    def create_styled_table(self, slide, rows: int, cols: int,
                            left: float, top: float,
                            width: float, height: float) -> Any:
        """
        Create a styled table with consistent formatting

        Args:
            slide: Slide to add table to
            rows: Number of rows
            cols: Number of columns
            left, top, width, height: Position and size in inches

        Returns:
            Table object
        """
        table_shape = slide.shapes.add_table(
            rows, cols,
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )

        table = table_shape.table

        # Apply default styling
        for row_idx in range(rows):
            for col_idx in range(cols):
                cell = table.cell(row_idx, col_idx)

                # Header row styling
                if row_idx == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.styles['colors']['table_header_bg']
                    self.format_text_frame(cell.text_frame, 'table_header')
                else:
                    self.format_text_frame(cell.text_frame, 'table_body')

                # Set margins
                cell.margin_left = Inches(0.05)
                cell.margin_right = Inches(0.05)
                cell.margin_top = Inches(0.02)
                cell.margin_bottom = Inches(0.02)

        return table

    def add_styled_textbox(self, slide, text: str,
                           left: float, top: float,
                           width: float, height: float,
                           style_name: str = 'body') -> Any:
        """
        Add a styled text box to slide

        Args:
            slide: Slide to add text box to
            text: Text content
            left, top, width, height: Position and size in inches
            style_name: Style to apply

        Returns:
            Text box shape
        """
        text_box = slide.shapes.add_textbox(
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )

        text_frame = text_box.text_frame
        text_frame.text = text
        text_frame.word_wrap = True

        self.format_text_frame(text_frame, style_name)

        return text_box

    def add_header_footer(self, slide, team_name: str, slide_title: str):
        """
        Add consistent header/footer to slide

        Args:
            slide: Slide to add header/footer to
            team_name: Team name for header
            slide_title: Title for the slide
        """
        # Header background
        header_rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(0.4)
        )
        header_rect.fill.solid()
        header_rect.fill.fore_color.rgb = self.styles['colors']['header_bg']
        header_rect.line.color.rgb = RGBColor(200, 200, 200)
        header_rect.line.width = Pt(0.5)

        # Team name (left)
        self.add_styled_textbox(
            slide, team_name,
            0.2, 0.05, 3, 0.3,
            style_name='heading'
        )

        # Slide title (right)
        title_box = self.add_styled_textbox(
            slide, slide_title,
            5, 0.05, 4.8, 0.3,
            style_name='heading'
        )
        title_box.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    def _hex_to_rgb(self, hex_color: str) -> RGBColor:
        """Convert hex color to RGBColor"""
        hex_color = hex_color.lstrip('#')
        r = int(hex_color[0:2], 16)
        g = int(hex_color[2:4], 16)
        b = int(hex_color[4:6], 16)
        return RGBColor(r, g, b)

    def get_style(self, style_type: str, style_name: str) -> Any:
        """
        Get a specific style configuration

        Args:
            style_type: Type of style (fonts, colors, spacing)
            style_name: Name of the specific style

        Returns:
            Style configuration
        """
        return self.styles.get(style_type, {}).get(style_name)

    def export_template_config(self, output_path: Path):
        """
        Export current template configuration to file

        Args:
            output_path: Path to save configuration
        """
        config = {
            'styles': self.styles,
            'template_config': self.template_config
        }

        with open(output_path, 'w') as f:
            if output_path.suffix == '.yaml':
                yaml.dump(config, f, default_flow_style=False)
            else:
                json.dump(config, f, indent=2)

        logger.info(f"Exported template config to: {output_path}")


# Singleton instance
_template_manager = None


def get_template_manager(template_path: Optional[Path] = None) -> TemplateManager:
    """
    Get or create the template manager instance

    Args:
        template_path: Path to template file

    Returns:
        TemplateManager instance
    """
    global _template_manager

    if _template_manager is None:
        _template_manager = TemplateManager(template_path)

    return _template_manager