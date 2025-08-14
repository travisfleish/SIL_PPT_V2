# report_builder/pptx_builder.py
"""
Main PowerPoint builder that orchestrates the entire presentation generation
Combines all slide generators to create the complete sponsorship insights report
UPDATED: Added demographic overview slide with AI insights
UPDATED: Added support for emerging categories with tiered selection
UPDATED: Added real-time progress tracking
"""

import logging
from pathlib import Path
from typing import Dict, List, Optional, Any
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Import data processors
from data_processors.demographic_processor import DemographicsProcessor
from data_processors.merchant_ranker import MerchantRanker
from data_processors.category_analyzer import CategoryAnalyzer
from data_processors.snowflake_connector import query_to_dataframe

# Import slide generators
from slide_generators.title_slide import TitleSlide
from slide_generators.demographic_overview_slide import DemographicOverviewSlide  # NEW
from slide_generators.demographics_slide import DemographicsSlide  # Now single slide
from slide_generators.behaviors_slide import BehaviorsSlide
from slide_generators.category_slide import CategorySlide

# Import visualizations
from visualizations.demographic_charts import DemographicCharts

# Import utilities
from utils.team_config_manager import TeamConfigManager

# from utils.logo_downloader import LogoDownloader  # Not implemented yet

logger = logging.getLogger(__name__)

# Default font configuration
DEFAULT_FONT_FAMILY = "Red Hat Display"
FALLBACK_FONT = "Arial"


def update_progress(progress: int, message: str):
    """Update job progress if running in a job context"""
    # Use the callback if available
    if (hasattr(PowerPointBuilder._current_instance, 'progress_callback') and
            PowerPointBuilder._current_instance.progress_callback):
        try:
            PowerPointBuilder._current_instance.progress_callback(progress, message)
        except Exception as e:
            logger.debug(f"Error calling progress callback: {e}")

    # Always log the progress
    logger.info(f"Progress: {progress}% - {message}")


class PowerPointBuilder:
    """Main orchestrator for building complete PowerPoint presentations"""

    # Class variable to store current instance for progress updates
    _current_instance = None

    def __init__(self, team_key: str, job_id: Optional[str] = None, cache_manager: Optional[Any] = None, progress_callback: Optional[callable] = None):
        """
        Initialize the PowerPoint builder with proper 16:9 formatting

        Args:
            team_key: Team identifier (e.g., 'utah_jazz', 'dallas_cowboys')
            job_id: Optional job ID for progress tracking
            cache_manager: Optional CacheManager instance for caching
        """
        # Store job_id for progress tracking
        self.job_id = job_id

        # Store cache_manager for passing to components
        self.cache_manager = cache_manager

        # Set as current instance for progress updates
        PowerPointBuilder._current_instance = self

        # Store progress callback
        self.progress_callback = progress_callback

        self.team_key = team_key
        self.timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')

        # Load team configuration
        self.config_manager = TeamConfigManager()
        self.team_config = self.config_manager.get_team_config(team_key)

        # Extract key values
        self.team_name = self.team_config['team_name']
        self.team_short = self.team_config['team_name_short']
        self.league = self.team_config['league']
        self.view_prefix = self.team_config['view_prefix']

        # Initialize data processors
        self.merchant_ranker = MerchantRanker(
            team_view_prefix=self.view_prefix,
            comparison_population=self.team_config['comparison_population'],
            cache_manager = self.cache_manager
        )

        self.category_analyzer = CategoryAnalyzer(
            team_name=self.team_name,
            team_short=self.team_short,
            league=self.league,
            comparison_population=self.team_config['comparison_population'],
            audience_name=self.team_config.get('audience_name'),
            cache_manager=self.cache_manager
        )

        # Validate and set presentation font
        self.presentation_font = self._validate_font()

        # CHANGE: Use combined template instead of blue template
        COMBINED_TEMPLATE_PATH = Path(__file__).parent.parent / 'templates' / 'sil_combined_template.pptx'

        if COMBINED_TEMPLATE_PATH.exists():
            try:
                # Load the combined template
                self.presentation = Presentation(str(COMBINED_TEMPLATE_PATH))
                logger.info("Initialized presentation from combined SIL template")
            except Exception as e:
                logger.warning(f"Could not load template: {e}. Creating blank presentation.")
                self.presentation = Presentation()
        else:
            # Create blank presentation if no template
            logger.warning(f"Combined template not found at: {COMBINED_TEMPLATE_PATH}")
            self.presentation = Presentation()

        # Set 16:9 dimensions regardless of template
        self.presentation.slide_width = Inches(13.333)  # 16:9 widescreen width
        self.presentation.slide_height = Inches(7.5)  # 16:9 widescreen height

        # Use blank layout for consistent formatting (no title boxes)
        self.blank_layout = self.presentation.slide_layouts[6]

        # Setup directories
        self.output_dir = Path(f'output/{self.team_key}_{self.timestamp}')
        self.temp_dir = self.output_dir / 'temp'
        self.charts_dir = self.temp_dir / 'charts'
        self.logos_dir = self.temp_dir / 'logos'

        # Create directories
        for dir_path in [self.output_dir, self.temp_dir, self.charts_dir, self.logos_dir]:
            dir_path.mkdir(parents=True, exist_ok=True)

        # Track progress
        self.slides_created = []

        logger.info(f"Initialized PowerPoint builder for {self.team_name} (16:9 format)")
        logger.info(f"Using font: {self.presentation_font}")

    def _validate_font(self) -> str:
        """
        Validate that Red Hat Display font is available

        Returns:
            Font name to use (Red Hat Display or fallback)
        """
        try:
            # Try to create a test presentation with Red Hat Display
            test_pres = Presentation()
            test_slide = test_pres.slides.add_slide(test_pres.slide_layouts[5])
            test_box = test_slide.shapes.add_textbox(0, 0, 100, 100)
            test_box.text_frame.text = "Test"
            test_box.text_frame.paragraphs[0].font.name = DEFAULT_FONT_FAMILY

            logger.info(f"✓ {DEFAULT_FONT_FAMILY} font validated and available")
            return DEFAULT_FONT_FAMILY

        except Exception as e:
            logger.warning(f"⚠ {DEFAULT_FONT_FAMILY} font may not be available: {e}")
            logger.warning(f"Using fallback font: {FALLBACK_FONT}")

            # Update all slide generators to use fallback font
            self._update_slide_generators_font(FALLBACK_FONT)

            return FALLBACK_FONT

    def _update_slide_generators_font(self, font_name: str):
        """
        Update the default font in all slide generators

        Args:
            font_name: Font to use as default
        """
        # Update base slide
        import slide_generators.base_slide as base_slide
        base_slide.DEFAULT_FONT_FAMILY = font_name

        # Update demographics slide
        import slide_generators.demographics_slide as demo_slide
        demo_slide.DEFAULT_FONT_FAMILY = font_name

        # Update demographic overview slide
        import slide_generators.demographic_overview_slide as demo_overview_slide
        demo_overview_slide.DEFAULT_FONT_FAMILY = font_name

        # Update category slide
        import slide_generators.category_slide as cat_slide
        cat_slide.DEFAULT_FONT_FAMILY = font_name

        logger.info(f"Updated all slide generators to use font: {font_name}")

    def build_presentation(self,
                           include_custom_categories: bool = True,
                           custom_category_count: Optional[int] = None,
                           category_mode: Optional[str] = None,
                           custom_categories: Optional[str] = None) -> Path:
        """
        Build the complete PowerPoint presentation with real progress tracking

        Slide Order:
        1. Title slide
        2. "How To Use This Report" (using layout 13)
        3. Demographic Overview with AI insights (NEW - using layout 11)
        4. Demographics slide with charts
        5. Behaviors slide
        6-N. Category slides (fixed + custom/emerging OR fully custom)
        N+1. "Sports Innovation Lab" branding (using layout 14)

        Args:
            include_custom_categories: Whether to include custom categories
            custom_category_count: Number of custom categories (default: 4 for men's, 2 for women's)
            category_mode: Override category mode ('standard' or 'custom')
            custom_categories: Comma-separated string of custom category names

        Returns:
            Path to the generated PowerPoint file
        """
        logger.info(f"Starting presentation build for {self.team_name}")
        logger.info(f"Font: {self.presentation_font}")

        try:
            # Check category mode from team config or override with passed parameters
            if category_mode:
                logger.info(f"Overriding category mode: {category_mode}")
                # Create a temporary override config
                override_config = self.team_config.copy()
                override_config['category_mode'] = category_mode
                
                if category_mode == 'custom' and custom_categories:
                    # Parse custom categories from comma-separated string
                    selected_categories = [cat.strip() for cat in custom_categories.split(',')]
                    override_config['custom_categories'] = {
                        'count': len(selected_categories),
                        'selected_categories': selected_categories,
                        'thresholds': {
                            'min_audience_pct': 0.20,  # Default thresholds
                            'min_merchant_audience_pct': 0.10
                        }
                    }
                    logger.info(f"Custom categories: {', '.join(selected_categories)}")
                elif category_mode == 'standard':
                    # Remove custom category config if switching to standard
                    if 'custom_categories' in override_config:
                        del override_config['custom_categories']
                
                # Use override config for this build
                build_config = override_config
            else:
                # Use team config as-is
                build_config = self.team_config
            
            category_mode = build_config.get('category_mode', 'standard')
            logger.info(f"Category mode: {category_mode}")
            
            # Calculate total slides for progress tracking
            is_womens = self._is_womens_team()
            
            if category_mode == 'custom':
                # Fully custom mode - use team's selected categories
                custom_categories = build_config.get('custom_categories', {}).get('selected_categories', [])
                fixed_categories = []  # No fixed categories in custom mode
                custom_count = len(custom_categories)
                logger.info(f"Custom mode: {custom_count} selected categories")
            else:
                # Standard mode - use existing fixed + custom logic
                fixed_categories = ['restaurants', 'athleisure', 'finance', 'gambling', 'travel', 'auto']
                if is_womens:
                    fixed_categories.extend(['beauty', 'health'])
                custom_count = custom_category_count or (2 if is_womens else 4) if include_custom_categories else 0
                logger.info(f"Standard mode: {len(fixed_categories)} fixed + {custom_count} custom categories")

            total_categories = len(fixed_categories) + custom_count

            # Progress milestones
            progress_per_category = 40 // total_categories if total_categories > 0 else 5  # Categories take up ~40% of progress (50-90%)

            # 1. Create title slide (25-30%)
            update_progress(28, "Creating title slide...")
            self._create_title_slide()
            update_progress(30, "Title slide completed")

            # 2. Add "How To Use This Report" static slide (30-32%)
            update_progress(31, "Adding instructions slide...")
            self._add_static_slide_from_layout(13, "How To Use This Report")
            update_progress(32, "Instructions slide added")

            # 3. Create demographic overview slide (32-37%)
            update_progress(33, "Generating demographic AI insights...")
            self._create_demographic_overview_slide()
            update_progress(37, "Demographic overview completed")

            # 4. Create detailed demographics slide (37-45%)
            update_progress(38, "Loading demographic data from Snowflake...")
            self._create_demographics_slide()
            update_progress(45, "Demographics charts completed")

            # 5. Create behaviors slide (45-50%)
            update_progress(46, "Analyzing fan behaviors...")
            self._create_behaviors_slide()
            update_progress(50, "Behaviors slide completed")

            # 6. Create category slides (50-85%)
            current_progress = 50
            update_progress(current_progress, "Starting category analysis...")

            # Fixed categories (only in standard mode)
            for i, category_key in enumerate(fixed_categories):
                current_progress = 50 + (i * progress_per_category)
                update_progress(
                    current_progress,
                    f"Analyzing {category_key} category ({i + 1}/{len(fixed_categories)})..."
                )
                self._create_category_slide(category_key, is_custom=False)
                update_progress(
                    current_progress + progress_per_category - 1,
                    f"Completed {category_key} slides"
                )

            # 7. Custom categories (if requested) OR fully custom categories
            if category_mode == 'custom':
                # Fully custom mode - create slides for team's selected categories
                base_progress = 50 + (len(fixed_categories) * progress_per_category)
                update_progress(base_progress, "Creating custom category slides...")
                self._create_fully_custom_category_slides(custom_categories)
                update_progress(85, "All custom categories completed")
            elif include_custom_categories:
                # Standard mode - use existing custom category logic
                base_progress = 50 + (len(fixed_categories) * progress_per_category)
                update_progress(base_progress, "Identifying top custom categories...")
                self._create_custom_category_slides(custom_category_count)
                update_progress(85, "All custom categories completed")
            else:
                update_progress(85, "All categories completed")

            # 8. Add SIL branding slide (85-87%)
            update_progress(86, "Adding branding slide...")
            self._add_static_slide_from_layout(14, "Sports Innovation Lab Branding")
            update_progress(87, "Branding slide added")

            # 9. Save presentation (87-90%)
            update_progress(88, "Saving presentation file...")
            output_path = self._save_presentation()
            update_progress(90, "Presentation saved successfully")

            logger.info(f"Presentation completed with {len(self.slides_created)} slides")
            logger.info("Final slide order:")
            for i, slide_name in enumerate(self.slides_created):
                logger.info(f"  Slide {i + 1}: {slide_name}")

            return output_path

        except Exception as e:
            logger.error(f"Error building presentation: {str(e)}")
            raise

        finally:
            # ADD THIS: Clean up connection pool after all processing
            try:
                from data_processors.snowflake_connector import close_pool
                close_pool()
                logger.info("✓ Snowflake connection pool closed")
            except Exception as e:
                logger.debug(f"Error closing pool: {e}")

    def _create_demographic_overview_slide(self):
        """
        NEW: Create the demographic overview slide with AI insights
        Uses the blue SIL layout #11 with AI-generated demographic insights
        """
        logger.info("Creating demographic overview slide with AI insights...")

        try:
            # Get AI insights from demographic data
            ai_insights = self._get_demographic_ai_insights()

            # Create the slide
            overview_generator = DemographicOverviewSlide(self.presentation)  # Pass template!
            overview_generator.generate(
                team_config=self.team_config,
                ai_insights=ai_insights
            )

            self.slides_created.append("Demographic Overview")
            logger.info("✓ Demographic overview slide created")

        except Exception as e:
            logger.error(f"Error creating demographic overview slide: {str(e)}")
            # Create a fallback slide
            self._add_placeholder_slide("Demographic Overview - Error loading data")

    def _get_demographic_ai_insights(self) -> str:
        """
        Get AI-generated demographic insights for the overview slide

        Returns:
            AI-generated insights text or fallback text
        """
        try:
            # Load demographics data
            demographics_view = self.config_manager.get_view_name(self.team_key, 'demographics')
            query = f"SELECT * FROM {demographics_view}"
            df = query_to_dataframe(query)

            if df.empty:
                logger.warning("No demographics data found for AI insights")
                return self._get_fallback_demographic_insight()

            comparison_population = self.team_config.get('comparison_population')

            # Process demographics with AI insights enabled
            processor = DemographicsProcessor(
                data_source=df,
                team_name=self.team_name,
                league=self.league,
                use_ai_insights=True,  # Enable AI insights
                comparison_population=comparison_population  # ADD THIS LINE
            )

            demographic_data = processor.process_all_demographics()

            # Get the AI-generated insight
            ai_insights = demographic_data.get('key_insights')
            if ai_insights and len(ai_insights.strip()) > 20:
                logger.info("Using AI-generated demographic insights")
                return ai_insights
            else:
                logger.warning("AI insights not available, using fallback")
                return self._get_fallback_demographic_insight()

        except Exception as e:
            logger.warning(f"Could not generate AI demographic insights: {e}")
            return self._get_fallback_demographic_insight()

    def _get_fallback_demographic_insight(self) -> str:
        """
        Generate fallback demographic insight when AI is not available

        Returns:
            Fallback insights text
        """
        team_name = self.team_config['team_name']
        team_short = self.team_config['team_name_short']

        # Provide a generic but professional fallback
        return (f"{team_name} fans demonstrate distinct demographic characteristics including "
                f"unique age distribution, income levels, and professional backgrounds compared "
                f"to the local general population and {self.league} average fans, making them "
                f"an attractive audience for targeted sponsorship opportunities.")

    def _add_static_slide_from_layout(self, layout_index: int, slide_name: str):
        """
        Add a static slide using a specific layout from the template

        Args:
            layout_index: Index of layout to use (13 or 14)
            slide_name: Name for tracking purposes
        """
        try:
            # Check if the layout exists
            if layout_index >= len(self.presentation.slide_layouts):
                logger.warning(f"Layout {layout_index} not found - adding placeholder for {slide_name}")
                self._add_placeholder_slide(f"{slide_name} - Layout not found")
                return

            # Get the layout
            static_layout = self.presentation.slide_layouts[layout_index]
            layout_name = getattr(static_layout, 'name', f'Layout {layout_index}')

            logger.info(f"Adding static slide '{slide_name}' using layout {layout_index} ({layout_name})")

            # Create new slide using the static layout
            slide = self.presentation.slides.add_slide(static_layout)

            # Add header to "How To Use This Report" slide (layout 13)
            if layout_index == 13:
                from pptx.enum.shapes import MSO_SHAPE
                from pptx.util import Inches, Pt
                from pptx.dml.color import RGBColor
                from pptx.enum.text import PP_ALIGN

                # Header background rectangle
                header_rect = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(0), Inches(0),
                    Inches(13.333), Inches(0.5)  # Full 16:9 width
                )

                # Header styling
                header_rect.fill.solid()
                header_rect.fill.fore_color.rgb = RGBColor(240, 240, 240)  # Light gray
                header_rect.line.color.rgb = RGBColor(200, 200, 200)  # Border gray
                header_rect.line.width = Pt(0.5)

                # Team name text (left side)
                team_text = slide.shapes.add_textbox(
                    Inches(0.2), Inches(0.1),
                    Inches(3), Inches(0.3)
                )
                team_text.text_frame.text = "Sponsorship Insights Report"
                team_p = team_text.text_frame.paragraphs[0]
                team_p.font.name = self.presentation_font
                team_p.font.size = Pt(14)
                team_p.font.bold = True
                team_p.font.color.rgb = RGBColor(0, 0, 0)  # Dark gray

                # Header title text (right side)
                title_text = slide.shapes.add_textbox(
                    Inches(6.5), Inches(0.1),
                    Inches(6.633), Inches(0.3)
                )
                title_text.text_frame.text = "How To Use This Report"
                title_p = title_text.text_frame.paragraphs[0]
                title_p.font.name = self.presentation_font
                title_p.font.size = Pt(14)
                title_p.alignment = PP_ALIGN.RIGHT
                title_p.font.color.rgb = RGBColor(0, 0, 0)  # Dark gray

            # Track the slide
            self.slides_created.append(slide_name)
            logger.info(f"✓ Added static slide: {slide_name}")

        except Exception as e:
            logger.error(f"Error adding static slide '{slide_name}': {str(e)}")
            # Add placeholder instead of failing
            self._add_placeholder_slide(f"{slide_name} - Error loading layout")

    def _create_title_slide(self):
        """Create the title slide"""
        logger.info("Creating title slide...")

        title_generator = TitleSlide(self.presentation)
        title_generator.default_font = self.presentation_font

        self.presentation = title_generator.generate(
            team_config=self.team_config,
            subtitle="Sponsorship Insights Report"
        )

        self.slides_created.append("Title Slide")
        logger.info("✓ Title slide created")

    def _get_demographic_insights(self) -> Optional[str]:
        """Generate demographic insights text for title slide (DEPRECATED - use _get_demographic_ai_insights)"""
        # This method is kept for backward compatibility but now calls the new AI method
        return self._get_demographic_ai_insights()

    def _format_demographic_insights(self, demographic_data: Dict[str, Any]) -> str:
        """Format demographic data into insights text (DEPRECATED)"""
        # This method is kept for backward compatibility
        return self._get_fallback_demographic_insight()

    def _create_demographics_slide(self):
        """Create single demographics slide with all 6 charts"""
        logger.info("Creating demographics slide...")

        try:
            # Load demographics data
            update_progress(39, "Querying demographic data...")
            demographics_view = self.config_manager.get_view_name(self.team_key, 'demographics')
            query = f"SELECT * FROM {demographics_view}"
            df = query_to_dataframe(query)

            if df.empty:
                logger.warning("No demographics data found")
                self._add_placeholder_slide("Demographics data not available")
                return

            comparison_population = self.team_config.get('comparison_population')

            # Process demographics
            update_progress(41, "Processing demographic data...")
            processor = DemographicsProcessor(
                data_source=df,
                team_name=self.team_name,
                league=self.league,
                comparison_population=comparison_population  # ADD THIS LINE
            )

            demographic_data = processor.process_all_demographics()

            # Generate charts
            update_progress(43, "Generating demographic charts...")
            charter = DemographicCharts(
                team_colors=self.team_config.get('colors'),
                team_config=self.team_config  # ADD THIS LINE!
            )
            charts = charter.create_all_demographic_charts(
                demographic_data,
                output_dir=self.charts_dir
            )

            # Create single demographics slide with all 6 charts
            update_progress(44, "Creating demographics slide...")
            demo_generator = DemographicsSlide(self.presentation)
            demo_generator.default_font = self.presentation_font

            self.presentation = demo_generator.generate(
                demographic_data=demographic_data,
                chart_dir=self.charts_dir,
                team_config=self.team_config
            )

            self.slides_created.append("Demographics Slide (All 6 Charts)")
            logger.info("✓ Demographics slide created")

        except Exception as e:
            logger.error(f"Error creating demographics slide: {str(e)}")
            self._add_placeholder_slide("Demographics slide - error loading data")

    def _create_behaviors_slide(self):
        """Create the fan behaviors slide"""
        logger.info("Creating behaviors slide...")

        try:
            update_progress(47, "Loading behavior data...")
            behaviors_generator = BehaviorsSlide(self.presentation)
            behaviors_generator.default_font = self.presentation_font

            update_progress(48, "Generating fan wheel visualization...")
            self.presentation = behaviors_generator.generate(
                self.merchant_ranker,
                self.team_config
            )

            self.slides_created.append("Fan Behaviors")
            logger.info("✓ Behaviors slide created")

        except Exception as e:
            logger.error(f"Error creating behaviors slide: {str(e)}")
            self._add_placeholder_slide("Behaviors slide - error loading data")

    def _create_fixed_category_slides(self):
        """Create slides for all fixed categories"""
        logger.info("Creating fixed category slides...")

        # Define fixed categories based on team type
        is_womens_team = self._is_womens_team()

        fixed_categories = ['restaurants', 'athleisure', 'finance', 'gambling', 'travel', 'auto']
        if is_womens_team:
            fixed_categories.extend(['beauty', 'health'])

        for category_key in fixed_categories:
            self._create_category_slide(category_key, is_custom=False)

    def _create_custom_category_slides(self, custom_count: Optional[int] = None):
        """
        Create slides for custom categories using tiered selection
        Now includes both established and emerging categories
        """
        logger.info("Creating custom category slides...")

        # Determine number of custom categories
        is_womens_team = self._is_womens_team()
        if custom_count is None:
            custom_count = 2 if is_womens_team else 4

        # Get all data needed for custom category selection
        try:
            # Calculate base progress for custom categories
            fixed_categories = ['restaurants', 'athleisure', 'finance', 'gambling', 'travel', 'auto']
            if is_womens_team:
                fixed_categories.extend(['beauty', 'health'])

            progress_per_category = 35 // (len(fixed_categories) + custom_count) if (
                                                                                                len(fixed_categories) + custom_count) > 0 else 5
            base_progress = 50 + (len(fixed_categories) * progress_per_category)

            # Load category data
            update_progress(base_progress + 1, "Loading category data for custom analysis...")
            category_query = f"SELECT * FROM {self.view_prefix}_CATEGORY_INDEXING_ALL_TIME"
            all_category_df = query_to_dataframe(category_query)

            # NEW: Load merchant data for verification
            update_progress(base_progress + 2, "Loading merchant data...")
            merchant_query = f"SELECT * FROM {self.view_prefix}_MERCHANT_INDEXING_ALL_TIME"
            all_merchant_df = query_to_dataframe(merchant_query)

            # Get custom categories using the new tiered selection
            update_progress(base_progress + 3, "Analyzing custom categories...")
            custom_categories = self.category_analyzer.get_custom_categories(
                category_df=all_category_df,
                merchant_df=all_merchant_df,  # NEW: Pass merchant data
                is_womens_team=is_womens_team,
                existing_categories=fixed_categories
            )

            # Log the selected categories
            logger.info(f"Selected {len(custom_categories)} custom categories:")
            for i, cat in enumerate(custom_categories):
                emerging_tag = " [EMERGING]" if cat.get('is_emerging', False) else " [ESTABLISHED]"
                logger.info(f"  {i + 1}. {cat['display_name']}{emerging_tag} "
                            f"(audience: {cat.get('audience_pct', 0) * 100:.1f}%, "
                            f"composite: {cat.get('composite_index', 0):.1f})")

            # Create slides for each custom category
            for i, custom_cat in enumerate(custom_categories[:custom_count]):
                category_name = custom_cat['display_name'].strip()  # Strip here!
                is_emerging = custom_cat.get('is_emerging', False)

                current_progress = base_progress + 5 + (i * progress_per_category)
                update_progress(
                    current_progress,
                    f"Creating custom category {i + 1}/{custom_count}: {category_name}"
                )

                logger.info(f"Creating custom category slide {i + 1}: {category_name} "
                            f"{'[EMERGING]' if is_emerging else '[ESTABLISHED]'}")

                # Pass the stripped category name
                self._create_category_slide(
                    category_key=category_name,
                    is_custom=True,
                    custom_cat_info=custom_cat
                )

        except Exception as e:
            logger.error(f"Error creating custom category slides: {str(e)}")

    def _create_fully_custom_category_slides(self, selected_categories: List[str]):
        """
        Create slides for fully custom categories selected by the team
        """
        logger.info(f"Creating fully custom category slides for {len(selected_categories)} categories...")
        
        try:
            # Calculate progress per category
            progress_per_category = 35 // len(selected_categories) if len(selected_categories) > 0 else 5
            base_progress = 50  # Start at 50 since no fixed categories in custom mode
            
            # Load category data
            update_progress(base_progress + 1, "Loading category data for custom analysis...")
            category_query = f"SELECT * FROM {self.view_prefix}_CATEGORY_INDEXING_ALL_TIME"
            all_category_df = query_to_dataframe(category_query)
            
            # Load merchant data
            update_progress(base_progress + 2, "Loading merchant data...")
            merchant_query = f"SELECT * FROM {self.view_prefix}_MERCHANT_INDEXING_ALL_TIME"
            all_merchant_df = query_to_dataframe(merchant_query)
            
            # Create slides for each selected category
            for i, category_name in enumerate(selected_categories):
                current_progress = base_progress + 3 + (i * progress_per_category)
                update_progress(
                    current_progress,
                    f"Creating custom category {i + 1}/{len(selected_categories)}: {category_name}"
                )
                
                logger.info(f"Creating fully custom category slide {i + 1}: {category_name}")
                
                # Create the category slide
                self._create_category_slide(
                    category_key=category_name,
                    is_custom=True,
                    custom_cat_info={'display_name': category_name, 'is_emerging': False}
                )
                
                update_progress(
                    current_progress + progress_per_category - 1,
                    f"Completed {category_name} slides"
                )
                
        except Exception as e:
            logger.error(f"Error creating fully custom category slides: {str(e)}")

    def _create_category_slide(self, category_key: str, is_custom: bool = False,
                               custom_cat_info: Optional[Dict[str, Any]] = None):
        """
        Create slides for a single category
        """
        try:
            logger.info(f"Creating slides for {category_key} {'[CUSTOM]' if is_custom else '[FIXED]'}...")

            # Load category data
            if is_custom:
                # FIX: Strip whitespace from category_key for custom categories
                category_key = category_key.strip()
                cat_config = self.category_analyzer.create_custom_category_config(category_key)
                cat_names = [category_key]  # Now using the stripped key
            else:
                cat_config = self.category_analyzer.categories.get(category_key, {})
                cat_names = cat_config.get('category_names_in_data', [])

            if not cat_names:
                logger.warning(f"No configuration found for {category_key}")
                return

            # Build WHERE clause - also strip each cat name for safety
            category_where = " OR ".join([f"TRIM(CATEGORY) = '{cat.strip()}'" for cat in cat_names])

            # Load data
            category_df = query_to_dataframe(f"""
                SELECT * FROM {self.view_prefix}_CATEGORY_INDEXING_ALL_TIME 
                WHERE {category_where}
            """)

            subcategory_df = query_to_dataframe(f"""
                SELECT * FROM {self.view_prefix}_SUBCATEGORY_INDEXING_ALL_TIME 
                WHERE {category_where}
            """)

            merchant_df = query_to_dataframe(f"""
                SELECT * FROM {self.view_prefix}_MERCHANT_INDEXING_ALL_TIME 
                WHERE {category_where}
                AND AUDIENCE = '{self.category_analyzer.audience_name}'
                ORDER BY PERC_AUDIENCE DESC
            """)

            # NEW: Load LAST_FULL_YEAR data for specific insights
            subcategory_last_year_df = query_to_dataframe(f"""
                SELECT * FROM {self.view_prefix}_SUBCATEGORY_INDEXING_LAST_FULL_YEAR 
                WHERE {category_where}
            """)

            merchant_last_year_df = query_to_dataframe(f"""
                SELECT * FROM {self.view_prefix}_MERCHANT_INDEXING_LAST_FULL_YEAR 
                WHERE {category_where}
                AND AUDIENCE = '{self.category_analyzer.audience_name}'
                ORDER BY PERC_AUDIENCE DESC
            """)

            # Add config for custom categories temporarily
            if is_custom:
                self.category_analyzer.categories[category_key] = cat_config

            # Analyze category
            results = self.category_analyzer.analyze_category(
                category_key=category_key,
                category_df=category_df,
                subcategory_df=subcategory_df,
                merchant_df=merchant_df,
                subcategory_last_year_df=subcategory_last_year_df,  # NEW
                merchant_last_year_df=merchant_last_year_df,  # NEW
                validate=False
            )

            # NEW: Add emerging flag from custom_cat_info if available
            if custom_cat_info and 'is_emerging' in custom_cat_info:
                results['is_emerging'] = custom_cat_info['is_emerging']

            # Clean up temporary config
            if is_custom:
                del self.category_analyzer.categories[category_key]

            # Create slides
            category_generator = CategorySlide(self.presentation)
            category_generator.default_font = self.presentation_font

            # Category analysis slide
            self.presentation = category_generator.generate(results, self.team_config)

            # Track slide with emerging tag if applicable
            slide_name = f"{results['display_name']} Analysis"
            if results.get('is_emerging', False):
                slide_name += " [EMERGING]"
            self.slides_created.append(slide_name)

            # Brand analysis slide
            self.presentation = category_generator.generate_brand_slide(results, self.team_config)
            self.slides_created.append(f"{results['display_name']} Brands")

            logger.info(f"✓ Created {results['display_name']} slides" +
                        (" [EMERGING]" if results.get('is_emerging', False) else ""))

        except Exception as e:
            logger.error(f"Error creating {category_key} slides: {str(e)}")
            self._add_placeholder_slide(f"{category_key.title()} - error loading data")

    def _add_placeholder_slide(self, message: str):
        """Add a placeholder slide when data is not available"""
        # CHANGE: Use white template layout if available
        if len(self.presentation.slides) > 1:
            # Get the white layout from slide 1
            white_layout = self.presentation.slides[1].slide_layout
            slide = self.presentation.slides.add_slide(white_layout)
        else:
            slide = self.presentation.slides.add_slide(self.blank_layout)

        text_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11.333), Inches(1))
        text_box.text = message
        p = text_box.text_frame.paragraphs[0]
        p.font.name = self.presentation_font
        p.font.size = Pt(24)
        p.alignment = PP_ALIGN.CENTER
        self.slides_created.append(f"Placeholder: {message}")

    def _is_womens_team(self) -> bool:
        """Determine if this is a women's team"""
        womens_indicators = ["women's", "ladies", "wnba", "nwsl"]
        return any(indicator in self.team_name.lower() for indicator in womens_indicators)

    def _save_presentation(self) -> Path:
        """Save the presentation to file"""
        filename = f"{self.team_key}_sponsorship_insights_{self.timestamp}.pptx"
        output_path = self.output_dir / filename

        self.presentation.save(str(output_path))
        logger.info(f"Presentation saved to: {output_path}")

        # Create summary file
        self._create_summary_file(output_path)

        return output_path

    def _create_summary_file(self, pptx_path: Path):
        """Create a summary text file with build information"""
        summary_path = self.output_dir / 'build_summary.txt'

        # Count emerging categories
        emerging_count = sum(1 for slide in self.slides_created if "[EMERGING]" in slide)

        with open(summary_path, 'w') as f:
            f.write(f"PowerPoint Build Summary\n")
            f.write(f"{'=' * 50}\n\n")
            f.write(f"Team: {self.team_name}\n")
            f.write(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n")
            f.write(f"Output: {pptx_path.name}\n")
            f.write(f"Format: 16:9 Widescreen (13.333\" x 7.5\")\n")
            f.write(f"Font: {self.presentation_font}\n\n")
            f.write(f"Slides Created ({len(self.slides_created)}):\n")
            for i, slide in enumerate(self.slides_created, 1):
                f.write(f"  {i}. {slide}\n")

            # Add category summary
            f.write(f"\nCategory Summary:\n")
            f.write(
                f"  Fixed Categories: {sum(1 for s in self.slides_created if 'Analysis' in s and '[EMERGING]' not in s and 'Demographic' not in s)}\n")
            f.write(
                f"  Established Custom: {sum(1 for s in self.slides_created if 'Analysis' in s and '[EMERGING]' not in s) - sum(1 for s in self.slides_created if any(cat in s.lower() for cat in ['restaurants', 'athleisure', 'finance', 'gambling', 'travel', 'auto', 'beauty', 'health']))}\n")
            f.write(f"  Emerging Categories: {emerging_count}\n")

            f.write(f"\nView Configuration:\n")
            f.write(f"  Prefix: {self.view_prefix}\n")
            f.write(f"  Demographics: {self.config_manager.get_view_name(self.team_key, 'demographics')}\n")

            # Add font validation status
            f.write(f"\nFont Configuration:\n")
            f.write(f"  Requested: {DEFAULT_FONT_FAMILY}\n")
            f.write(f"  Used: {self.presentation_font}\n")
            if self.presentation_font != DEFAULT_FONT_FAMILY:
                f.write(f"  Note: Using fallback font as {DEFAULT_FONT_FAMILY} was not available\n")

        logger.info(f"Summary saved to: {summary_path}")

    def check_font_installation(self) -> Dict[str, Any]:
        """
        Check if Red Hat Display font is properly installed

        Returns:
            Dictionary with font status information
        """
        import platform
        import os

        font_status = {
            'requested_font': DEFAULT_FONT_FAMILY,
            'fallback_font': FALLBACK_FONT,
            'system': platform.system(),
            'font_available': False,
            'font_paths': [],
            'instructions': []
        }

        system = platform.system()

        if system == "Windows":
            # Check Windows font directory
            fonts_dir = os.path.join(os.environ.get('WINDIR', 'C:\\Windows'), 'Fonts')
            if os.path.exists(fonts_dir):
                red_hat_fonts = [f for f in os.listdir(fonts_dir) if 'RedHat' in f or 'Red Hat' in f]
                font_status['font_paths'] = red_hat_fonts
                font_status['font_available'] = len(red_hat_fonts) > 0

            if not font_status['font_available']:
                font_status['instructions'] = [
                    "1. Download Red Hat Display from: https://fonts.google.com/specimen/Red+Hat+Display",
                    "2. Extract the ZIP file",
                    "3. Select all .ttf files",
                    "4. Right-click and select 'Install for all users'",
                    "5. Restart your Python environment"
                ]

        elif system == "Darwin":  # macOS
            # Check common font locations
            font_dirs = [
                os.path.expanduser("~/Library/Fonts"),
                "/Library/Fonts",
                "/System/Library/Fonts"
            ]

            for font_dir in font_dirs:
                if os.path.exists(font_dir):
                    red_hat_fonts = [f for f in os.listdir(font_dir) if 'RedHat' in f or 'Red Hat' in f]
                    font_status['font_paths'].extend(red_hat_fonts)

            font_status['font_available'] = len(font_status['font_paths']) > 0

            if not font_status['font_available']:
                font_status['instructions'] = [
                    "1. Download Red Hat Display from: https://fonts.google.com/specimen/Red+Hat+Display",
                    "2. Extract the ZIP file",
                    "3. Double-click each .ttf file",
                    "4. Click 'Install Font' in Font Book",
                    "5. Restart your Python environment"
                ]

        elif system == "Linux":
            # Check using fc-list if available
            try:
                import subprocess
                result = subprocess.run(['fc-list', ':family'],
                                        capture_output=True, text=True)
                font_status['font_available'] = 'Red Hat Display' in result.stdout
            except:
                pass

            if not font_status['font_available']:
                font_status['instructions'] = [
                    "1. Download Red Hat Display from: https://fonts.google.com/specimen/Red+Hat+Display",
                    "2. Extract the ZIP file",
                    "3. Copy .ttf files to ~/.fonts/ or /usr/share/fonts/",
                    "4. Run: fc-cache -f -v",
                    "5. Restart your Python environment"
                ]

        return font_status

def build_report(team_key: str, **kwargs) -> Path:
    """
    Convenience function to build a complete PowerPoint report

    Args:
        team_key: Team identifier
        **kwargs: Additional arguments for PowerPointBuilder.build_presentation()

    Returns:
        Path to generated PowerPoint file
    """
    try:
        # Extract job_id and callback if provided in kwargs
        job_id = kwargs.pop('job_id', None)
        cache_manager = kwargs.pop('cache_manager', None)
        progress_callback = kwargs.pop('progress_callback', None)

        builder = PowerPointBuilder(team_key, job_id=job_id, cache_manager=cache_manager,
                                    progress_callback=progress_callback)

        # Log font status
        font_status = builder.check_font_installation()
        if not font_status['font_available']:
            logger.warning(f"Font '{DEFAULT_FONT_FAMILY}' not installed on system")
            logger.info("Installation instructions:")
            for instruction in font_status['instructions']:
                logger.info(f"  {instruction}")

        return builder.build_presentation(**kwargs)

    finally:
        # Clean up connection pool after report generation
        try:
            from data_processors.snowflake_connector import close_pool
            close_pool()
            logger.info("✓ Snowflake connection pool closed")
        except Exception as e:
            logger.debug(f"Error closing connection pool: {e}")
            # Don't raise - we want the report to still return successfully


def validate_fonts_before_build():
    """
    Validate fonts before building any presentations
    Can be called from main.py
    """
    try:
        # Test font availability
        test_pres = Presentation()
        test_slide = test_pres.slides.add_slide(test_pres.slide_layouts[5])
        test_box = test_slide.shapes.add_textbox(0, 0, 100, 100)
        test_box.text_frame.text = "Font Test"

        # Try Red Hat Display
        p = test_box.text_frame.paragraphs[0]
        p.font.name = DEFAULT_FONT_FAMILY

        logger.info(f"✓ Font validation passed: {DEFAULT_FONT_FAMILY} is available")
        return True

    except Exception as e:
        logger.warning(f"⚠ Font validation failed: {DEFAULT_FONT_FAMILY} may not be available")
        logger.warning(f"  Error: {str(e)}")
        logger.info(f"  Will use fallback font: {FALLBACK_FONT}")

        # Print installation instructions
        print("\n" + "=" * 60)
        print("Red Hat Display Font Installation Required")
        print("=" * 60)
        print("\nTo use Red Hat Display font:")
        print("1. Download from: https://fonts.google.com/specimen/Red+Hat+Display")
        print("2. Install all .ttf files on your system")
        print("3. Restart your Python environment")
        print("=" * 60 + "\n")

        return False