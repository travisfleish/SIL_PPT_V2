"""
Centralized font configuration for PowerPoint presentations
"""

from pptx.util import Pt
from typing import Optional

# Default font family for all text
DEFAULT_FONT_FAMILY = "Red Hat Display"

# Fallback fonts if Red Hat Display is not available
FALLBACK_FONTS = ["Arial", "Helvetica", "Sans Serif"]


def apply_font_to_text_frame(text_frame, font_name: Optional[str] = None,
                             font_size: Optional[int] = None, bold: bool = False):
    """
    Apply font styling to a text frame

    Args:
        text_frame: PowerPoint text frame object
        font_name: Font family name (defaults to DEFAULT_FONT_FAMILY)
        font_size: Font size in points
        bold: Whether text should be bold
    """
    if font_name is None:
        font_name = DEFAULT_FONT_FAMILY

    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = font_name
            if font_size:
                run.font.size = Pt(font_size)
            if bold is not None:
                run.font.bold = bold


def apply_font_to_paragraph(paragraph, font_name: Optional[str] = None,
                            font_size: Optional[int] = None, bold: bool = False):
    """
    Apply font styling to a single paragraph
    """
    if font_name is None:
        font_name = DEFAULT_FONT_FAMILY

    paragraph.font.name = font_name
    if font_size:
        paragraph.font.size = Pt(font_size)
    if bold is not None:
        paragraph.font.bold = bold