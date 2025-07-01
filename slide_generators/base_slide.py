# slide_generators/base_slide.py
"""
Base class for all slide generators
Provides common functionality for creating PowerPoint slides
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path
import logging
from typing import Optional, Tuple

logger = logging.getLogger(__name__)

# Default font configuration
DEFAULT_FONT_FAMILY = "Red Hat Display"
FALLBACK_FONT = "Arial"


class BaseSlide:
    """Base class for all slide generators"""

    def __init__(self, presentation: Optional[Presentation] = None):
        """
        Initialize base slide generator with proper 16:9 formatting

        Args:
            presentation: Existing presentation to add slides to
        """
        if presentation is None:
            self.presentation = Presentation()
            # Set 16:9 widescreen aspect ratio
            self.presentation.slide_width = Inches(13.333)  # 16:9 width
            self.presentation.slide_height = Inches(7.5)  # 16:9 height
        else:
            self.presentation = presentation

        # Set default font
        self.default_font = DEFAULT_FONT_FAMILY

        # IMPORTANT: Setup layouts including SIL custom layouts
        self._setup_layouts()

    def _setup_layouts(self):
        """Setup slide layout references based on template structure"""
        layout_count = len(self.presentation.slide_layouts)
        logger.info(f"Available slide layouts: {layout_count}")

        # Log all available layouts for debugging
        for i, layout in enumerate(self.presentation.slide_layouts):
            logger.debug(f"Layout {i}: {layout.name}")

        # Always keep reference to standard blank layout
        self.blank_layout = self.presentation.slide_layouts[6]  # Blank layout

        # SIL Template has custom layouts at indices 11 and 12:
        # Layout 11: Blank blue SIL slide (for title/intro)
        # Layout 12: Blank white SIL slide (for content pages)

        # Check if we have the SIL custom layouts
        if layout_count >= 13:  # SIL template has 13 layouts (0-12)
            # Use the SIL custom layouts
            self.title_layout = self.presentation.slide_layouts[11]  # Blue SIL layout
            self.content_layout = self.presentation.slide_layouts[12]  # White SIL layout
            logger.info("Using SIL custom layouts - Blue (#11) and White (#12)")
        else:
            # Fallback to standard layouts if SIL template not loaded properly
            logger.warning("SIL custom layouts not found, using standard layouts")
            self.title_layout = self.presentation.slide_layouts[0]  # Standard title
            self.content_layout = self.presentation.slide_layouts[1]  # Standard content

        logger.info(f"Using layouts - Title: {self.title_layout.name}, Content: {self.content_layout.name}")

    def add_title_slide(self):
        """Add a new slide using the title layout (SIL blue or standard)"""
        slide = self.presentation.slides.add_slide(self.title_layout)
        logger.debug(f"Added title slide using layout: {self.title_layout.name}")
        return slide

    def add_content_slide(self):
        """Add a new slide using the content layout (SIL white or standard)"""
        slide = self.presentation.slides.add_slide(self.content_layout)
        logger.debug(f"Added content slide using layout: {self.content_layout.name}")
        return slide

    def add_blank_slide(self):
        """Add a new slide using the blank layout"""
        slide = self.presentation.slides.add_slide(self.blank_layout)
        logger.debug(f"Added blank slide using layout: {self.blank_layout.name}")
        return slide

    def hex_to_rgb(self, hex_color: str) -> RGBColor:
        """
        Convert hex color to RGBColor

        Args:
            hex_color: Hex color string (e.g., '#FF0000' or 'FF0000')

        Returns:
            RGBColor object
        """
        # Remove # if present
        hex_color = hex_color.lstrip('#')

        # Convert to RGB
        r = int(hex_color[0:2], 16)
        g = int(hex_color[2:4], 16)
        b = int(hex_color[4:6], 16)

        return RGBColor(r, g, b)

    def add_text_box(self, slide, text: str,
                     left: float, top: float,
                     width: float, height: float,
                     font_size: int = 14,
                     bold: bool = False,
                     alignment: PP_ALIGN = PP_ALIGN.LEFT,
                     font_name: Optional[str] = None) -> None:
        """
        Add a text box to slide with common formatting

        Args:
            slide: Slide object
            text: Text content
            left, top, width, height: Position and size in inches
            font_size: Font size in points
            bold: Whether text should be bold
            alignment: Text alignment
            font_name: Font family (defaults to Red Hat Display)
        """
        if font_name is None:
            font_name = self.default_font

        text_box = slide.shapes.add_textbox(
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )

        text_frame = text_box.text_frame
        text_frame.text = text
        text_frame.word_wrap = True

        p = text_frame.paragraphs[0]
        p.font.name = font_name
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.alignment = alignment

    def add_image(self, slide, image_path: Path,
                  left: float, top: float,
                  width: Optional[float] = None,
                  height: Optional[float] = None) -> None:
        """
        Add an image to slide

        Args:
            slide: Slide object
            image_path: Path to image file
            left, top: Position in inches
            width, height: Size in inches (maintains aspect if only one specified)
        """
        if width and height:
            slide.shapes.add_picture(
                str(image_path),
                Inches(left), Inches(top),
                width=Inches(width), height=Inches(height)
            )
        elif width:
            slide.shapes.add_picture(
                str(image_path),
                Inches(left), Inches(top),
                width=Inches(width)
            )
        elif height:
            slide.shapes.add_picture(
                str(image_path),
                Inches(left), Inches(top),
                height=Inches(height)
            )
        else:
            slide.shapes.add_picture(
                str(image_path),
                Inches(left), Inches(top)
            )

    def save(self, output_path: Path) -> Path:
        """
        Save presentation to file

        Args:
            output_path: Where to save the presentation

        Returns:
            Path to saved file
        """
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)

        self.presentation.save(str(output_path))
        logger.info(f"Presentation saved to {output_path}")

        return output_path

    def get_slide_dimensions(self) -> Tuple[float, float]:
        """
        Get slide dimensions in inches

        Returns:
            (width, height) in inches
        """
        width = self.presentation.slide_width / 914400  # EMUs to inches
        height = self.presentation.slide_height / 914400

        return width, height

    def apply_font_to_text_frame(self, text_frame, font_name: Optional[str] = None,
                                 font_size: Optional[int] = None, bold: Optional[bool] = None):
        """
        Apply font styling to all paragraphs in a text frame

        Args:
            text_frame: PowerPoint text frame object
            font_name: Font family name (defaults to Red Hat Display)
            font_size: Font size in points
            bold: Whether text should be bold
        """
        if font_name is None:
            font_name = self.default_font

        for paragraph in text_frame.paragraphs:
            if font_name:
                paragraph.font.name = font_name
            if font_size is not None:
                paragraph.font.size = Pt(font_size)
            if bold is not None:
                paragraph.font.bold = bold