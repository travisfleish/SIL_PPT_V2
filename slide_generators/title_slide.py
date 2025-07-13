# slide_generators/title_slide.py
"""
Generate title slide for PowerPoint presentations
Creates the opening slide with team branding and report title
"""

from pathlib import Path
from typing import Dict, Optional, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import logging

from .base_slide import BaseSlide

logger = logging.getLogger(__name__)


class TitleSlide(BaseSlide):
    """Generate the title slide for the presentation"""

    def __init__(self, presentation: Presentation = None):
        """
        Initialize title slide generator

        Args:
            presentation: Existing presentation to add slide to
        """
        super().__init__(presentation)

    def generate(self,
                 team_config: Dict[str, Any],
                 subtitle: str = "Sponsorship Insights Report",
                 slide_index: Optional[int] = None) -> Presentation:
        """
        Generate the title slide

        Args:
            team_config: Team configuration with name and colors
            subtitle: Subtitle for the presentation
            slide_index: Where to insert slide (None = append)

        Returns:
            Updated presentation object
        """
        # Extract team info
        team_name = team_config.get('team_name', 'Team')
        team_colors = team_config.get('colors', {})

        # Use the title layout (SIL blue layout #11)
        slide = self.add_title_slide()
        logger.info("Added title slide using SIL blue layout")

        # The SIL blue layout already has the background, so we don't need to add it
        # Only add background if using fallback layout
        if self.title_layout == self.presentation.slide_layouts[0]:  # Standard title layout
            # Add background color manually for non-SIL templates
            background = slide.background
            fill = background.fill
            fill.solid()

            # Use team primary color or default
            if 'primary' in team_colors:
                primary_color = team_colors['primary']
                if isinstance(primary_color, str):
                    fill.fore_color.rgb = self._hex_to_rgb(primary_color)
                else:
                    fill.fore_color.rgb = primary_color
            else:
                fill.fore_color.rgb = RGBColor(0, 34, 68)  # Default blue

        # Add team logo image (positioned where title would be)
        self._add_team_logo(slide, team_config)

        # Add subtitle
        subtitle_box = slide.shapes.add_textbox(
            Inches(2), Inches(3.8),  # Moved up with more space from logo
            Inches(9.333), Inches(1)  # Wider for 16:9
        )
        subtitle_box.text_frame.text = subtitle
        subtitle_box.text_frame.word_wrap = True

        # Format subtitle
        for paragraph in subtitle_box.text_frame.paragraphs:
            paragraph.font.name = self.default_font  # Red Hat Display
            paragraph.font.size = Pt(48)
            paragraph.font.bold = True
            paragraph.font.italic = True
            paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White text
            paragraph.alignment = PP_ALIGN.CENTER

        # Add descriptive text
        desc_text = "A purchase-based, data-backed view into who fans are and the brands they buy."
        desc_box = slide.shapes.add_textbox(
            Inches(3.5), Inches(4.8),  # Moved up with spacing from subtitle
            Inches(6.333), Inches(0.8)  # Narrower width
        )
        desc_box.text_frame.text = desc_text
        desc_box.text_frame.word_wrap = True

        # Format descriptive text
        for paragraph in desc_box.text_frame.paragraphs:
            paragraph.font.name = self.default_font  # Red Hat Display
            paragraph.font.size = Pt(18)
            paragraph.font.italic = True
            paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White text
            paragraph.alignment = PP_ALIGN.CENTER

        # Add SIL logo in corner
        self._add_sil_logo(slide)

        logger.info(f"Generated title slide for {team_name}")
        return self.presentation

    def _add_team_logo(self, slide, team_config: Dict[str, Any]):
        """Add team logo image to the slide"""
        try:
            # Get team name for file lookup
            team_name = team_config.get('team_name', '').lower().replace(' ', '_')

            # Construct path to team logo
            # Expected format: assets/logos/teams/utah_jazz.png
            # Get the project root directory
            from pathlib import Path
            import os

            # Try to find the project root
            current_file = Path(__file__).resolve()
            project_root = current_file.parent.parent  # Go up from slide_generators/title_slide.py

            logo_path = project_root / 'assets' / 'logos' / 'teams' / f"{team_name}.png"

            # Check if file exists
            if not logo_path.exists():
                logger.warning(f"Team logo not found at {logo_path}, falling back to text")
                self._add_team_name_text(slide, team_config)
                return

            # Add the logo image
            # Position it where the team name would go
            logo_height = Inches(2.5)  # Larger size for main logo

            # Add the logo first to get its dimensions
            try:
                logo = slide.shapes.add_picture(
                    str(logo_path),
                    Inches(0),  # Temporary position
                    Inches(0),
                    height=logo_height
                )

                # Center the logo horizontally where the title would be
                slide_width = Inches(13.333)  # 16:9 slide width
                logo_x = (slide_width - logo.width) / 2
                logo_y = Inches(1.0)  # Position higher up on slide

                # Move logo to centered position
                logo.left = int(logo_x)
                logo.top = int(logo_y)

                logger.info(f"Added team logo from {logo_path}")

            except Exception as e:
                logger.error(f"Error adding logo image: {e}")
                self._add_team_name_text(slide, team_config)

        except Exception as e:
            logger.error(f"Error in _add_team_logo: {e}")
            self._add_team_name_text(slide, team_config)

    def _add_team_name_text(self, slide, team_config: Dict[str, Any]):
        """Add team name as text (fallback when logo not available)"""
        team_name = team_config.get('team_name', 'Team')

        # Add main title (team name)
        title_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(1.5),  # Positioned higher up
            Inches(10.333), Inches(1.5)  # Wider for 16:9
        )
        title_box.text_frame.text = team_name
        title_box.text_frame.word_wrap = True

        # Format title
        for paragraph in title_box.text_frame.paragraphs:
            paragraph.font.name = self.default_font  # Red Hat Display
            paragraph.font.size = Pt(48)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White text
            paragraph.alignment = PP_ALIGN.CENTER

    def _add_sil_logo(self, slide):
        """Add SIL logo to bottom left corner"""
        try:
            # Get the project root directory
            from pathlib import Path

            current_file = Path(__file__).resolve()
            project_root = current_file.parent.parent  # Go up from slide_generators/title_slide.py

            # Path to SIL logo
            sil_logo_path = project_root / 'assets' / 'logos' / 'general' / 'SIL_white.png'

            if sil_logo_path.exists():
                # Add the SIL logo image
                logo_height = Inches(1.0)  # Reasonable size for corner logo

                try:
                    sil_logo = slide.shapes.add_picture(
                        str(sil_logo_path),
                        Inches(0.5),  # Left margin
                        Inches(6.0),  # Bottom position
                        height=logo_height
                    )
                    logger.info(f"Added SIL logo from {sil_logo_path}")

                except Exception as e:
                    logger.error(f"Error adding SIL logo: {e}")
                    self._add_sil_text_fallback(slide)
            else:
                logger.warning(f"SIL logo not found at {sil_logo_path}")
                self._add_sil_text_fallback(slide)

        except Exception as e:
            logger.error(f"Error in _add_sil_logo: {e}")
            self._add_sil_text_fallback(slide)

    def _add_sil_text_fallback(self, slide):
        """Add SIL text as fallback when logo not available"""
        sil_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(6.8),  # Bottom left
            Inches(2), Inches(0.5)
        )
        sil_box.text_frame.text = "Sports Innovation Lab"

        # Format SIL text
        for paragraph in sil_box.text_frame.paragraphs:
            paragraph.font.name = self.default_font  # Red Hat Display
            paragraph.font.size = Pt(12)
            paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White text

    def _hex_to_rgb(self, hex_color: str) -> RGBColor:
        """Convert hex color to RGBColor"""
        hex_color = hex_color.lstrip('#')
        r = int(hex_color[0:2], 16)
        g = int(hex_color[2:4], 16)
        b = int(hex_color[4:6], 16)
        return RGBColor(r, g, b)