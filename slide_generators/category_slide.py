# slide_generators/category_slide.py
"""
Generate category analysis slides for PowerPoint presentations
Formats category spending data into professional slides matching SIL template
"""

from pathlib import Path
from typing import Dict, Optional, Any, List, Tuple
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
import pandas as pd
import logging
import re

from .base_slide import BaseSlide
from data_processors.category_analyzer import CategoryAnalyzer, CategoryMetrics

logger = logging.getLogger(__name__)


# Formatting utility functions
def format_percentage_no_decimal(value):
    """Format percentage without decimals, handle 0% as EQUAL"""
    if isinstance(value, str):
        # Extract number from string like "10.5%"
        try:
            num = float(value.replace('%', '').strip())
        except:
            return value
    else:
        num = float(value)

    # Handle 0% special case
    if num == 0:
        return "EQUAL"

    # Round to nearest whole number
    return f"{int(round(num))}%"


def format_currency_no_cents(value):
    """Format currency without cents and with commas"""
    if isinstance(value, str):
        # Extract number from string like "$2,536.17"
        try:
            num = float(value.replace('$', '').replace(',', '').strip())
        except:
            return value
    else:
        num = float(value)

    return f"${int(round(num)):,}"


def clean_text_references(text):
    """Remove (Excl. Jazz Fans) and similar references"""
    # Remove various exclusion patterns
    patterns = [
        r'\(Excl\. [^)]*\)',
        r'\(excl\. [^)]*\)',
        r'\(Excluding [^)]*\)',
        r'\(excluding [^)]*\)'
    ]

    for pattern in patterns:
        text = re.sub(pattern, '', text)

    # Clean up extra spaces
    text = re.sub(r'\s+', ' ', text).strip()

    return text


def format_gen_pop_references(text):
    """Convert 'gen pop' to 'local gen pop' without duplication"""
    # First handle cases where "local gen pop" might already exist (fix duplications)
    text = re.sub(r'\blocal\s+local\s+gen\s+pop\b', 'local gen pop', text, flags=re.IGNORECASE)

    # Handle "Local Gen Pop" (capitalized) -> "local gen pop"
    text = re.sub(r'\bLocal\s+Gen\s+Pop\b', 'local gen pop', text)

    # Then handle standalone "gen pop" that doesn't already have "local"
    text = re.sub(r'(?<!local\s)\bgen pop\b', 'local gen pop', text, flags=re.IGNORECASE)

    return text


def clean_subcategory_duplication(text, category_name):
    """Remove category name duplication in subcategory references"""
    # Convert category name to different case variations for matching
    category_variations = [
        category_name,
        category_name.lower(),
        category_name.upper(),
        category_name.title()
    ]

    # Remove "Category - Subcategory" patterns and just keep "Subcategory"
    for cat_var in category_variations:
        # Pattern: "Category - Subcategory" -> "Subcategory"
        pattern = rf'\b{re.escape(cat_var)}\s*-\s*'
        text = re.sub(pattern, '', text, flags=re.IGNORECASE)

    return text.strip()


def remove_unnecessary_comparisons(text):
    """Remove comparison language when no actual comparison is provided"""
    # Remove "vs. the Local Gen Pop" when it's just stating a fact without comparison
    # Look for patterns like "spend $X vs. the Local Gen Pop" without actual comparison

    # Pattern: "spend $amount vs. the Local Gen Pop" -> "spend $amount"
    text = re.sub(r'spend\s+(\$[\d,]+)\s+vs\.\s+the\s+Local\s+Gen\s+Pop', r'spend \1', text, flags=re.IGNORECASE)

    # Pattern: "spend $amount per year vs. the Local Gen Pop" -> "spend $amount per year"
    text = re.sub(r'spend\s+(\$[\d,]+)\s+per\s+year\s+vs\.\s+the\s+Local\s+Gen\s+Pop', r'spend \1 per year', text,
                  flags=re.IGNORECASE)

    return text


def process_insight_text_enhanced(insight_text: str, category_name: str = "") -> str:
    """Enhanced process insight text with all formatting rules including subcategory fixes"""
    # Apply all formatting rules
    text = clean_text_references(insight_text)
    text = format_gen_pop_references(text)

    # Fix capitalization
    text = text.replace(' MORE ', ' more ')

    # Clean subcategory duplication if category name provided
    if category_name:
        text = clean_subcategory_duplication(text, category_name)

    # Remove unnecessary comparison language
    text = remove_unnecessary_comparisons(text)

    # Format any embedded percentages and currency
    def format_pct_match(match):
        return format_percentage_no_decimal(match.group(0))

    text = re.sub(r'\d+\.?\d*%', format_pct_match, text)

    # Find and format currency
    def format_curr_match(match):
        return format_currency_no_cents(match.group(0))

    text = re.sub(r'\$[\d,]+\.?\d*', format_curr_match, text)

    return text


def process_insight_text(insight_text: str) -> str:
    """Process insight text with all formatting rules - legacy function for compatibility"""
    return process_insight_text_enhanced(insight_text)


def format_subcategory_text(category_name: str, subcategory_name: str) -> str:
    """Format subcategory text to avoid duplication"""
    # If subcategory already contains the category name, just return subcategory
    if category_name.lower() in subcategory_name.lower():
        return subcategory_name

    # Otherwise return the full name
    return f"{category_name} - {subcategory_name}"


class CategorySlide(BaseSlide):
    """Generate category analysis slides with insights, subcategory stats, and merchant rankings"""

    def __init__(self, presentation: Presentation = None):
        """
        Initialize category slide generator

        Args:
            presentation: Existing presentation to add slide to
        """
        super().__init__(presentation)

        # Colors for the slide - UPDATED with EQUAL color
        self.colors = {
            'header_bg': RGBColor(240, 240, 240),
            'header_border': RGBColor(200, 200, 200),
            'table_header': RGBColor(217, 217, 217),
            'table_border': RGBColor(0, 0, 0),
            'positive': RGBColor(0, 176, 80),  # Green
            'negative': RGBColor(255, 0, 0),  # Red
            'equal': RGBColor(184, 134, 11),  # Dark yellow for EQUAL
            'neutral': RGBColor(0, 0, 0)  # Black
        }

    def generate(self,
                 analysis_results: Dict[str, Any],
                 team_config: Dict[str, Any],
                 slide_index: Optional[int] = None) -> Presentation:
        """
        Generate the complete category analysis slide

        Args:
            analysis_results: Results from CategoryAnalyzer.analyze_category()
            team_config: Team configuration with colors, names, etc.
            slide_index: Where to insert slide (None = append)

        Returns:
            Updated presentation object
        """
        # Extract team info
        team_name = team_config.get('team_name', 'Team')
        team_short = team_config.get('team_name_short', team_name.split()[-1])

        # Use the content layout (SIL white layout #12)
        slide = self.add_content_slide()
        logger.info(f"Added category slide for {analysis_results['display_name']} using SIL white layout")

        # Add header
        self._add_header(slide, team_name, analysis_results['slide_title'])

        # Add title - UPDATED to use "Category Analysis: [CATEGORY]"
        category_title = f"Category Analysis: {analysis_results['display_name']}"
        self._add_title(slide, category_title)

        # Add category insights (left side) - UPDATED with formatting fixes
        self._add_category_insights(slide, analysis_results, team_short, team_config)

        # Add category metrics table (top right) - pass results for category name
        self._add_category_table(slide, analysis_results)

        # Add subcategory table (middle right) - adjusted for 16:9
        self._add_subcategory_table(slide, analysis_results['subcategory_stats'])

        logger.info(f"Generated {analysis_results['display_name']} slide")
        return self.presentation

    def generate_brand_slide(self,
                             analysis_results: Dict[str, Any],
                             team_config: Dict[str, Any],
                             slide_index: Optional[int] = None) -> Presentation:
        """
        Generate the brand analysis slide (second slide for category)

        Args:
            analysis_results: Results from CategoryAnalyzer
            team_config: Team configuration
            slide_index: Where to insert slide

        Returns:
            Updated presentation
        """
        # Extract team info
        team_name = team_config.get('team_name', 'Team')
        team_short = team_config.get('team_name_short', team_name.split()[-1])
        category_name = analysis_results['display_name']

        # Use the content layout (SIL white layout #12)
        slide = self.add_content_slide()
        logger.info(f"Added brand slide for {analysis_results['display_name']} using SIL white layout")

        # UPDATED brand slide title format based on category
        if category_name.upper() == "QSR":
            brand_title = f"Top QSR Brands for {team_name} Fans"
        else:
            brand_title = f"Top {category_name} Brands for {team_name} Fans"

        # Add header - brand slide uses updated title format
        header_title = f"Sponsor Spending Analysis: {category_name} Brands"
        self._add_header(slide, team_name, header_title)

        # Add title with updated format
        self._add_title(slide, brand_title)

        # Add brand logos (numbered circles) - adjusted for 16:9
        self._add_brand_logos(slide, analysis_results['merchant_stats'])

        # Add brand insights (left side) - UPDATED with formatting
        self._add_brand_insights(slide, analysis_results, team_name, team_short, category_name)

        # Add brand table (right side) - adjusted for 16:9
        self._add_brand_table(slide, analysis_results['merchant_stats'])

        # Add sponsorship recommendation - UPDATED with formatting
        self._add_sponsor_recommendation(slide, analysis_results['recommendation'], team_config)

        logger.info(f"Generated {analysis_results['display_name']} brand slide")
        return self.presentation

    def _add_header(self, slide, team_name: str, slide_title: str):
        """Add header with team name and slide title"""
        # Team name (left)
        team_text = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.05),
            Inches(4), Inches(0.3)
        )
        team_text.text_frame.text = team_name
        p = team_text.text_frame.paragraphs[0]
        p.font.name = self.default_font  # Red Hat Display
        p.font.size = Pt(12)

        # Slide title (right) - adjusted position for 16:9
        title_text = slide.shapes.add_textbox(
            Inches(6.333), Inches(0.05),  # Moved right for 16:9
            Inches(6.8), Inches(0.3)  # Wider for 16:9
        )
        title_text.text_frame.text = slide_title
        p = title_text.text_frame.paragraphs[0]
        p.font.name = self.default_font  # Red Hat Display
        p.alignment = PP_ALIGN.RIGHT
        p.font.size = Pt(12)

    def _add_title(self, slide, title: str):
        """Add main slide title"""
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.6),
            Inches(12.333), Inches(0.5)  # Adjusted for 16:9 width
        )
        title_box.text_frame.text = title
        p = title_box.text_frame.paragraphs[0]
        p.font.name = self.default_font  # Red Hat Display
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(85, 85, 85)

    def _add_category_insights(self, slide, results: Dict[str, Any], team_short: str, team_config: Dict[str, Any]):
        """Add category insights section - UPDATED WITH ALL FORMATTING FIXES"""
        # Insights title
        insights_title = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5),
            Inches(4), Inches(0.3)
        )
        insights_title.text_frame.text = "Category Insights:"
        p = insights_title.text_frame.paragraphs[0]
        p.font.name = self.default_font
        p.font.size = Pt(14)
        p.font.bold = True

        # Get category name for subcategory cleaning
        category_name = results.get('display_name', '')

        # Separate insights into regular and NBA comparison
        regular_insights = []
        nba_insights = []

        for insight in results['insights']:
            # Clean and format the insight text with enhanced processing
            cleaned_insight = process_insight_text_enhanced(insight, category_name)

            if "NBA" in cleaned_insight and "compared to" in cleaned_insight:
                nba_insights.append(cleaned_insight)
            else:
                regular_insights.append(cleaned_insight)

        # Regular insights box
        insights_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(1.9),
            Inches(4.5), Inches(3.0)  # Reduced height to make room for NBA section
        )

        text_frame = insights_box.text_frame
        text_frame.word_wrap = True

        # Add regular insights with BULLETS instead of numbers
        for i, insight in enumerate(regular_insights[:4]):
            p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
            p.text = f"• {insight}"  # Changed from f"{i+1}. {insight}"
            p.font.name = self.default_font
            p.font.size = Pt(11)
            p.line_spacing = 1.2

        # Add NBA comparison section if there are NBA insights
        if nba_insights:
            # NBA comparison label
            nba_label = slide.shapes.add_textbox(
                Inches(0.5), Inches(5.2),
                Inches(4), Inches(0.3)
            )
            nba_label.text_frame.text = f"{team_short} Fans vs. {team_config.get('league', 'NBA')} Fans:"
            p = nba_label.text_frame.paragraphs[0]
            p.font.name = self.default_font
            p.font.size = Pt(14)
            p.font.bold = True

            # NBA insights box
            nba_box = slide.shapes.add_textbox(
                Inches(0.7), Inches(5.6),
                Inches(4.5), Inches(1.5)
            )

            nba_text_frame = nba_box.text_frame
            nba_text_frame.word_wrap = True

            # Add NBA insights with BULLETS instead of numbers
            for i, insight in enumerate(nba_insights[:2]):
                p = nba_text_frame.add_paragraph() if i > 0 else nba_text_frame.paragraphs[0]
                p.text = f"• {insight}"  # Changed from numbered
                p.font.name = self.default_font
                p.font.size = Pt(11)
                p.line_spacing = 1.2

    def _add_category_table(self, slide, results: Dict[str, Any]):
        """Add category metrics table (adjusted for 16:9) with formatting fixes"""
        # Extract metrics from results
        metrics = results['category_metrics']

        # Table position - adjusted to prevent bleeding
        left = Inches(5.8)  # Moved left to fit better
        top = Inches(1.5)
        width = Inches(6.8)  # Adjusted width to fit within slide
        height = Inches(0.8)  # Reduced height for better proportions

        # Create table
        table_shape = slide.shapes.add_table(2, 4, left, top, width, height)
        table = table_shape.table

        # Set column widths - better distribution
        table.columns[0].width = Inches(1.3)  # Category
        table.columns[1].width = Inches(1.6)  # Percent of Fans
        table.columns[2].width = Inches(2.0)  # How likely
        table.columns[3].width = Inches(1.9)  # Purchases

        # Header row - UPDATED with "local gen pop"
        headers = ['Category', 'Percent of Fans\nWho Spend', 'How likely fans are to\nspend vs. local gen pop',
                   'How many more purchases\nper fan vs. local gen pop']

        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            self._format_header_cell(cell)

        # Data row - extract category name properly and apply formatting
        category_name = results.get('display_name', 'Category')
        table.cell(1, 0).text = category_name

        # Apply formatting to the metric values
        table.cell(1, 1).text = self._format_metric_value(metrics.format_percent_fans())
        table.cell(1, 2).text = self._format_metric_value(metrics.format_likelihood())
        table.cell(1, 3).text = self._format_metric_value(metrics.format_purchases())

        # Format data cells
        for i in range(4):
            self._format_data_cell(table.cell(1, i))

    def _add_subcategory_table(self, slide, subcategory_stats: pd.DataFrame):
        """Add subcategory statistics table (adjusted for 16:9) with formatting fixes"""
        if subcategory_stats.empty:
            return

        # Table position - adjusted to prevent bleeding
        left = Inches(5.8)  # Moved left to fit better
        top = Inches(2.7)  # Adjusted vertical position
        width = Inches(6.8)  # Adjusted width to fit within slide

        # Create table with header + data rows
        rows = min(len(subcategory_stats), 5) + 1  # Max 5 subcategories + header
        table_shape = slide.shapes.add_table(rows, 4, left, top, width, Inches(0.3 * rows))
        table = table_shape.table

        # Set column widths - better distribution
        table.columns[0].width = Inches(1.6)  # Sub-Category
        table.columns[1].width = Inches(1.5)  # Percent of Fans
        table.columns[2].width = Inches(1.9)  # How likely
        table.columns[3].width = Inches(1.8)  # Purchases

        # Headers - UPDATED with "local gen pop"
        headers = ['Sub-Category', 'Percent of Fans\nWho Spend', 'How likely fans are to\nspend vs. local gen pop',
                   'How many more purchases\nper fan vs. local gen pop']

        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            self._format_header_cell(cell)

        # Data rows with formatting
        for row_idx, (_, row) in enumerate(subcategory_stats.iterrows(), 1):
            if row_idx >= rows:
                break

            # Apply subcategory text formatting to avoid duplication
            subcategory_text = row['Subcategory']
            table.cell(row_idx, 0).text = subcategory_text

            # Apply formatting to metric values
            table.cell(row_idx, 1).text = self._format_metric_value(row['Percent of Fans Who Spend'])
            table.cell(row_idx, 2).text = self._format_metric_value(row['How likely fans are to spend vs. gen pop'])
            table.cell(row_idx, 3).text = self._format_metric_value(row['Purchases per fan vs. gen pop'])

            # Format cells
            for col in range(4):
                self._format_data_cell(table.cell(row_idx, col))

    def _add_brand_logos(self, slide, merchant_stats: Tuple[pd.DataFrame, List[str]]):
        """Add brand logo placeholders (numbered circles) - adjusted for 16:9"""
        merchant_df, top_merchants = merchant_stats

        if merchant_df.empty:
            return

        # Position for logos - adjusted spacing for 16:9
        start_x = Inches(0.5)
        y = Inches(1.2)
        spacing = Inches(2.4)  # Wider spacing for 16:9

        # Add numbered circles for top 5 brands
        for i in range(min(5, len(merchant_df))):
            x = start_x + (i * spacing)

            # Circle
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                x, y,
                Inches(1.2), Inches(1.2)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = RGBColor(255, 255, 255)
            circle.line.color.rgb = RGBColor(200, 200, 200)
            circle.line.width = Pt(2)

            # Number
            text_box = slide.shapes.add_textbox(
                x, y + Inches(0.4),
                Inches(1.2), Inches(0.4)
            )
            text_box.text_frame.text = str(i + 1)
            p = text_box.text_frame.paragraphs[0]
            p.font.name = self.default_font  # Red Hat Display
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(48)
            p.font.color.rgb = RGBColor(150, 150, 150)

    def _add_brand_insights(self, slide, results: Dict[str, Any], team_name: str, team_short: str, category_name: str):
        """Add brand-specific insights with updated formatting"""
        # Top Brand Insights section
        insights_title = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.8),
            Inches(4), Inches(0.3)
        )
        insights_title.text_frame.text = "Top Brand Insights"
        p = insights_title.text_frame.paragraphs[0]
        p.font.name = self.default_font  # Red Hat Display
        p.font.size = Pt(14)
        p.font.bold = True

        # Insights box
        insights_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(3.2),
            Inches(4.5), Inches(1.6)  # Reduced height to make room for Top Brand Target
        )

        text_frame = insights_box.text_frame
        text_frame.word_wrap = True

        # UPDATED: Format insight #2 with standardized text
        merchant_insights = results.get('merchant_insights', [])

        for i, insight in enumerate(merchant_insights[:4]):
            p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]

            # Format insight #2 (index 1) with standardized format
            if i == 1 and 'purchases per year' in insight.lower():
                # Extract brand and number from insight for standardization
                formatted_insight = self._format_insight_two(insight, team_short, category_name)
                p.text = f"• {formatted_insight}"
            else:
                # Apply general formatting
                formatted_insight = process_insight_text(insight)
                p.text = f"• {formatted_insight}"

            p.font.name = self.default_font  # Red Hat Display
            p.font.size = Pt(10)
            p.line_spacing = 1.0

    def _add_brand_table(self, slide, merchant_stats: Tuple[pd.DataFrame, List[str]]):
        """Add brand ranking table - adjusted for 16:9 with formatting fixes"""
        merchant_df, _ = merchant_stats

        if merchant_df.empty:
            return

        # Create table - adjusted position and size for 16:9
        rows = min(len(merchant_df) + 1, 6)  # Header + max 5 brands
        cols = 5
        left = Inches(5.833)  # Moved right for 16:9
        top = Inches(1.2)
        width = Inches(7.0)  # Wider for 16:9
        height = Inches(0.35 * rows)

        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        # Headers
        headers = ['Rank (by\npercent of\nfans who\nspend)', 'Brand',
                   'Percent of\nFans Who\nSpend',
                   'How likely\nfans are to\nspend vs.\nlocal gen pop',
                   'Purchases\nPer Fan\n(vs. Local\nGen Pop)']

        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            self._format_header_cell(cell, small=True)

        # Data rows with formatting
        for row_idx, (_, row) in enumerate(merchant_df.iterrows(), 1):
            if row_idx >= rows:
                break

            table.cell(row_idx, 0).text = str(row['Rank'])
            table.cell(row_idx, 1).text = row['Brand']

            # Apply formatting to metric values
            table.cell(row_idx, 2).text = self._format_metric_value(row['Percent of Fans Who Spend'])
            table.cell(row_idx, 3).text = self._format_metric_value(row['How likely fans are to spend vs. gen pop'])
            table.cell(row_idx, 4).text = self._format_metric_value(row['Purchases Per Fan (vs. Gen Pop)'])

            # Format cells
            for col in range(5):
                self._format_data_cell(table.cell(row_idx, col), small=True)

    def _add_sponsor_recommendation(self, slide, recommendation: Optional[Dict[str, Any]], team_config: Dict[str, Any]):
        """Add Top Brand Target and sponsorship recommendation with updated formatting"""
        if not recommendation:
            return

        # Top Brand Target header
        target_title = slide.shapes.add_textbox(
            Inches(0.5), Inches(5.0),
            Inches(4), Inches(0.3)
        )
        target_title.text_frame.text = "Top Brand Target"
        p = target_title.text_frame.paragraphs[0]
        p.font.name = self.default_font  # Red Hat Display
        p.font.size = Pt(14)
        p.font.bold = True

        # Recommendation content
        rec_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(5.4),
            Inches(6), Inches(1.2)
        )

        text_frame = rec_box.text_frame
        text_frame.word_wrap = True

        # UPDATED: Standardized recommendation format with composite index
        team_name = team_config.get('team_name', 'Team')
        merchant_name = recommendation.get('merchant', 'Brand')
        composite_index = recommendation.get('composite_index', 'N/A')

        # First bullet - main recommendation with composite index
        p1 = text_frame.paragraphs[0]
        p1.text = f"• The {team_name} should target {merchant_name} for a sponsorship based on having the highest composite index of {composite_index}"
        p1.font.name = self.default_font
        p1.font.size = Pt(11)
        p1.line_spacing = 1.2

        # Second bullet - standardized explanation
        p2 = text_frame.add_paragraph()
        p2.text = "• The composite index indicates a brand with significant likelihood for more fans to be spending more frequently, and at a higher spend per fan vs. other brands"
        p2.font.name = self.default_font
        p2.font.size = Pt(11)
        p2.line_spacing = 1.2

    def _format_header_cell(self, cell, small: bool = False):
        """Format table header cell"""
        cell.fill.solid()
        cell.fill.fore_color.rgb = self.colors['table_header']

        # Format text
        text_frame = cell.text_frame
        text_frame.margin_left = Inches(0.05)
        text_frame.margin_right = Inches(0.05)
        text_frame.margin_top = Inches(0.03)
        text_frame.margin_bottom = Inches(0.03)
        text_frame.word_wrap = True

        # Format all paragraphs in the cell (for multi-line headers)
        for paragraph in text_frame.paragraphs:
            paragraph.font.name = self.default_font  # Red Hat Display
            paragraph.font.size = Pt(8) if small else Pt(10)  # Smaller font for compact tables
            paragraph.font.bold = True
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.line_spacing = 1.0  # Tighter line spacing for headers

        # Vertical alignment
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    def _format_data_cell(self, cell, small: bool = False):
        """Format table data cell with updated color handling"""
        # Format text
        text_frame = cell.text_frame
        text_frame.margin_left = Inches(0.05)
        text_frame.margin_right = Inches(0.05)
        text_frame.margin_top = Inches(0.03)
        text_frame.margin_bottom = Inches(0.03)
        text_frame.word_wrap = True

        # Format all paragraphs in the cell
        for paragraph in text_frame.paragraphs:
            paragraph.font.name = self.default_font  # Red Hat Display
            paragraph.font.size = Pt(8) if small else Pt(11)  # Consistent data font size
            paragraph.alignment = PP_ALIGN.CENTER

            # UPDATED color coding with EQUAL handling
            if 'More' in cell.text or 'more' in cell.text:
                paragraph.font.color.rgb = self.colors['positive']  # Green
            elif 'Less' in cell.text or 'less' in cell.text or 'fewer' in cell.text:
                paragraph.font.color.rgb = self.colors['negative']  # Red
            elif 'EQUAL' in cell.text or 'Equal' in cell.text:
                paragraph.font.color.rgb = self.colors['equal']  # Dark yellow

        # Vertical alignment
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    def _format_metric_value(self, value: str) -> str:
        """Apply formatting to metric values (percentages, currency, etc.)"""
        if not value or pd.isna(value):
            return "N/A"

        # Clean text references first
        formatted_value = clean_text_references(str(value))
        formatted_value = format_gen_pop_references(formatted_value)

        # Handle percentage formatting
        if '%' in formatted_value:
            formatted_value = format_percentage_no_decimal(formatted_value)

        # Handle currency formatting
        if '$' in formatted_value:
            formatted_value = format_currency_no_cents(formatted_value)

        return formatted_value

    def _format_insight_two(self, insight: str, team_short: str, category_name: str) -> str:
        """Format insight #2 with standardized format"""
        # Extract relevant information using regex

        # Try to extract brand name and number of purchases
        brand_match = re.search(r'at ([A-Z\s&\']+)', insight, re.IGNORECASE)
        purchases_match = re.search(r'(\d+)\s+purchases?\s+per\s+year', insight, re.IGNORECASE)

        if brand_match and purchases_match:
            brand = brand_match.group(1)
            purchases = purchases_match.group(1)

            return f"{team_short} fans make an average of {purchases} purchases per year at {brand}—more than any other top {category_name} brand"

        # Fallback to original insight with general formatting
        return process_insight_text(insight)


# Convenience functions
def create_category_slides(analysis_results: Dict[str, Any],
                           team_config: Dict[str, Any],
                           presentation: Optional[Presentation] = None) -> Presentation:
    """
    Create both category slides (analysis + brands) for a category

    Args:
        analysis_results: Results from CategoryAnalyzer
        team_config: Team configuration
        presentation: Existing presentation (creates new if None)

    Returns:
        Presentation with both slides added
    """
    generator = CategorySlide(presentation)

    # Generate category analysis slide
    presentation = generator.generate(analysis_results, team_config)

    # Generate brand analysis slide
    presentation = generator.generate_brand_slide(analysis_results, team_config)

    return presentation