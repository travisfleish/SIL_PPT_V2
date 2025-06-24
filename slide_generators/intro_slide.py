# slide_generators/intro_slide.py
"""
Generate introduction/overview slide for PowerPoint presentations
This is a placeholder for potential future use
"""

from pathlib import Path
from typing import Dict, Optional, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import logging

from .base_slide import BaseSlide

logger = logging.getLogger(__name__)


class IntroSlide(BaseSlide):
    """
    Generate an introduction/overview slide

    This could be used for:
    - Executive summary
    - Table of contents
    - Key findings overview
    - How to use this report
    """

    def __init__(self, presentation: Presentation = None):
        """
        Initialize intro slide generator

        Args:
            presentation: Existing presentation to add slide to
        """
        super().__init__(presentation)

    def generate(self,
                 team_config: Dict[str, Any],
                 content_type: str = "overview",
                 slide_index: Optional[int] = None) -> Presentation:
        """
        Generate the intro slide

        Args:
            team_config: Team configuration
            content_type: Type of intro content ("overview", "toc", "how_to_use")
            slide_index: Where to insert slide (None = append)

        Returns:
            Updated presentation object
        """
        # Add slide
        slide = self.presentation.slides.add_slide(self.blank_layout)

        # Add header
        self._add_header(slide, team_config['team_name'])

        # Add content based on type
        if content_type == "overview":
            self._add_overview_content(slide, team_config)
        elif content_type == "toc":
            self._add_table_of_contents(slide, team_config)
        elif content_type == "how_to_use":
            self._add_how_to_use(slide, team_config)
        else:
            self._add_placeholder_content(slide)

        logger.info(f"Generated intro slide ({content_type}) for {team_config['team_name']}")
        return self.presentation

    def _add_header(self, slide, team_name: str):
        """Add standard header"""
        header_rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(0.5)
        )
        header_rect.fill.solid()
        header_rect.fill.fore_color.rgb = RGBColor(240, 240, 240)
        header_rect.line.color.rgb = RGBColor(200, 200, 200)
        header_rect.line.width = Pt(0.5)

        # Team name
        team_text = slide.shapes.add_textbox(
            Inches(0.2), Inches(0.1),
            Inches(4), Inches(0.3)
        )
        team_text.text_frame.text = team_name
        team_text.text_frame.paragraphs[0].font.size = Pt(14)
        team_text.text_frame.paragraphs[0].font.bold = True

    def _add_overview_content(self, slide, team_config: Dict[str, Any]):
        """Add executive overview content"""
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.8),
            Inches(9), Inches(0.5)
        )
        title_box.text_frame.text = "Executive Overview"
        title_box.text_frame.paragraphs[0].font.size = Pt(24)
        title_box.text_frame.paragraphs[0].font.bold = True

        # Key findings placeholder
        content_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5),
            Inches(9), Inches(5)
        )

        content = f"""Key Findings for {team_config['team_name']}:

• Fan Demographics: Unique audience profile with distinct characteristics
• Fan Behaviors: High-value communities with strong spending patterns  
• Category Opportunities: Multiple sponsorship categories with strong fan affinity
• Custom Insights: Data-driven recommendations for partnership targets

This report provides detailed analysis to support sponsorship sales and partnership development."""

        content_box.text_frame.text = content
        content_box.text_frame.word_wrap = True

        for paragraph in content_box.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.line_spacing = 1.5

    def _add_table_of_contents(self, slide, team_config: Dict[str, Any]):
        """Add table of contents"""
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.8),
            Inches(9), Inches(0.5)
        )
        title_box.text_frame.text = "Report Contents"
        title_box.text_frame.paragraphs[0].font.size = Pt(24)
        title_box.text_frame.paragraphs[0].font.bold = True

        # TOC items
        toc_items = [
            "1. Fan Demographics - Understanding Your Audience",
            "2. Fan Behaviors - How Fans Engage with Brands",
            "3. Restaurant Analysis - Food & Beverage Opportunities",
            "4. Athleisure Analysis - Apparel & Footwear Partners",
            "5. Finance Analysis - Banking & Payment Solutions",
            "6. Gambling Analysis - Gaming & Betting Partners",
            "7. Travel Analysis - Transportation & Hospitality",
            "8. Auto Analysis - Automotive Partners",
            "9. Custom Categories - Additional Opportunities"
        ]

        content_box = slide.shapes.add_textbox(
            Inches(1), Inches(1.8),
            Inches(8), Inches(5)
        )

        content_box.text_frame.text = "\n".join(toc_items)
        content_box.text_frame.word_wrap = True

        for paragraph in content_box.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.line_spacing = 1.8

    def _add_how_to_use(self, slide, team_config: Dict[str, Any]):
        """Add 'How to Use This Report' content"""
        # This matches the content from the sample PPT
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.8),
            Inches(9), Inches(0.5)
        )
        title_box.text_frame.text = "How To Use This Report:"
        title_box.text_frame.paragraphs[0].font.size = Pt(24)
        title_box.text_frame.paragraphs[0].font.bold = True

        # Three columns
        sections = [
            {
                'title': 'Identify The Right\nSponsorship Targets',
                'content': 'This report will uncover a set of categories custom to your team where your fans have unique and high value spending behaviors. Create a dedicated sales effort to secure partners from these categories, and use this report to prioritize those brands for outreach.'
            },
            {
                'title': 'Get data rich collateral that\nconverts',
                'content': 'This report includes spending habits by your fans on individual brands - including how many of your fans spend, how often they spend, and how that compares to your local market population. Use these data points to craft tailored marketing materials that hook and convert.'
            },
            {
                'title': 'Measure the effectiveness of\nsponsorship with real fan\nspending',
                'content': "You'll find valuable spending data for your current sponsors included in this report. Use this to justify the ROI of their investment and/or create benchmarks to continue to measure attribution and ROI of your partnership."
            }
        ]

        # Add each section
        for i, section in enumerate(sections):
            left = Inches(0.5 + i * 3.2)

            # Section title
            title_box = slide.shapes.add_textbox(
                left, Inches(1.8),
                Inches(2.8), Inches(0.8)
            )
            title_box.text_frame.text = section['title']
            title_box.text_frame.word_wrap = True
            p = title_box.text_frame.paragraphs[0]
            p.font.size = Pt(14)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER

            # Section content
            content_box = slide.shapes.add_textbox(
                left, Inches(2.8),
                Inches(2.8), Inches(3.5)
            )
            content_box.text_frame.text = section['content']
            content_box.text_frame.word_wrap = True

            for paragraph in content_box.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.line_spacing = 1.2
                paragraph.alignment = PP_ALIGN.LEFT

    def _add_placeholder_content(self, slide):
        """Add placeholder content"""
        text_box = slide.shapes.add_textbox(
            Inches(1), Inches(3),
            Inches(8), Inches(1)
        )
        text_box.text_frame.text = "Introduction slide placeholder"
        text_box.text_frame.paragraphs[0].font.size = Pt(24)
        text_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


# Convenience function
def create_intro_slide(team_config: Dict[str, Any],
                       presentation: Optional[Presentation] = None,
                       content_type: str = "overview") -> Presentation:
    """
    Create an intro slide for the presentation

    Args:
        team_config: Team configuration
        presentation: Existing presentation (creates new if None)
        content_type: Type of intro content

    Returns:
        Presentation with intro slide
    """
    generator = IntroSlide(presentation)
    return generator.generate(team_config, content_type)