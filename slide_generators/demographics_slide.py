# slide_generators/demographics_slide.py
"""
Generate complete demographics slide for PowerPoint presentations
Includes all charts arranged according to the reference layout
Updated with optimized layout for gender horizontal bars
Updated legend: removed KEY label, centered items, updated text
"""

from pathlib import Path
from typing import Dict, Optional, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import logging

logger = logging.getLogger(__name__)

# Default font
DEFAULT_FONT_FAMILY = "Red Hat Display"


class DemographicsSlide:
    """Generate the complete demographics slide with all charts"""

    def __init__(self, presentation: Presentation = None):
        """
        Initialize demographics slide generator

        Args:
            presentation: Existing presentation to add slide to
        """
        if presentation is None:
            self.presentation = Presentation()
            # Set 16:9 widescreen aspect ratio
            self.presentation.slide_width = Inches(13.333)
            self.presentation.slide_height = Inches(7.5)
        else:
            self.presentation = presentation

        self.default_font = DEFAULT_FONT_FAMILY
        self.team_colors = {}

        # Setup layouts
        self._setup_layouts()

    def _setup_layouts(self):
        """Setup slide layout references"""
        layout_count = len(self.presentation.slide_layouts)
        logger.info(f"Available slide layouts: {layout_count}")

        # Always keep reference to standard blank layout
        self.blank_layout = self.presentation.slide_layouts[6]  # Blank layout

        # Check if we have the SIL custom layouts (12 = white content layout)
        if layout_count >= 13:
            self.content_layout = self.presentation.slide_layouts[12]  # SIL white content
            logger.info("Using SIL white content layout")
        else:
            self.content_layout = self.blank_layout  # Fallback to blank
            logger.warning("SIL layouts not found, using blank layout")

    def add_content_slide(self):
        """Add a content slide using the appropriate layout"""
        return self.presentation.slides.add_slide(self.content_layout)

    def generate(self,
                 demographic_data: Dict[str, Any],
                 chart_dir: Path,
                 team_config: Dict[str, Any],
                 slide_index: Optional[int] = None) -> Presentation:
        """
        Generate the complete demographics slide

        Args:
            demographic_data: Processed demographic data
            chart_dir: Directory containing chart images
            team_config: Team configuration
            slide_index: Where to insert slide (None = append)

        Returns:
            Updated presentation object
        """
        # Extract team info
        team_name = team_config.get('team_name', 'Team')
        team_short = team_config.get('team_name_short', team_name.split()[-1])
        colors = team_config.get('colors', {})

        # Store team colors for legend
        self.team_colors = colors

        # Use the content layout
        slide = self.add_content_slide()
        logger.info("Added demographics slide using content layout")

        # Add header
        self._add_header(slide, team_name)

        # Add chart section headers with black backgrounds
        self._add_chart_headers(slide)

        # Add demographic charts with optimized spacing
        self._add_charts(slide, chart_dir)

        # Add centered legend without KEY label
        self._add_legend_box(slide, team_name, team_short, demographic_data.get('league', 'League'))

        logger.info(f"Generated demographics slide for {team_name}")
        return self.presentation

    def _add_header(self, slide, team_name: str):
        """Add the standard header bar at the top of the slide"""
        # Header bar dimensions (full width at top)
        header_rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(13.333), Inches(0.5)  # Full 16:9 width
        )

        # Style the header bar with light gray background
        header_rect.fill.solid()
        header_rect.fill.fore_color.rgb = RGBColor(240, 240, 240)  # Light gray
        header_rect.line.color.rgb = RGBColor(200, 200, 200)
        header_rect.line.width = Pt(0.5)

        # Add team name on the left side of header
        team_text = slide.shapes.add_textbox(
            Inches(0.2), Inches(0.1),
            Inches(3), Inches(0.3)
        )
        team_text.text_frame.text = team_name
        team_p = team_text.text_frame.paragraphs[0]
        team_p.font.name = self.default_font
        team_p.font.size = Pt(14)
        team_p.font.bold = True
        team_p.font.color.rgb = RGBColor(0, 0, 0)  # Black text

        # Add page title on the right side of header
        title_text = slide.shapes.add_textbox(
            Inches(6.5), Inches(0.1),
            Inches(6.633), Inches(0.3)
        )
        title_text.text_frame.text = f"Fan Demographics: How Are {team_name} Fans Unique"
        title_p = title_text.text_frame.paragraphs[0]
        title_p.font.name = self.default_font
        title_p.font.size = Pt(14)
        title_p.alignment = PP_ALIGN.RIGHT
        title_p.font.color.rgb = RGBColor(0, 0, 0)  # Black text

    def _add_chart_headers(self, slide):
        """Add black header bars for each chart section - with vertically centered text"""
        headers = [
            # Top row headers - back to original alignment
            ("GENDER", 0.5, 0.85, 1.5),  # Back to 1.5" width
            ("HOUSEHOLD INCOME", 2.2, 0.85, 4.8),  # Original position
            ("OCCUPATION CATEGORY", 7.2, 0.85, 5.6),  # Original position

            # Bottom row headers - original positions
            ("ETHNICITY", 0.5, 3.95, 4.5),
            ("GENERATION", 5.2, 3.95, 4.5),
            ("CHILDREN IN HOUSEHOLD", 9.8, 3.95, 3.0)
        ]

        for text, left, top, width in headers:
            # Black background bar
            header_rect = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(left), Inches(top),
                Inches(width), Inches(0.25)
            )
            header_rect.fill.solid()
            header_rect.fill.fore_color.rgb = RGBColor(0, 0, 0)
            header_rect.line.fill.background()

            # White text - adjusted positioning for vertical centering
            text_box = slide.shapes.add_textbox(
                Inches(left + 0.1), Inches(top),  # Start at top of bar
                Inches(width - 0.2), Inches(0.25)  # Full height of bar
            )
            text_box.text_frame.text = text

            # Configure text frame for vertical centering
            text_box.text_frame.margin_top = Pt(0)
            text_box.text_frame.margin_bottom = Pt(0)
            text_box.text_frame.margin_left = Pt(0)
            text_box.text_frame.margin_right = Pt(0)
            text_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE  # Vertical center

            p = text_box.text_frame.paragraphs[0]
            p.font.name = self.default_font
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER

    def _add_charts(self, slide, chart_dir: Path):
        """Add all demographic charts with realigned spacing for horizontal gender bars"""
        chart_dir = Path(chart_dir)

        # REALIGNED LAYOUT: Back to original alignment now that gender uses horizontal bars
        chart_positions = [
            # Top row - Gender back to narrow width, aligning with bottom row
            ('gender_chart', 0.5, 1.1, 1.5, 2.5),  # Back to 1.5" width
            ('income_chart', 2.2, 1.1, 4.8, 2.5),  # Original position
            ('occupation_chart', 7.2, 1.1, 5.6, 2.5),  # Original position

            # Bottom row - Ethnicity chart is 2.8" tall but placed as if 2.5" to maintain alignment
            ('ethnicity_chart', 0.5, 4.2, 4.5, 2.5),  # Will bleed 0.3" below its space
            ('generation_chart', 5.2, 4.2, 4.5, 2.5),  # Original height
            ('children_chart', 9.8, 4.2, 3.0, 2.5)  # Original height
        ]

        for chart_name, left, top, width, height in chart_positions:
            # Try both regular and hires versions
            chart_path = chart_dir / f'{chart_name}_hires.png'
            if not chart_path.exists():
                chart_path = chart_dir / f'{chart_name}.png'

            if chart_path.exists():
                try:
                    pic = slide.shapes.add_picture(
                        str(chart_path),
                        Inches(left), Inches(top),
                        width=Inches(width), height=Inches(height)
                    )
                    logger.info(f"Added {chart_name} to slide: {width}x{height} inches at ({left}, {top})")
                except Exception as e:
                    logger.warning(f"Could not add {chart_name}: {e}")
            else:
                logger.warning(f"Chart not found: {chart_path}")
                self._add_chart_placeholder(slide, chart_name, left, top, width, height)

    def _add_chart_placeholder(self, slide, chart_name: str, left: float, top: float,
                               width: float, height: float):
        """Add a placeholder when chart is not found"""
        placeholder = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = RGBColor(245, 245, 245)
        placeholder.line.color.rgb = RGBColor(200, 200, 200)
        placeholder.line.width = Pt(1)

        # Add placeholder text
        text_box = slide.shapes.add_textbox(
            Inches(left + 0.5), Inches(top + height / 2 - 0.2),
            Inches(width - 1), Inches(0.4)
        )
        text_box.text_frame.text = f"{chart_name.replace('_', ' ').title()} Not Found"
        text_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def _add_legend_box(self, slide, team_name: str, team_short: str, league: str):
        """Add legend with colored squares in a single horizontal row, centered at bottom"""
        # NO KEY LABEL - just the legend items centered at bottom

        # Positioning for horizontal centered layout
        legend_top = Inches(6.9)  # Bottom of slide
        square_size = Inches(0.15)  # Small square size
        text_offset = Inches(0.1)  # Reduced space between square and text
        item_spacing = Inches(0.5)  # Much tighter spacing between items

        # Get colors from team config
        colors = self.team_colors or {
            'primary': '#002244',
            'secondary': '#FFB612',
            'accent': '#8B8B8B'
        }

        # Legend items
        legend_items = [
            (f"{team_name} Fans", colors.get('primary', '#002244')),
            (f"Local Gen Pop (Excluding {team_short} Fans)", colors.get('secondary', '#FFB612')),
            (f"{league} Fans (Excluding {team_short} Fans)", colors.get('accent', '#8B8B8B'))
        ]

        # Center the middle item (Local Gen Pop) on the slide
        slide_width = 13.333  # Standard slide width
        middle_text = legend_items[1][0]

        # Calculate position for middle item to be centered
        # For middle item: square + text should be centered
        middle_item_width = 3.5  # Approximate width of middle text
        middle_square_left = (slide_width - (square_size.inches + text_offset.inches + middle_item_width)) / 2

        # Calculate positions for all three items based on fixed spacing from center
        positions = [
            middle_square_left - item_spacing.inches - square_size.inches - text_offset.inches - 2.5,
            # Left item (2.5" text width)
            middle_square_left,  # Center item
            middle_square_left + square_size.inches + text_offset.inches + middle_item_width + item_spacing.inches
            # Right item
        ]

        # Text widths for each item
        text_widths = [2.5, 3.5, 3.5]

        # Draw each legend item
        for i, ((label, color_hex), square_left, text_width) in enumerate(zip(legend_items, positions, text_widths)):
            # Create colored square - vertically centered with text
            # Text height is 0.25", square is 0.15", so offset square by (0.25 - 0.15) / 2 = 0.05"
            square_vertical_offset = (0.25 - square_size.inches) / 2
            square = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(square_left), legend_top + Inches(square_vertical_offset),
                square_size, square_size
            )
            square.fill.solid()
            square.fill.fore_color.rgb = self._hex_to_rgb(color_hex)
            square.line.fill.background()  # No border on square

            # Add text label - ensure vertical alignment
            text_left = square_left + square_size.inches + text_offset.inches
            text_box = slide.shapes.add_textbox(
                Inches(text_left), legend_top,
                Inches(text_width), Inches(0.25)
            )
            text_frame = text_box.text_frame
            text_frame.text = label
            text_frame.margin_top = Pt(0)
            text_frame.margin_bottom = Pt(0)
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE  # Vertically center text

            p = text_frame.paragraphs[0]
            p.font.name = self.default_font
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(0, 0, 0)  # BLACK text
            p.alignment = PP_ALIGN.LEFT

    def _add_insights_text(self, slide, insight: str):
        """Add insights text at bottom of slide"""
        text_box = slide.shapes.add_textbox(
            Inches(4.5), Inches(6.9),
            Inches(8.5), Inches(0.6)
        )
        text_box.text_frame.text = insight
        text_box.text_frame.word_wrap = True

        p = text_box.text_frame.paragraphs[0]
        p.font.name = self.default_font
        p.font.size = Pt(14)
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT

    def _hex_to_rgb(self, hex_color: str) -> RGBColor:
        """Convert hex color to RGBColor"""
        hex_color = hex_color.lstrip('#')
        return RGBColor(
            int(hex_color[0:2], 16),
            int(hex_color[2:4], 16),
            int(hex_color[4:6], 16)
        )

    def save(self, output_path: Path) -> Path:
        """Save presentation"""
        output_path = Path(output_path)
        self.presentation.save(str(output_path))
        return output_path


# Convenience function
def create_demographics_slide(demographic_data: Dict[str, Any],
                              chart_dir: Path,
                              team_config: Dict[str, Any],
                              presentation: Optional[Presentation] = None) -> Presentation:
    """
    Create a demographics slide

    Args:
        demographic_data: Processed demographic data
        chart_dir: Directory with chart images
        team_config: Team configuration
        presentation: Existing presentation (creates new if None)

    Returns:
        Presentation with demographics slide
    """
    generator = DemographicsSlide(presentation)
    return generator.generate(demographic_data, chart_dir, team_config)