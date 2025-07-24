"""
Test script for category slide generation with mock data
Allows experimentation with formatting without needing database connection
"""

import sys
import os
from pathlib import Path
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import logging

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Add project root to path so we can import our modules
project_root = Path(__file__).parent
sys.path.insert(0, str(project_root))

# Mock the imports we need
class MockCategoryMetrics:
    """Mock CategoryMetrics class to avoid importing the real one"""
    def __init__(self, percent_audience, percent_index, composite_index, spc, ppc):
        self.percent_audience = percent_audience
        self.percent_index = percent_index
        self.composite_index = composite_index
        self.spc = spc
        self.ppc = ppc

    def format_percent_fans(self):
        return f"{self.percent_audience:.1f}%"

    def format_likelihood(self):
        if self.percent_index > 0:
            return f"{int(self.percent_index)}% More likely"
        elif self.percent_index < 0:
            return f"{abs(int(self.percent_index))}% Less likely"
        else:
            return "EQUAL"

    def format_purchases(self):
        if self.ppc > 0:
            return f"{int(self.ppc)}% More purchases"
        elif self.ppc < 0:
            return f"{abs(int(self.ppc))}% Fewer purchases"
        else:
            return "EQUAL"

class MockBaseSlide:
    """Mock BaseSlide class"""
    def __init__(self, presentation=None):
        self.presentation = presentation or Presentation()
        self.default_font = 'Red Hat Display'

    def add_content_slide(self):
        """Add a blank slide"""
        # Use layout 6 which is typically blank
        slide_layout = self.presentation.slide_layouts[6]
        return self.presentation.slides.add_slide(slide_layout)

# Now we can include the actual CategorySlide code
# Copy the entire CategorySlide class from your file here
# (I'll include a simplified version that maintains all the key functionality)

# First, copy all the utility functions from the original file
def format_percentage_no_decimal(value):
    """Format percentage without decimals, handle 0% as EQUAL"""
    if isinstance(value, str):
        try:
            num = float(value.replace('%', '').strip())
        except:
            return value
    else:
        num = float(value)

    if num == 0:
        return "EQUAL"

    return f"{int(round(num))}%"

def format_percent_of_fans(value):
    """Format 'Percent of Fans Who Spend' column with <5% for small values"""
    if isinstance(value, str):
        try:
            num = float(value.replace('%', '').strip())
        except:
            return value
    else:
        num = float(value)

    if num < 5:
        return "<5%"

    return f"{int(round(num))}%"

def format_currency_no_cents(value):
    """Format currency without cents and with commas"""
    if isinstance(value, str):
        try:
            num = float(value.replace('$', '').replace(',', '').strip())
        except:
            return value
    else:
        num = float(value)

    return f"${int(round(num)):,}"

# Import your CategorySlide class here (you'll need to copy the entire class)
# For this example, I'll create a minimal version that focuses on the key elements

from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE

class CategorySlide(MockBaseSlide):
    """Simplified CategorySlide for testing"""

    def __init__(self, presentation=None):
        super().__init__(presentation)

        # Mock LogoManager
        self.logo_manager = None

        # Colors for the slide
        self.colors = {
            'header_bg': RGBColor(240, 240, 240),
            'header_border': RGBColor(200, 200, 200),
            'table_header': RGBColor(0, 0, 0),
            'table_border': RGBColor(0, 0, 0),
            'positive': RGBColor(0, 176, 80),
            'negative': RGBColor(255, 0, 0),
            'equal': RGBColor(184, 134, 11),
            'neutral': RGBColor(0, 0, 0),
            'emerging_bg': RGBColor(217, 217, 217)
        }

    def generate(self, analysis_results, team_config, slide_index=None):
        """Generate the complete category analysis slide"""
        # Extract team info
        team_name = team_config.get('team_name', 'Team')
        team_short = team_config.get('team_name_short', team_name.split()[-1])

        # Check if this is an emerging category
        is_emerging = analysis_results.get('is_emerging', False)

        # Use the content layout
        slide = self.add_content_slide()
        logger.info(f"Added category slide for {analysis_results['display_name']}")

        # Add header
        self._add_header(slide, team_name, analysis_results['slide_title'])

        # Add title
        if is_emerging:
            self._add_emerging_category_title(slide, analysis_results['display_name'])
        else:
            category_title = f"Category Analysis: {analysis_results['display_name']}"
            self._add_title(slide, category_title)

        # Add category insights
        self._add_category_insights(slide, analysis_results, team_short, team_config, is_emerging=is_emerging)

        # Add category metrics table
        self._add_category_table(slide, analysis_results)

        # Add subcategory table
        self._add_subcategory_table(slide, analysis_results['subcategory_stats'])

        return self.presentation

    def generate_brand_slide(self, analysis_results, team_config, slide_index=None):
        """Generate the brand analysis slide (second slide for category)"""
        # Extract team info
        team_name = team_config.get('team_name', 'Team')
        team_short = team_config.get('team_name_short', team_name.split()[-1])
        category_name = analysis_results['display_name']

        # Use the content layout
        slide = self.add_content_slide()
        logger.info(f"Added brand slide for {analysis_results['display_name']}")

        # Brand slide title format
        if category_name.upper() == "QSR":
            brand_title = f"Top QSR Brands for {team_name} Fans"
        else:
            brand_title = f"Top {category_name} Brands for {team_name} Fans"

        # Add header
        header_title = f"Sponsor Spending Analysis: {category_name} Brands"
        self._add_header(slide, team_name, header_title)

        # Add title
        self._add_title(slide, brand_title)

        # Add brand logos (numbered circles)
        self._add_brand_logos(slide, analysis_results['merchant_stats'])

        # Add brand insights (left side)
        self._add_brand_insights(slide, analysis_results, team_name, team_short, category_name)

        # Add brand table (right side)
        self._add_brand_table(slide, analysis_results['merchant_stats'])

        # Add sponsorship recommendation
        self._add_sponsor_recommendation(slide, analysis_results['recommendation'], team_config)

        return self.presentation

    def _add_header(self, slide, team_name, slide_title):
        """Add header with team name and slide title"""
        # Header background
        header_rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(13.333), Inches(0.5)
        )
        header_rect.fill.solid()
        header_rect.fill.fore_color.rgb = self.colors['header_bg']
        header_rect.line.color.rgb = self.colors['header_border']
        header_rect.line.width = Pt(0.5)

        # Team name (left)
        team_text = slide.shapes.add_textbox(
            Inches(0.2), Inches(0.1),
            Inches(3), Inches(0.3)
        )
        team_text.text_frame.text = team_name
        p = team_text.text_frame.paragraphs[0]
        p.font.name = self.default_font
        p.font.size = Pt(14)
        p.font.bold = True

        # Slide indicator (right)
        slide_text = slide.shapes.add_textbox(
            Inches(6.5), Inches(0.1),
            Inches(6.633), Inches(0.3)
        )
        slide_text.text_frame.text = slide_title
        p = slide_text.text_frame.paragraphs[0]
        p.font.name = self.default_font
        p.alignment = PP_ALIGN.RIGHT
        p.font.size = Pt(14)

    def _add_title(self, slide, title):
        """Add main slide title"""
        width = Inches(5.3)

        title_box = slide.shapes.add_textbox(
            Inches(0.6), Inches(1.3),
            width, Inches(1.8)
        )

        text_frame = title_box.text_frame
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.NONE

        # Remove margins
        text_frame.margin_left = Inches(0)
        text_frame.margin_right = Inches(0)
        text_frame.margin_top = Inches(0)
        text_frame.margin_bottom = Inches(0)

        # Split title for category slides
        if "Category Analysis:" in title:
            category_name = title.replace("Category Analysis: ", "")
            text_frame.text = f"Category Analysis:\n{category_name}"
        else:
            text_frame.text = title

        # Format all paragraphs
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = "Red Hat Display"
                run.font.size = Pt(26)
                run.font.bold = True
                run.font.italic = True
                run.font.color.rgb = RGBColor(0, 0, 0)
            paragraph.line_spacing = 1.0

    def _add_emerging_category_title(self, slide, category_name):
        """Add special title formatting for emerging categories with logo"""

        # Add the emerging arrow logo first - need to go up to project root
        # Try multiple possible paths
        possible_paths = [
            # From project root
            Path(__file__).parent.parent.parent / "assets" / "logos" / "general" / "emerging_arrow.png",
            # Relative from current directory
            Path("../../assets/logos/general/emerging_arrow.png"),
            # Direct path
            Path("/Users/travisfleisher/PycharmProjects/PPT_Generator_SIL/assets/logos/general/emerging_arrow.png"),
            # Original path (in case running from different location)
            Path("assets/logos/general/emerging_arrow.png")
        ]

        logo_path = None
        for path in possible_paths:
            if path.exists():
                logo_path = path
                break

        # Extensive debugging for logo path
        logger.info("=" * 60)
        logger.info("LOGO PATH DEBUGGING:")
        logger.info(f"Script file location: {Path(__file__)}")
        logger.info(f"Script parent dir: {Path(__file__).parent}")
        logger.info(f"Project root should be: {Path(__file__).parent.parent.parent}")
        logger.info("Trying paths:")
        for i, path in enumerate(possible_paths):
            logger.info(f"  Path {i+1}: {path}")
            logger.info(f"    - Absolute: {path.absolute()}")
            logger.info(f"    - Exists: {path.exists()}")

        if logo_path:
            logger.info(f"✅ Found logo at: {logo_path}")
        else:
            logger.warning("❌ Logo not found in any of the expected locations")

        logger.info("=" * 60)

        # Check if logo exists and add it
        if logo_path and logo_path.exists():
            try:
                # Add logo at top left, similar to reference image
                pic = slide.shapes.add_picture(
                    str(logo_path),
                    Inches(2.2),  # Left margin
                    Inches(0.8),  # Top position
                    height=Inches(0.4)  # Height - width will scale proportionally
                )
                logger.info("✅ Successfully added emerging arrow logo to slide")
            except Exception as e:
                logger.error(f"❌ Failed to add emerging arrow logo: {e}")
                logger.error(f"Exception type: {type(e).__name__}")
                import traceback
                logger.error(f"Traceback: {traceback.format_exc()}")
        else:
            logger.warning(f"⚠️ Emerging arrow logo not found")
            logger.info("Adding placeholder chevron shape instead")
            # Add a placeholder shape to represent the logo
            logo_placeholder = slide.shapes.add_shape(
                MSO_SHAPE.CHEVRON,
                Inches(0.6), Inches(0.8),
                Inches(0.5), Inches(0.5)
            )
            logo_placeholder.fill.solid()
            logo_placeholder.fill.fore_color.rgb = RGBColor(50, 50, 50)
            logo_placeholder.line.color.rgb = RGBColor(0, 0, 0)
            logger.info("✅ Added chevron placeholder for missing logo")

        # Main title: "Emerging Category:"
        title_box = slide.shapes.add_textbox(
            Inches(0.6), Inches(1.3),
            Inches(5.3), Inches(0.5)
        )

        text_frame = title_box.text_frame
        text_frame.text = "Emerging Category:"

        p = text_frame.paragraphs[0]
        p.font.name = "Red Hat Display"
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.italic = True
        p.font.color.rgb = RGBColor(0, 0, 0)

        # Category name on second line
        category_box = slide.shapes.add_textbox(
            Inches(0.6), Inches(1.725),
            Inches(5.3), Inches(0.6)
        )

        text_frame = category_box.text_frame
        text_frame.text = category_name

        p = text_frame.paragraphs[0]
        p.font.name = "Red Hat Display"
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.italic = True
        p.font.color.rgb = RGBColor(0, 0, 0)

        # Explanatory subtext
        subtext_box = slide.shapes.add_textbox(
            Inches(0.6), Inches(2.3),
            Inches(5.3), Inches(0.6)
        )

        text_frame = subtext_box.text_frame
        text_frame.word_wrap = True
        text_frame.text = (
            "The Emerging Category is where at least 10% of your fans are "
            "spending, but where there isn't one clear brand leader, and the "
            "category has room to grow"
        )

        p = text_frame.paragraphs[0]
        p.font.name = "Red Hat Display"
        p.font.size = Pt(12)
        p.font.italic = True
        p.font.color.rgb = RGBColor(0, 0, 0)

    def _add_category_insights(self, slide, results, team_short, team_config, is_emerging=False):
        """Add category insights section"""
        # Adjust vertical position based on whether it's an emerging category
        title_y = Inches(3.2) if is_emerging else Inches(2.6)
        insights_y = Inches(3.6) if is_emerging else Inches(3.0)
        nba_label_y = Inches(5.8) if is_emerging else Inches(5.2)
        nba_box_y = Inches(6.2) if is_emerging else Inches(5.6)

        # Insights title
        insights_title = slide.shapes.add_textbox(
            Inches(0.5), title_y,
            Inches(4), Inches(0.3)
        )
        insights_title.text_frame.text = "Category Insights:"
        p = insights_title.text_frame.paragraphs[0]
        p.font.name = self.default_font
        p.font.size = Pt(14)
        p.font.bold = True

        # Regular insights box
        insights_box = slide.shapes.add_textbox(
            Inches(0.7), insights_y,
            Inches(4.5), Inches(3.0)
        )

        text_frame = insights_box.text_frame
        text_frame.word_wrap = True

        # Add regular insights
        regular_insights = results['insights'][:4]
        for i, insight in enumerate(regular_insights):
            p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
            p.text = f"• {insight}"
            p.font.name = self.default_font
            p.font.size = Pt(12)
            p.line_spacing = 1.2

        # Add NBA comparison section if needed
        if len(results['insights']) > 4:
            # NBA comparison label
            nba_label = slide.shapes.add_textbox(
                Inches(0.5), nba_label_y,
                Inches(4), Inches(0.3)
            )
            league = team_config.get('league', 'NBA')
            nba_label.text_frame.text = f"{team_short} Fans vs. {league} Fans:"
            p = nba_label.text_frame.paragraphs[0]
            p.font.name = self.default_font
            p.font.size = Pt(14)
            p.font.bold = True

            # NBA insights box
            nba_box = slide.shapes.add_textbox(
                Inches(0.7), nba_box_y,
                Inches(4.5), Inches(1.2)
            )

            nba_text_frame = nba_box.text_frame
            nba_text_frame.word_wrap = True

            # Add NBA insights
            for i, insight in enumerate(results['insights'][4:6]):
                p = nba_text_frame.add_paragraph() if i > 0 else nba_text_frame.paragraphs[0]
                p.text = f"• {insight}"
                p.font.name = self.default_font
                p.font.size = Pt(12)
                p.line_spacing = 1.2

    def _add_category_table(self, slide, results):
        """Add category metrics table"""
        metrics = results['category_metrics']
        is_emerging = results.get('is_emerging', False)

        # Table position
        left = Inches(6.2)
        top = Inches(1.4)
        width = Inches(6.8)

        # Table height
        row_height = Inches(0.7)
        num_rows = 2
        height = row_height * num_rows

        # Create table
        table_shape = slide.shapes.add_table(2, 4, left, top, width, height)
        table = table_shape.table

        # Set column widths
        table.columns[0].width = Inches(1.6)
        table.columns[1].width = Inches(1.5)
        table.columns[2].width = Inches(1.7)
        table.columns[3].width = Inches(1.6)

        # Header row
        headers = ['Category', 'Percent of Fans\nWho Spend', 'Likelihood to Spend\n(vs. Local Gen Pop)',
                   'Purchases Per Fan\n(vs. Local Gen Pop)']

        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            self._format_header_cell(cell)

        # Data row
        category_name = results.get('display_name', 'Category')
        table.cell(1, 0).text = category_name

        percent_value = metrics.format_percent_fans()
        table.cell(1, 1).text = format_percent_of_fans(percent_value)

        table.cell(1, 2).text = metrics.format_likelihood()
        table.cell(1, 3).text = metrics.format_purchases()

        # Format data cells
        for i in range(4):
            cell = table.cell(1, i)
            self._format_data_cell(cell)

            if is_emerging and i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = self.colors['emerging_bg']

    def _add_subcategory_table(self, slide, subcategory_stats):
        """Add subcategory statistics table"""
        if subcategory_stats.empty:
            return

        # Table position
        left = Inches(6.2)
        top = Inches(3.0)
        width = Inches(6.8)

        # Create table
        rows = min(len(subcategory_stats), 5) + 1
        row_height = Inches(0.7)
        table_height = row_height * rows

        table_shape = slide.shapes.add_table(rows, 4, left, top, width, table_height)
        table = table_shape.table

        # Set column widths
        table.columns[0].width = Inches(1.6)
        table.columns[1].width = Inches(1.5)
        table.columns[2].width = Inches(1.7)
        table.columns[3].width = Inches(1.6)

        # Headers
        headers = ['Sub-Category', 'Percent of Fans\nWho Spend', 'Likelihood to Spend\n(vs. Local Gen Pop)',
                   'Purchases Per Fan\n(vs. Local Gen Pop)']

        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            self._format_header_cell(cell)

        # Data rows
        for row_idx, (_, row) in enumerate(subcategory_stats.iterrows(), 1):
            if row_idx >= rows:
                break

            table.cell(row_idx, 0).text = row['Subcategory']
            table.cell(row_idx, 1).text = format_percent_of_fans(row['Percent of Fans Who Spend'])
            table.cell(row_idx, 2).text = row['Likelihood to spend (vs. Local Gen Pop)']
            table.cell(row_idx, 3).text = row['Purchases Per Fan (vs. Gen Pop)']

            for col in range(4):
                self._format_data_cell(table.cell(row_idx, col))

        # Add explanatory text
        explanation_top = top + table_height + Inches(0.1)
        explanation_box = slide.shapes.add_textbox(
            left, explanation_top,
            width, Inches(0.3)
        )

        text_frame = explanation_box.text_frame
        text_frame.word_wrap = True
        text_frame.text = "Subcategories shown in descending order by composite index"

        p = text_frame.paragraphs[0]
        p.font.name = self.default_font
        p.font.size = Pt(10)
        p.font.italic = True
        p.font.color.rgb = RGBColor(100, 100, 100)
        p.alignment = PP_ALIGN.LEFT

    def _format_header_cell(self, cell, small=False):
        """Format table header cell"""
        cell.fill.solid()
        cell.fill.fore_color.rgb = self.colors['table_header']

        text_frame = cell.text_frame
        text_frame.margin_left = Inches(0.05)
        text_frame.margin_right = Inches(0.05)
        text_frame.margin_top = Inches(0.03)
        text_frame.margin_bottom = Inches(0.03)
        text_frame.word_wrap = True

        for paragraph in text_frame.paragraphs:
            paragraph.font.name = self.default_font
            paragraph.font.size = Pt(8) if small else Pt(10)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.line_spacing = 1.0

        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    def _add_brand_logos(self, slide, merchant_stats):
        """Add brand logos as numbered circles"""
        merchant_df, top_merchants = merchant_stats

        if merchant_df.empty:
            return

        # Align with brand table dimensions
        table_left = Inches(6.6)
        table_width = Inches(6.0)

        # Calculate logo positions
        num_logos = min(5, len(merchant_df))
        display_size = Inches(1.0)

        # Calculate spacing
        available_space = table_width - (display_size * num_logos)
        spacing_between = available_space / (num_logos - 1) if num_logos > 1 else 0

        # Position for logos
        y = Inches(1.2)

        # Add numbered circles for top 5 brands
        for i in range(num_logos):
            x = table_left + (i * (display_size + spacing_between))

            # Add circle
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                x, y,
                display_size, display_size
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = RGBColor(255, 255, 255)
            circle.line.color.rgb = RGBColor(200, 200, 200)
            circle.line.width = Pt(1)

            # Add ranking number below each logo
            number_box = slide.shapes.add_textbox(
                x, y + display_size + Inches(0.05),
                display_size, Inches(0.3)
            )
            number_box.text_frame.text = str(i + 1)
            p = number_box.text_frame.paragraphs[0]
            p.font.name = self.default_font
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = RGBColor(100, 100, 100)

    def _add_brand_insights(self, slide, results, team_name, team_short, category_name):
        """Add brand-specific insights"""
        # Top Brand Insights section
        insights_title = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.4),
            Inches(4), Inches(0.3)
        )
        insights_title.text_frame.text = "Top Brand Insights"
        p = insights_title.text_frame.paragraphs[0]
        p.font.name = self.default_font
        p.font.size = Pt(14)
        p.font.bold = True

        # Insights box
        insights_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(2.8),
            Inches(5.5), Inches(1.6)
        )

        text_frame = insights_box.text_frame
        text_frame.word_wrap = True

        # Get merchant insights
        merchant_insights = results.get('merchant_insights', [])

        # Format insights with labels
        labels = ["Highest % of Fans: ", "Most Purchases per Fan: ", "Highest Spend per Fan: ", "Highest % of Fans Index vs NBA: "]

        for i, (label, insight) in enumerate(zip(labels, merchant_insights)):
            p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
            p.line_spacing = 1.2

            # Add bullet and label in bold
            run1 = p.add_run()
            run1.text = f"• {label}"
            run1.font.name = self.default_font
            run1.font.size = Pt(12)
            run1.font.bold = True

            # Add the insight text
            run2 = p.add_run()
            run2.text = insight
            run2.font.name = self.default_font
            run2.font.size = Pt(12)
            run2.font.bold = False

    def _add_brand_table(self, slide, merchant_stats):
        """Add brand ranking table"""
        merchant_df, _ = merchant_stats

        if merchant_df.empty:
            return

        # Create table
        rows = min(len(merchant_df) + 1, 6)
        cols = 5
        left = Inches(6.6)
        top = Inches(3.0)
        width = Inches(6.0)

        row_height = Inches(0.6)
        height = row_height * rows

        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table

        # Column widths
        table.columns[0].width = Inches(1.0)
        table.columns[1].width = Inches(1.4)
        table.columns[2].width = Inches(1.2)
        table.columns[3].width = Inches(1.2)
        table.columns[4].width = Inches(1.2)

        # Headers
        headers = ['Rank (by\npercent of\nfans who\nspend)', 'Brand',
                   'Percent of\nFans Who\nSpend',
                   'Likelihood to Spend\n(vs. Local Gen Pop)',
                   'Purchases Per Fan\n(vs. Local Gen Pop)']

        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            self._format_header_cell(cell, small=True)

        # Data rows
        for row_idx, (_, row) in enumerate(merchant_df.iterrows(), 1):
            if row_idx >= rows:
                break

            table.cell(row_idx, 0).text = str(row['Rank'])
            table.cell(row_idx, 1).text = row['Brand']
            table.cell(row_idx, 2).text = format_percent_of_fans(row['Percent of Fans Who Spend'])
            table.cell(row_idx, 3).text = row['Likelihood to spend (vs. Local Gen Pop)']
            table.cell(row_idx, 4).text = row['Purchases Per Fan (vs. Gen Pop)']

            for col in range(5):
                self._format_data_cell(table.cell(row_idx, col), small=True)

    def _add_sponsor_recommendation(self, slide, recommendation, team_config):
        """Add Hot Brand Target and sponsorship recommendation"""
        if not recommendation:
            return

        # Hot Brand Target header
        target_title = slide.shapes.add_textbox(
            Inches(0.5), Inches(5.0),
            Inches(2.0), Inches(0.3)
        )
        target_title.text_frame.text = "Hot Brand Target:"
        p = target_title.text_frame.paragraphs[0]
        p.font.name = self.default_font
        p.font.size = Pt(14)
        p.font.bold = True

        # Add X logo placeholder
        logo_size = Inches(0.5)
        logo_x = Inches(2.4)
        logo_y = Inches(4.925)

        # Add circle for logo
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            logo_x, logo_y,
            logo_size, logo_size
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = RGBColor(255, 255, 255)
        circle.line.color.rgb = RGBColor(200, 200, 200)
        circle.line.width = Pt(0.5)

        # Add X in the circle
        x_text = slide.shapes.add_textbox(
            logo_x, logo_y,
            logo_size, logo_size
        )
        x_text.text_frame.text = "X"
        p = x_text.text_frame.paragraphs[0]
        p.font.name = self.default_font
        p.font.size = Pt(24)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        p.font.color.rgb = RGBColor(0, 0, 200)
        # Center vertically
        x_text.text_frame.margin_top = Inches(0.08)

        # Recommendation content
        rec_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(5.5),
            Inches(5.5), Inches(1.2)
        )

        text_frame = rec_box.text_frame
        text_frame.word_wrap = True

        # Format recommendation
        team_short = team_config.get('team_name_short', 'Team')
        merchant_name = recommendation.get('merchant', 'Brand')
        composite_index = int(recommendation.get('composite_index', 0))

        # First bullet with bold brand name
        p1 = text_frame.paragraphs[0]

        run1 = p1.add_run()
        run1.text = f"• The {team_short} should target "
        run1.font.name = self.default_font
        run1.font.size = Pt(12)
        run1.font.bold = False

        run2 = p1.add_run()
        run2.text = merchant_name
        run2.font.name = self.default_font
        run2.font.size = Pt(12)
        run2.font.bold = True

        run3 = p1.add_run()
        run3.text = f" for a sponsorship based on having the highest composite index of {composite_index}"
        run3.font.name = self.default_font
        run3.font.size = Pt(12)
        run3.font.bold = False

        p1.line_spacing = 1.2

        # Second bullet
        p2 = text_frame.add_paragraph()
        p2.text = "• The composite index indicates a brand with significant likelihood for more fans to be spending more frequently, and at a higher spend per fan vs. other brands"
        p2.font.name = self.default_font
        p2.font.size = Pt(12)
        p2.line_spacing = 1.2
        p2.left_indent = Inches(0.25)
        p2.first_line_indent = Inches(0)

    def _format_data_cell(self, cell, small=False):
        """Format table data cell"""
        text_frame = cell.text_frame
        text_frame.margin_left = Inches(0.05)
        text_frame.margin_right = Inches(0.05)
        text_frame.margin_top = Inches(0.03)
        text_frame.margin_bottom = Inches(0.03)
        text_frame.word_wrap = True

        for paragraph in text_frame.paragraphs:
            paragraph.font.name = self.default_font
            paragraph.font.size = Pt(12)
            paragraph.alignment = PP_ALIGN.CENTER

            # Color coding
            if 'More' in cell.text or 'more' in cell.text:
                paragraph.font.color.rgb = self.colors['positive']
                paragraph.font.bold = True
            elif 'Less' in cell.text or 'less' in cell.text or 'fewer' in cell.text:
                paragraph.font.color.rgb = self.colors['negative']
                paragraph.font.bold = False
            elif 'EQUAL' in cell.text or 'Equal' in cell.text:
                paragraph.font.color.rgb = self.colors['equal']
                paragraph.font.bold = False
            else:
                paragraph.font.bold = False

        cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def create_mock_data(is_emerging=False):
    """Create mock data for testing"""

    if is_emerging:
        # Mock data for emerging category (Sportstainment)
        subcategory_data = pd.DataFrame({
            'Subcategory': ['Pickleball & Racket Sports', 'Indoor Entertainment', 'Golf Resorts', 'Ski Resorts'],
            'Percent of Fans Who Spend': [4.5, 5.0, 10.0, 9.0],
            'Likelihood to spend (vs. Local Gen Pop)': ['614% More likely', '484% More likely', '423% More likely', '382% More likely'],
            'Purchases Per Fan (vs. Gen Pop)': ['18% More purchases', '19% More purchases', '52% More purchases', '38% More purchases']
        })

        # Mock merchant/brand data for emerging category
        merchant_data = pd.DataFrame({
            'Rank': [1, 2, 3, 4, 5],
            'Brand': ['TopGolf', 'iFLY', 'Life Time Fitness', 'Main Event', 'Dave & Busters'],
            'Percent of Fans Who Spend': [3.2, 2.8, 2.1, 1.9, 1.5],
            'Likelihood to spend (vs. Local Gen Pop)': ['108% More likely', '216% More likely', '89% More likely', '124% More likely', '156% More likely'],
            'Purchases Per Fan (vs. Gen Pop)': ['12% More purchases', '25% More purchases', '18% More purchases', '8% More purchases', '15% More purchases']
        })

        # Mock category metrics for emerging
        category_metrics = MockCategoryMetrics(
            percent_audience=20.0,
            percent_index=370,
            composite_index=125,
            spc=52,
            ppc=52
        )

        # Mock recommendation for emerging
        recommendation = {
            'merchant': 'Life Time Fitness',
            'composite_index': 189
        }

        # Mock analysis results for emerging category
        analysis_results = {
            'display_name': 'Sportstainment',
            'slide_title': 'Sportstainment Sponsor Analysis',
            'is_emerging': True,  # This is the emerging category
            'category_metrics': category_metrics,
            'subcategory_stats': subcategory_data,
            'merchant_stats': (merchant_data, merchant_data['Brand'].tolist()),
            'insights': [
                "Jazz fans are 370% more likely to spend on Sportstainment than the local gen pop",
                "Jazz fans make an average of 52% more purchases per fan on Sportstainment than the local gen pop",
                "Jazz fans are more than 7.1X more likely to spend on Pickleball & Racket Sports vs. the local gen pop",
                "Jazz fans spend an average of $469 per fan per year on Ski Resorts",
                "Utah Jazz fans are 89% more likely to spend on Ski Resorts when compared to the NBA average"
            ],
            'merchant_insights': [
                "3% of Jazz fans spend at TopGolf",
                "Jazz fans make an average of 6 purchases per year at Life Time Fitness",
                "Jazz fans spent an average of $892 per fan on Life Time Fitness per year",
                "Jazz fans are 216% more likely to spend on iFLY than the average NBA fan"
            ],
            'recommendation': recommendation
        }

        # Mock team config
        team_config = {
            'team_name': 'Utah Jazz',
            'team_name_short': 'Jazz',
            'league': 'NBA'
        }

    else:
        # Original mock data for regular category
        subcategory_data = pd.DataFrame({
            'Subcategory': ['Hotels', 'Airlines', 'Car Rentals', 'Vacation Rentals', 'Cruises'],
            'Percent of Fans Who Spend': [22.5, 18.3, 15.7, 12.4, 8.9],
            'Likelihood to spend (vs. Local Gen Pop)': ['15% More likely', '12% More likely', '8% More likely', 'EQUAL', '5% Less likely'],
            'Purchases Per Fan (vs. Gen Pop)': ['20% More purchases', '15% More purchases', '10% More purchases', '5% More purchases', 'EQUAL']
        })

        # Mock merchant/brand data
        merchant_data = pd.DataFrame({
            'Rank': [1, 2, 3, 4, 5],
            'Brand': ['American Airlines', 'Delta', 'United Airlines', 'Budget', 'Allegiant Air'],
            'Percent of Fans Who Spend': [8.0, 4.8, 4.2, 3.5, 2.9],
            'Likelihood to spend (vs. Local Gen Pop)': ['16% Less likely', '13% Less likely', '9% More likely', '62% More likely', '108% More likely'],
            'Purchases Per Fan (vs. Gen Pop)': ['47% Less purchases', '43% Less purchases', '55% Less purchases', '18% More purchases', '9% More purchases']
        })

        # Mock category metrics
        category_metrics = MockCategoryMetrics(
            percent_audience=35.2,
            percent_index=18,
            composite_index=125,
            spc=22,
            ppc=15
        )

        # Mock recommendation
        recommendation = {
            'merchant': 'Celebrity Cruises',
            'composite_index': 194
        }

        # Mock analysis results
        analysis_results = {
            'display_name': 'Travel',
            'slide_title': 'Sponsor Spending Analysis: Travel',
            'is_emerging': False,
            'category_metrics': category_metrics,
            'subcategory_stats': subcategory_data,
            'merchant_stats': (merchant_data, merchant_data['Brand'].tolist()),
            'insights': [
                "Panthers fans spend 22% more on travel than the local gen pop",
                "Hotels represent the highest subcategory with 23% of fans spending",
                "Panthers fans make 15% more travel purchases per year compared to local gen pop",
                "Airlines show strong growth potential with 18% fan engagement",
                "Panthers fans are 25% more likely to spend on travel compared to NBA fans",
                "Average travel spend per Panthers fan is $2,450 annually vs NBA average of $1,960"
            ],
            'merchant_insights': [
                "8% of Carolina Panthers fans spend at American Airlines",
                "Panthers fans make an average of 4 purchases per year at Allegiant Air",
                "Panthers fans spent an average of $1,499 per fan on Allegiant Air per year",
                "Panthers fans are 125% more likely to spend on American Airlines than the average NBA fan"
            ],
            'recommendation': recommendation
        }

        # Mock team config
        team_config = {
            'team_name': 'Carolina Panthers',
            'team_name_short': 'Panthers',
            'league': 'NFL'
        }

    return analysis_results, team_config


def main():
    """Main function to generate test slides"""
    print("Generating test category and brand slides...")

    # Set up more detailed logging for debugging
    logging.basicConfig(
        level=logging.DEBUG,
        format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
    )

    # Log current environment
    logger.info("=" * 60)
    logger.info("ENVIRONMENT INFO:")
    logger.info(f"Script location: {Path(__file__).absolute()}")
    logger.info(f"Current working directory: {Path.cwd()}")
    logger.info(f"Python executable: {sys.executable}")
    logger.info("=" * 60)

    # Ensure the assets directory structure exists
    assets_path = Path("assets/logos/general")
    logger.info(f"Creating directory structure: {assets_path}")
    assets_path.mkdir(parents=True, exist_ok=True)

    # Check what's in the directory
    logger.info("Checking assets directory contents:")
    if assets_path.exists():
        logger.info(f"Directory {assets_path} exists")
        contents = list(assets_path.iterdir())
        if contents:
            logger.info(f"Contents of {assets_path}:")
            for item in contents:
                logger.info(f"  - {item.name}")
        else:
            logger.info(f"Directory {assets_path} is empty")
    else:
        logger.warning(f"Directory {assets_path} does not exist after creation attempt!")

    # Create presentation
    presentation = Presentation()

    # Create slide generator
    generator = CategorySlide(presentation)

    print("\n1. Creating REGULAR category slides (Travel)...")
    # Create regular category data
    regular_data, regular_team = create_mock_data(is_emerging=False)

    # Generate regular category analysis slide
    presentation = generator.generate(regular_data, regular_team)

    # Generate regular brand analysis slide
    presentation = generator.generate_brand_slide(regular_data, regular_team)

    print("\n2. Creating EMERGING category slides (Sportstainment)...")
    # Create emerging category data
    emerging_data, emerging_team = create_mock_data(is_emerging=True)

    # Generate emerging category analysis slide
    presentation = generator.generate(emerging_data, emerging_team)

    # Generate emerging brand analysis slide
    presentation = generator.generate_brand_slide(emerging_data, emerging_team)

    # Save the presentation
    output_path = "test_category_slides_with_emerging.pptx"
    presentation.save(output_path)

    print(f"\n✅ Test slides saved to: {output_path}")
    print("\nThe presentation now contains 4 slides:")
    print("  - Slide 1: Regular Category Analysis (Travel)")
    print("  - Slide 2: Regular Brand Analysis (Travel)")
    print("  - Slide 3: EMERGING Category Analysis (Sportstainment)")
    print("  - Slide 4: EMERGING Brand Analysis (Sportstainment)")

    print("\nNOTE: Check the console output above for detailed logo path debugging.")
    print("      If the logo file is not found, make sure to place it at:")
    print(f"      {Path('assets/logos/general/emerging_arrow.png').absolute()}")

    print("\nYou can now:")
    print("1. Open the PowerPoint file to see all slides")
    print("2. Compare regular vs emerging category formatting")
    print("3. Check if the logo appears on the emerging category slide")
    print("4. Adjust formatting parameters as needed")


if __name__ == "__main__":
    main()