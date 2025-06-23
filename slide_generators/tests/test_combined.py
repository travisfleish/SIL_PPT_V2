# test_combined_slides.py
"""
Test script to create a PowerPoint with Demographics and Behaviors slides
"""

import sys
from pathlib import Path

sys.path.append(str(Path(__file__).parent))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Import data processors
from data_processors.merchant_ranker import MerchantRanker
from data_processors.demographic_processor import DemographicsProcessor
from data_processors.snowflake_connector import query_to_dataframe
from utils.team_config_manager import TeamConfigManager

# Import slide generators
from slide_generators.behaviors_slide import BehaviorsSlide
# Remove the old import - we'll import it where needed

# Import visualizations
from visualizations.demographic_charts import DemographicCharts

import logging

logging.basicConfig(level=logging.INFO)


def create_title_slide(presentation, team_config):
    """Create a simple title slide"""
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2),
        Inches(8), Inches(1.5)
    )
    title_box.text = f"{team_config['team_name']}\nSponsorship Insights Report"

    # Format title
    for i, paragraph in enumerate(title_box.text_frame.paragraphs):
        paragraph.font.size = Pt(40 if i == 0 else 32)
        paragraph.font.bold = True
        paragraph.alignment = PP_ALIGN.CENTER

    return slide


def test_combined_presentation(team_key: str = 'utah_jazz'):
    """
    Create a PowerPoint with Demographics and Behaviors slides

    Args:
        team_key: Team identifier (utah_jazz or dallas_cowboys)
    """
    print(f"\n{'=' * 60}")
    print(f"COMBINED PRESENTATION TEST - {team_key.replace('_', ' ').title()}")
    print(f"{'=' * 60}")

    try:
        # 1. Setup
        print("\n1. Loading configuration...")
        config_manager = TeamConfigManager()
        team_config = config_manager.get_team_config(team_key)
        print(f"   ✅ Team: {team_config['team_name']}")

        # 2. Create presentation
        print("\n2. Creating presentation...")
        presentation = Presentation()

        # 3. Add title slide
        print("\n3. Adding title slide...")
        create_title_slide(presentation, team_config)
        print("   ✅ Title slide added")

        # 4. Add demographics slide
        print("\n4. Creating demographics slide...")
        try:
            # Fetch demographic data
            print("   - Fetching demographic data from Snowflake...")
            demographics_view = config_manager.get_view_name(team_key, 'demographics')
            query = f"SELECT * FROM {demographics_view}"
            df = query_to_dataframe(query)

            if not df.empty:
                print(f"   - Processing {len(df):,} demographic records...")

                # Process demographics
                processor = DemographicsProcessor(
                    data_source=df,
                    team_name=team_config['team_name'],
                    league=team_config['league']
                )
                demographic_data = processor.process_all_demographics()

                # Generate charts
                print("   - Generating demographic charts...")
                output_dir = Path(f"{team_key}_temp_charts")
                output_dir.mkdir(exist_ok=True)

                charter = DemographicCharts(team_colors=team_config.get('colors'))
                charts = charter.create_all_demographic_charts(
                    demographic_data,
                    output_dir=output_dir
                )

                # Create demographics slide using the proper generator
                from slide_generators.demographics_slide import DemographicsSlide
                demo_generator = DemographicsSlide(presentation)
                presentation = demo_generator.generate(
                    demographic_data=demographic_data,
                    chart_dir=output_dir,
                    team_config=team_config
                )

                print("   ✅ Demographics slide added with all charts")
            else:
                print("   ⚠️  No demographic data found, adding placeholder...")
                # Add placeholder slide
                slide = presentation.slides.add_slide(presentation.slide_layouts[5])
                text_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
                text_box.text = "Demographics data not available"

        except Exception as e:
            print(f"   ⚠️  Could not create demographics slide: {e}")
            import traceback
            traceback.print_exc()
            # Add placeholder slide
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])
            text_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
            text_box.text = "Demographics slide - error loading data"

        # 5. Add behaviors slide
        print("\n5. Creating behaviors slide...")
        try:
            # Initialize merchant ranker
            ranker = MerchantRanker(team_view_prefix=team_config['view_prefix'])

            # Create behaviors slide
            behaviors_generator = BehaviorsSlide(presentation)
            presentation = behaviors_generator.generate(ranker, team_config)

            print("   ✅ Behaviors slide added")

        except Exception as e:
            print(f"   ⚠️  Could not create behaviors slide: {e}")
            import traceback
            traceback.print_exc()

        # 6. Save presentation
        output_path = Path(f"{team_key}_combined_presentation.pptx")
        presentation.save(str(output_path))

        print(f"\n✅ SUCCESS! Presentation saved to: {output_path.absolute()}")
        print(f"\nPresentation contains:")
        print(f"  • Slide 1: Title slide")
        print(f"  • Slide 2: Demographics overview")
        print(f"  • Slide 3: Fan behaviors (wheel + chart)")

        return output_path

    except Exception as e:
        print(f"\n❌ ERROR: {type(e).__name__}: {str(e)}")
        import traceback
        traceback.print_exc()
        return None


def main():
    """Main test function"""
    print("\n🎯 COMBINED PRESENTATION GENERATOR TEST")
    print("This will create a PowerPoint with:")
    print("  1. Title slide")
    print("  2. Demographics slide")
    print("  3. Fan behaviors slide")

    # Test with Utah Jazz
    jazz_result = test_combined_presentation('utah_jazz')

    if jazz_result:
        # Ask if user wants to test Dallas Cowboys too
        user_input = input("\n\nAlso create for Dallas Cowboys? (y/n): ")
        if user_input.lower() == 'y':
            cowboys_result = test_combined_presentation('dallas_cowboys')

    print("\n✨ Test complete!")
    print("\nNote: The demographics slide currently shows a placeholder.")
    print("To see the full demographic charts, check the generated chart directory.")


if __name__ == "__main__":
    main()