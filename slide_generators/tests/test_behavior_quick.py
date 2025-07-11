# behaviors_slide_fresh_test.py
"""
FRESH behaviors slide test script with correct positioning to match reference
Starting completely from scratch to ensure proper layout
"""

from pathlib import Path
from typing import Dict, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import matplotlib.pyplot as plt
import numpy as np


class FreshBehaviorsSlideTest:
    """Fresh test class with correct positioning from scratch"""

    def __init__(self):
        self.default_font = "Arial"

    def create_mock_fan_wheel(self, output_path: Path) -> Path:
        """Create mock fan wheel matching the real design"""
        fig, ax = plt.subplots(figsize=(8, 8), facecolor='white')

        # Navy blue background circle
        circle_bg = plt.Circle((0.5, 0.5), 0.45, color='#002244', transform=ax.transAxes)
        ax.add_patch(circle_bg)

        # Black center circle
        center_circle = plt.Circle((0.5, 0.5), 0.15, color='black', transform=ax.transAxes)
        ax.add_patch(center_circle)

        # Center text
        ax.text(0.5, 0.52, 'THE JAZZ', transform=ax.transAxes, ha='center', va='center',
                fontsize=12, fontweight='bold', color='white')
        ax.text(0.5, 0.48, 'FAN', transform=ax.transAxes, ha='center', va='center',
                fontsize=10, color='white')

        # White spokes
        angles = np.linspace(0, 2 * np.pi, 12, endpoint=False)
        for angle in angles:
            x_start = 0.5 + 0.15 * np.cos(angle)
            y_start = 0.5 + 0.15 * np.sin(angle)
            x_end = 0.5 + 0.42 * np.cos(angle)
            y_end = 0.5 + 0.42 * np.sin(angle)
            ax.plot([x_start, x_end], [y_start, y_end], 'white', linewidth=2, transform=ax.transAxes)

        # Brand circles around perimeter
        brand_positions = [
            (0.5, 0.85), (0.72, 0.75), (0.85, 0.5), (0.72, 0.25),
            (0.5, 0.15), (0.28, 0.25), (0.15, 0.5), (0.28, 0.75)
        ]

        for x, y in brand_positions:
            # White circles for brands
            brand_circle = plt.Circle((x, y), 0.04, color='white', transform=ax.transAxes)
            ax.add_patch(brand_circle)
            # Yellow connector dots
            connector_circle = plt.Circle((x, y), 0.015, color='#FFD700', transform=ax.transAxes)
            ax.add_patch(connector_circle)

        ax.set_xlim(0, 1)
        ax.set_ylim(0, 1)
        ax.set_aspect('equal')
        ax.axis('off')

        plt.tight_layout()
        plt.savefig(output_path, dpi=150, bbox_inches='tight', facecolor='white', pad_inches=0.1)
        plt.close()

        return output_path

    def create_mock_community_chart(self, output_path: Path) -> Path:
        """Create mock community chart matching the real design"""
        fig, ax = plt.subplots(figsize=(10, 6), facecolor='white')

        # Data matching reference
        communities = ['Live Entertainment Seekers', 'Movie Buffs', 'Theme Parkers',
                       'Fitness Enthusiasts', 'Casual Outdoor Enthusiasts', 'Sports Merchandise Shopper',
                       'College Sports', 'Trend Setters', 'Youth Sports', 'Gambler']

        scores = [59.7, 34.4, 32.8, 29.5, 24.8, 22.6, 22.5, 21.3, 21.1, 20.9]

        # Horizontal bar chart
        y_pos = np.arange(len(communities))
        bars = ax.barh(y_pos, scores, color='#4472C4', alpha=0.8, height=0.7)

        # Value labels
        for i, (bar, score) in enumerate(zip(bars, scores)):
            width = bar.get_width()
            ax.text(width + 1, bar.get_y() + bar.get_height() / 2,
                    f'{score}%', ha='left', va='center', fontsize=10, fontweight='bold')

        # Style to match reference
        ax.set_yticks(y_pos)
        ax.set_yticklabels(communities, fontsize=11)
        ax.invert_yaxis()

        # Remove spines
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        ax.spines['left'].set_visible(False)

        # X-axis
        ax.set_xlim(0, max(scores) + 10)
        ax.set_xticks([0, 100, 200, 300, 400, 500, 600, 700])
        ax.set_xlabel('Percent Fan Audience', fontsize=12, fontweight='bold')
        ax.tick_params(axis='x', labelsize=10)
        ax.tick_params(axis='y', left=False)

        # Grid
        ax.grid(axis='x', alpha=0.3, linestyle='-', linewidth=0.5)
        ax.set_axisbelow(True)

        # Legend
        legend_elements = [
            plt.Rectangle((0, 0), 1, 1, facecolor='#4472C4', alpha=0.8, label='% Team Fans'),
            plt.Rectangle((0, 0), 1, 1, facecolor='#FFD966', alpha=0.8, label='Team Fan Index')
        ]
        ax.legend(handles=legend_elements, loc='lower right', frameon=False, fontsize=10)

        plt.tight_layout()
        plt.savefig(output_path, dpi=150, bbox_inches='tight', facecolor='white')
        plt.close()

        return output_path

    def create_test_slide(self) -> Presentation:
        """Create test slide with CORRECT positioning"""
        # New presentation
        pres = Presentation()
        pres.slide_width = Inches(13.333)
        pres.slide_height = Inches(7.5)

        # Blank slide
        slide = pres.slides.add_slide(pres.slide_layouts[6])

        # Create mock images
        temp_dir = Path('temp_fresh_images')
        temp_dir.mkdir(exist_ok=True)

        fan_wheel_path = self.create_mock_fan_wheel(temp_dir / 'fresh_fan_wheel.png')
        community_chart_path = self.create_mock_community_chart(temp_dir / 'fresh_community_chart.png')

        # CALCULATION: Left side vertical centering (with reduced explanation height)
        # Available space: 7.5 - 0.5 = 7.0"
        # Total left content height: 1.0 (insight) + 0.3 (title) + 3.5 (chart) + 0.8 (explanation) = 5.6"
        # Centering offset: (7.0 - 5.6) / 2 = 0.7"
        # Starting position: 0.5 (header) + 0.7 (centering) = 1.2"

        # CALCULATION: Fan wheel vertical centering with equal white space
        # Available space: 7.5 - 0.5 = 7.0"
        # Fan wheel + title height: 0.3 (title) + 5.8 (wheel) = 6.1"
        # White space total: 7.0 - 6.1 = 0.9"
        # Equal white space top/bottom: 0.9 / 2 = 0.45"
        # Fan wheel title top: 0.5 (header) + 0.45 (white space) = 0.95"
        # Fan wheel top: 0.95 + 0.3 (title height) = 1.25"

        # CALCULATION: Fan wheel horizontal centering
        # Community chart right edge: 0.5 + 5.5 = 6.0"
        # Page right edge: 13.333"
        # Center point: (6.0 + 13.333) / 2 = 9.667"
        # Fan wheel left position: 9.667 - (5.8 / 2) = 6.767"

        # Add elements with calculated positions
        self._add_header(slide)
        self._add_insight_text(slide)
        self._add_chart_title_left(slide)
        self._add_fan_wheel_title(slide)
        self._add_community_chart(slide, community_chart_path)
        self._add_fan_wheel(slide, fan_wheel_path)
        self._add_explanation_text(slide)

        return pres

    def _add_header(self, slide):
        """Add header bar"""
        # Gray header bar
        header_rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(13.333), Inches(0.5)
        )
        header_rect.fill.solid()
        header_rect.fill.fore_color.rgb = RGBColor(240, 240, 240)
        header_rect.line.color.rgb = RGBColor(200, 200, 200)
        header_rect.line.width = Pt(0.5)

        # Team name left
        team_text = slide.shapes.add_textbox(
            Inches(0.2), Inches(0.1),
            Inches(3), Inches(0.3)
        )
        team_text.text_frame.text = "Utah Jazz"
        p = team_text.text_frame.paragraphs[0]
        p.font.name = self.default_font
        p.font.size = Pt(14)
        p.font.bold = True

        # Title right
        title_text = slide.shapes.add_textbox(
            Inches(6.5), Inches(0.1),
            Inches(6.633), Inches(0.3)
        )
        title_text.text_frame.text = "Fan Behaviors: How Are Utah Jazz Fans Unique"
        p = title_text.text_frame.paragraphs[0]
        p.font.name = self.default_font
        p.font.size = Pt(14)
        p.alignment = PP_ALIGN.RIGHT

    def _add_insight_text(self, slide):
        """Add large insight text at TOP of left side - VERTICALLY CENTERED"""
        text_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.2),  # Adjusted for new centering
            Inches(5.5), Inches(1.0)  # Large width, reasonable height
        )
        text_box.text_frame.text = "Jazz fans are entertainment seekers who are cost conscious!"
        text_box.text_frame.word_wrap = True

        p = text_box.text_frame.paragraphs[0]
        p.font.name = self.default_font
        p.font.size = Pt(18)  # Large font for prominence
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT

    def _add_chart_title_left(self, slide):
        """Add chart title on LEFT side, below insight - VERTICALLY CENTERED"""
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.2),  # Adjusted for new centering
            Inches(5.5), Inches(0.3)
        )
        title_box.text_frame.text = "Top Ten Utah Jazz Fan Communities"
        p = title_box.text_frame.paragraphs[0]
        p.font.name = self.default_font
        p.font.size = Pt(14)
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT  # Left aligned like reference

    def _add_fan_wheel_title(self, slide):
        """Add fan wheel title - Equal white space top and bottom"""
        # Fan wheel center: 9.667", title should be centered on this
        title_center = 9.667
        title_width = 5.5
        title_left = title_center - (title_width / 2)  # 7.167"

        title_box = slide.shapes.add_textbox(
            Inches(title_left), Inches(0.95),  # Equal white space: 0.45" from header
            Inches(title_width), Inches(0.3)
        )
        title_box.text_frame.text = "Top Community Fan Purchases"
        p = title_box.text_frame.paragraphs[0]
        p.font.name = self.default_font
        p.font.size = Pt(14)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

    def _add_community_chart(self, slide, image_path: Path):
        """Add community chart on LEFT side - VERTICALLY CENTERED"""
        left = Inches(0.5)
        top = Inches(2.5)  # Adjusted for new centering
        width = Inches(5.5)

        slide.shapes.add_picture(str(image_path), left, top, width=width)

    def _add_fan_wheel(self, slide, image_path: Path):
        """Add fan wheel - Equal white space top and bottom"""
        width = Inches(5.8)  # 5.8" diameter as requested
        # Community chart right edge: 6.0", Page right edge: 13.333"
        # Center point: (6.0 + 13.333) / 2 = 9.667"
        # Fan wheel left: 9.667 - (5.8 / 2) = 6.767"
        left = Inches(6.767)
        top = Inches(1.25)  # Equal white space: title at 0.95 + 0.3 = 1.25

        slide.shapes.add_picture(str(image_path), left, top, width=width)

    def _add_explanation_text(self, slide):
        """Add explanation text - REDUCED HEIGHT to fit content only"""
        text_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(6.0),  # Adjusted for new centering
            Inches(5.5), Inches(0.8)  # REDUCED from 1.2 to 0.8 to fit text only
        )
        explanation = (
            "The top ten fan communities are ranked according to a composite index score "
            "of likelihood to purchase, likelihood to make more purchases per fan versus "
            "the local general population, and likelihood to spend more per fan."
        )
        text_box.text_frame.text = explanation
        text_box.text_frame.word_wrap = True

        p = text_box.text_frame.paragraphs[0]
        p.font.name = self.default_font
        p.font.size = Pt(10)
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.2


def fresh_test():
    """Run the fresh test"""
    print("🆕 Creating FRESH behaviors slide test with correct positioning...")

    tester = FreshBehaviorsSlideTest()
    presentation = tester.create_test_slide()

    # Save
    output_dir = Path('fresh_test_output')
    output_dir.mkdir(exist_ok=True)

    output_path = output_dir / 'fresh_behaviors_slide_test.pptx'
    presentation.save(str(output_path))

    print(f"✅ Fresh test slide created: {output_path}")
    print("📝 This should now match your reference layout!")

    return output_path


if __name__ == "__main__":
    # Run fresh test
    output_file = fresh_test()

    print(f"\n📋 FRESH TEST RESULTS:")
    print(f"1. Open: {output_file}")
    print(f"2. Layout should now match reference:")
    print(f"   ✓ Large insight text at top left")
    print(f"   ✓ Chart title below insight (left aligned)")
    print(f"   ✓ Community chart below title")
    print(f"   ✓ Fan wheel on right (5.8\" diameter)")
    print(f"   ✓ Explanation text at bottom left")
    print(f"3. If still not right, check positioning values")