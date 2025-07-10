# slide_generators/demographics_slide.py
"""
Generate complete demographics slide for PowerPoint presentations
Includes all charts arranged according to the reference layout
Updated with optimized gender pie chart spacing and wider bar charts
"""

from pathlib import Path
from typing import Dict, Optional, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import logging

from .base_slide import BaseSlide

logger = logging.getLogger(__name__)

# Default font
DEFAULT_FONT_FAMILY = "Red Hat Display"


class DemographicsSlide(BaseSlide):
    """Generate the complete demographics slide with all charts"""

    def __init__(self, presentation: Presentation = None):
        """
        Initialize demographics slide generator

        Args:
            presentation: Existing presentation to add slide to
        """
        super().__init__(presentation)

    def generate(self,
                 demographic_data: Dict[str, Any],
                 chart_dir: Path,
                 team_config: Dict[str, Any],
                 slide_index: Optional[int] = None) -> Presentation:
        """
        Generate the complete demographics slide

        Args:
            demographic_data: Processed demographic data
            chart_dir: Directory containing chart images
            team_config: Team configuration
            slide_index: Where to insert slide (None = append)

        Returns:
            Updated presentation object
        """
        # Extract team info
        team_name = team_config.get('team_name', 'Team')
        team_short = team_config.get('team_name_short', team_name.split()[-1])
        colors = team_config.get('colors', {})

        # Use the content layout (SIL white layout #12)
        slide = self.add_content_slide()
        logger.info("Added demographics slide using SIL white layout")

        # Add header
        self._add_header(slide, team_name)

        # Add chart section headers with black backgrounds
        self._add_chart_headers(slide)

        # Add demographic charts in 2x3 grid with optimized spacing
        self._add_charts(slide, chart_dir)

        # Add KEY/legend box at bottom-left
        self._add_legend_box(slide, team_name, team_short,
                             team_config.get('league', 'League'))

        logger.info("Demographics slide generation complete")
        return self.presentation

    def _add_header(self, slide, team_name: str):
        """Add team name and slide title header"""
        # Team name (top-left)
        team_text = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.2),
            Inches(3), Inches(0.5)
        )
        team_text.text_frame.text = team_name
        p = team_text.text_frame.paragraphs[0]
        p.font.name = DEFAULT_FONT_FAMILY
        p.font.size = Pt(18)
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT

        # Slide title (top-right)
        title_text = slide.shapes.add_textbox(
            Inches(7), Inches(0.2),
            Inches(6), Inches(0.5)
        )
        title_text.text_frame.text = "FAN DEMOGRAPHICS: HOW ARE UTAH JAZZ FANS UNIQUE"
        p = title_text.text_frame.paragraphs[0]
        p.font.name = DEFAULT_FONT_FAMILY
        p.font.size = Pt(16)
        p.font.bold = True
        p.alignment = PP_ALIGN.RIGHT

    def _add_chart_headers(self, slide):
        """Add black header bars above each chart section"""
        # Header configurations: (text, left, top, width)
        headers = [
            # Top row headers
            ("GENDER", 0.5, 0.95, 2.5),  # NARROWER: 2.5 instead of 3.5
            ("HOUSEHOLD INCOME", 3.2, 0.95, 5.5),  # WIDER: 5.5 instead of 5.0, starts earlier
            ("OCCUPATION CATEGORY", 9.0, 0.95, 4.2),  # WIDER: 4.2 instead of 3.8, starts earlier

            # Bottom row headers
            ("ETHNICITY", 0.5, 3.65, 2.5),  # NARROWER to match
            ("GENERATION", 3.2, 3.65, 5.5),  # WIDER to match
            ("CHILDREN IN HOUSEHOLD", 9.0, 3.65, 4.2)  # WIDER to match
        ]

        for text, left, top, width in headers:
            # Black background rectangle
            header_rect = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(left), Inches(top),
                Inches(width), Inches(0.25)
            )
            header_rect.fill.solid()
            header_rect.fill.fore_color.rgb = RGBColor(0, 0, 0)
            header_rect.line.fill.background()

            # White text
            text_box = slide.shapes.add_textbox(
                Inches(left + 0.1), Inches(top + 0.02),
                Inches(width - 0.2), Inches(0.2)
            )
            text_box.text_frame.text = text
            p = text_box.text_frame.paragraphs[0]
            p.font.name = DEFAULT_FONT_FAMILY
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER

    def _add_charts(self, slide, chart_dir: Path):
        """Add all demographic charts with optimized spacing"""
        chart_dir = Path(chart_dir)

        # UPDATED: Better space distribution
        # Gender gets minimal space, Income/Occupation get more space
        chart_positions = [
            # Top row - Redistributed widths
            ('gender_chart', 0.5, 1.2, 2.5, 2.2),  # NARROWER: 2.5 instead of 3.5
            ('income_chart', 3.2, 1.2, 5.5, 2.2),  # WIDER: 5.5 instead of 5.0, starts earlier
            ('occupation_chart', 9.0, 1.2, 4.2, 2.2),  # WIDER: 4.2 instead of 3.8, starts earlier

            # Bottom row - Same redistribution
            ('ethnicity_chart', 0.5, 3.9, 2.5, 2.2),  # NARROWER to match
            ('generation_chart', 3.2, 3.9, 5.5, 2.2),  # WIDER to match
            ('children_chart', 9.0, 3.9, 4.2, 2.2)  # WIDER to match
        ]

        for chart_name, left, top, width, height in chart_positions:
            # Try both regular and hires versions
            chart_path = chart_dir / f'{chart_name}_hires.png'
            if not chart_path.exists():
                chart_path = chart_dir / f'{chart_name}.png'

            if chart_path.exists():
                try:
                    slide.shapes.add_picture(
                        str(chart_path),
                        Inches(left), Inches(top),
                        width=Inches(width), height=Inches(height)
                    )
                    logger.info(f"Added {chart_name} to slide")
                except Exception as e:
                    logger.warning(f"Could not add {chart_name}: {e}")
            else:
                logger.warning(f"Chart not found: {chart_path}")
                # Add placeholder
                placeholder = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(left), Inches(top),
                    Inches(width), Inches(height)
                )
                placeholder.fill.solid()
                placeholder.fill.fore_color.rgb = RGBColor(245, 245, 245)
                placeholder.line.color.rgb = RGBColor(100, 100, 100)
                placeholder.line.width = Pt(1)

    def _add_legend_box(self, slide, team_name: str, team_short: str, league: str):
        """Add KEY legend box at bottom-left"""
        # Box position - adjusted for new layout
        left = Inches(0.8)
        top = Inches(6.2)
        width = Inches(5.5)  # Slightly wider to accommodate longer text
        height = Inches(0.8)

        # Create box with border
        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, top, width, height
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(255, 255, 255)
        box.line.color.rgb = RGBColor(0, 0, 0)
        box.line.width = Pt(1)

        # Add legend text
        text_box = slide.shapes.add_textbox(
            left + Inches(0.1), top + Inches(0.05),
            width - Inches(0.2), height - Inches(0.1)
        )

        # Format legend text lines
        legend_lines = [
            "KEY",
            f"- {team_name}",
            f"- {team_short} Gen Pop (state level, excluding {team_short} Fans)",
            f"- {league} Fans Total (excluding {team_short.lower()} fans)"
        ]

        text_frame = text_box.text_frame
        text_frame.clear()

        for i, line in enumerate(legend_lines):
            if i > 0:
                text_frame.add_paragraph()

            p = text_frame.paragraphs[i]
            p.text = line
            p.font.name = DEFAULT_FONT_FAMILY

            if i == 0:  # "KEY" header
                p.font.size = Pt(12)
                p.font.bold = True
            else:  # Legend items
                p.font.size = Pt(10)
                p.font.bold = False

            p.alignment = PP_ALIGN.LEFT