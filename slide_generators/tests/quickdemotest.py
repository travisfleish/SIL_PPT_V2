#!/usr/bin/env python3
"""
quick_test_slide.py
Minimal test script to generate the demographic overview slide
Can run independently without full data pipeline
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
from datetime import datetime

# Import your slide generator (adjust path as needed)
try:
    from slide_generators.demographic_overview_slide import DemographicOverviewSlide

    print("✅ Successfully imported DemographicOverviewSlide")
except ImportError as e:
    print(f"❌ Import error: {e}")
    print("🔧 Make sure the slide_generators module is in your Python path")
    print("🔧 Run this script from your project root directory")
    exit(1)


def quick_test():
    """Quick test of the demographic overview slide"""

    print("\n" + "=" * 50)
    print("🧪 QUICK TEST: Demographic Overview Slide")
    print("=" * 50)

    # Mock team configuration (replace with real config if available)
    team_configs = [
        {
            'team_name': 'Utah Jazz',
            'team_name_short': 'Jazz',
            'league': 'NBA',
            'colors': {
                'primary': '#002B5C',  # Navy
                'secondary': '#00471B',  # Green
                'accent': '#F9A01B'  # Gold
            }
        },
        {
            'team_name': 'Dallas Cowboys',
            'team_name_short': 'Cowboys',
            'league': 'NFL',
            'colors': {
                'primary': '#041E42',  # Navy
                'secondary': '#869397',  # Silver
                'accent': '#FFFFFF'  # White
            }
        }
    ]

    # Sample AI insights for testing
    sample_insights = {
        'Utah Jazz': ("Utah Jazz fans are significantly younger, with 79% being Millennials/Gen X "
                      "compared to 45% in the Utah general population and 76% of NBA average fans, "
                      "have higher household income with 44% earning $100K+ compared to 29% in the "
                      "Utah general population and 38% of NBA average fans, are predominantly male "
                      "at 59% versus 49% in the Utah general population and 52% of NBA average fans, "
                      "and are largely working professionals at 50% compared to 34% in the Utah "
                      "general population and 34% for NBA average fans."),

        'Dallas Cowboys': ("Dallas Cowboys fans are younger and higher-earning professionals who "
                           "are more likely to be parents compared to the Texas general population, "
                           "with 68% being Millennials/Gen X versus 58% in the general population "
                           "and 61% earning $75K+ compared to 48% in the general population, making "
                           "them an attractive demographic for premium lifestyle and family-oriented "
                           "sponsorship opportunities.")
    }

    # Create output directory
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    output_dir = Path(f'test_slides_{timestamp}')
    output_dir.mkdir(exist_ok=True)

    print(f"📁 Output directory: {output_dir}")

    # Test each team configuration
    for i, team_config in enumerate(team_configs, 1):
        team_name = team_config['team_name']
        print(f"\n{i}. Testing {team_name}")
        print(f"   League: {team_config['league']}")

        try:
            # Create slide generator
            slide_generator = DemographicOverviewSlide()

            # Get sample insight for this team
            ai_insights = sample_insights.get(team_name,
                                              f"{team_name} fans have unique demographic characteristics compared to the general population.")

            print(f"   Insight: {ai_insights[:80]}...")

            # Generate the slide
            slide_generator.generate(
                team_config=team_config,
                ai_insights=ai_insights
            )

            # Save the presentation
            output_file = output_dir / f"{team_config['team_name_short'].lower()}_demographic_overview.pptx"
            slide_generator.presentation.save(str(output_file))

            print(f"   ✅ Saved: {output_file}")

            # Basic validation
            presentation = slide_generator.presentation
            if len(presentation.slides) > 0:
                slide = presentation.slides[0]
                shape_count = len(slide.shapes)
                print(f"   📊 Slide has {shape_count} shapes")

                # Count text shapes
                text_shapes = [s for s in slide.shapes if hasattr(s, 'text_frame') and s.has_text_frame]
                print(f"   📝 Found {len(text_shapes)} text elements")

        except Exception as e:
            print(f"   ❌ Error: {e}")
            import traceback
            traceback.print_exc()

    print(f"\n" + "=" * 50)
    print("✅ Test completed!")
    print(f"📁 Check your files in: {output_dir}")
    print("🔍 Open the .pptx files in PowerPoint to review")
    print("=" * 50)

    return output_dir


def manual_slide_inspection(output_dir: Path):
    """Provide manual inspection instructions"""
    print(f"\n📋 MANUAL INSPECTION CHECKLIST:")
    print(f"1. Open the files in: {output_dir}")
    print(f"2. Check that each slide has:")
    print(f"   ✓ Blue background (SIL layout #11)")
    print(f"   ✓ Gray header bar at the top")
    print(f"   ✓ Team name in header (left side)")
    print(f"   ✓ 'Fan Demographics: How Are [Team] Fans Unique' in header (right side)")
    print(f"   ✓ Large team name + 'Fan Demographic Overview' title")
    print(f"   ✓ AI insights text on the left side")
    print(f"   ✓ Circular image placeholder on the right side")
    print(f"3. Verify text is white and readable on blue background")
    print(f"4. Check that layout matches your reference slide design")


if __name__ == "__main__":
    # Run the quick test
    output_dir = quick_test()

    # Provide inspection guidance
    manual_slide_inspection(output_dir)

    print(f"\n🎯 NEXT STEPS:")
    print(f"1. If slides look good, integrate into your main PowerPoint builder")
    print(f"2. If issues found, adjust the slide generator code")
    print(f"3. Test with real team configurations and AI insights")
    print(f"4. Add to your presentation build pipeline")