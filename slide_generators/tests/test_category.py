#!/usr/bin/env python3
"""
Test script to generate category slides for PowerPoint
Tests the CategorySlide generator with mock data
"""

import pandas as pd
import numpy as np
from pathlib import Path
from pptx import Presentation
from datetime import datetime
import logging
import sys

# Add project root to path (adjust if your structure is different)
project_root = Path(__file__).parent
if project_root not in sys.path:
    sys.path.insert(0, str(project_root))

# Import the required modules
from slide_generators.category_slide import CategorySlide
from data_processors.category_analyzer import CategoryMetrics

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)


def create_mock_category_data():
    """Create realistic mock data for testing category slides"""

    # Category metrics
    category_metrics = CategoryMetrics(
        percent_fans=0.68,  # 68% of fans spend
        percent_likely=127,  # 127% more likely
        percent_purchases=85,  # 85% more purchases
        composite_index=487.5,
        total_spend=2500000,
        spc=175.50,
        audience_count=25000,
        comparison_population="Local Gen Pop (Excl. Jazz)"
    )

    # Subcategory stats - formatted as they would be from the analyzer
    subcategory_data = [
        {
            'Subcategory': 'Fast Food',
            'Percent of Fans Who Spend': '72%',
            'Likelihood to spend (vs. Local Gen Pop)': '135% More',  # UPDATED column name
            'Purchases Per Fan (vs. Gen Pop)': '98% More'
        },
        {
            'Subcategory': 'Fast Casual Dining',
            'Percent of Fans Who Spend': '65%',
            'Likelihood to spend (vs. Local Gen Pop)': '118% More',  # UPDATED column name
            'Purchases Per Fan (vs. Gen Pop)': '76% More'
        },
        {
            'Subcategory': 'Full Service Restaurants',
            'Percent of Fans Who Spend': '58%',
            'Likelihood to spend (vs. Local Gen Pop)': '95% More',  # UPDATED column name
            'Purchases Per Fan (vs. Gen Pop)': '62% More'
        },
        {
            'Subcategory': 'Coffee & Tea',
            'Percent of Fans Who Spend': '52%',
            'Likelihood to spend (vs. Local Gen Pop)': '88% More',  # UPDATED column name
            'Purchases Per Fan (vs. Gen Pop)': '54% More'
        }
    ]
    subcategory_stats = pd.DataFrame(subcategory_data)

    # Category insights
    insights = [
        "Jazz Fans are 127% MORE likely to spend on Restaurants than the Local Gen Pop (Excl. Jazz)",
        "Jazz Fans make an average of 85% more purchases per fan on Restaurants than the Local Gen Pop (Excl. Jazz)",
        "Jazz Fans are more than 4.3X more likely to spend on Fast Food vs. the Local Gen Pop (Excl. Jazz)",
        "Jazz fans spend an average of $892 per fan per year on Fast Food",
        "Jazz fans are 1% more likely to spend on QSR & Fast Casual when compared to the NBA average"
    ]

    # Merchant stats - as returned by the analyzer
    merchant_data = [
        {
            'Rank': 1,
            'Brand': 'Chick-fil-A',
            'Percent of Fans Who Spend': '45.2%',
            'Likelihood to spend (vs. Local Gen Pop)': '156% More',  # UPDATED column name
            'Purchases Per Fan (vs. Gen Pop)': '112% More'
        },
        {
            'Rank': 2,
            'Brand': 'In-N-Out Burger',
            'Percent of Fans Who Spend': '42.8%',
            'Likelihood to spend (vs. Local Gen Pop)': '148% More',  # UPDATED column name
            'Purchases Per Fan (vs. Gen Pop)': '98% More'
        },
        {
            'Rank': 3,
            'Brand': 'Chipotle',
            'Percent of Fans Who Spend': '38.5%',
            'Likelihood to spend (vs. Local Gen Pop)': '135% More',  # UPDATED column name
            'Purchases Per Fan (vs. Gen Pop)': '87% More'
        },
        {
            'Rank': 4,
            'Brand': 'Starbucks',
            'Percent of Fans Who Spend': '36.2%',
            'Likelihood to spend (vs. Local Gen Pop)': '128% More',  # UPDATED column name
            'Purchases Per Fan (vs. Gen Pop)': '92% More'
        },
        {
            'Rank': 5,
            'Brand': "McDonald's",
            'Percent of Fans Who Spend': '34.7%',
            'Likelihood to spend (vs. Local Gen Pop)': '95% More',  # UPDATED column name
            'Purchases Per Fan (vs. Gen Pop)': '68% More'
        }
    ]
    merchant_df = pd.DataFrame(merchant_data)
    top_merchants = [row['Brand'] for row in merchant_data]

    # Merchant insights
    merchant_insights = [
        "45.2% of Utah Jazz fans spent at Chick-fil-A",
        "Utah Jazz fans make an average of 52 purchases per year at In-N-Out Burger—more than any other top Restaurant brand",
        "Utah Jazz fans spent an average of $1,245 per fan on Chick-fil-A per year",
        "Utah Jazz fans are 245% more likely to spend on Chipotle than NBA Fans"
    ]

    # Sponsorship recommendation
    recommendation = {
        'merchant': 'Chick-fil-A',
        'composite_index': 542.8,
        'explanation': 'The Jazz should target Chick-fil-A for a sponsorship based on having the highest composite index of 543',
        'sub_explanation': 'The composite index indicates a brand with significant likelihood for more fans to be spending more frequently, and at a higher spend per fan vs. other brands',
        'full_recommendation': {
            'main': 'The Utah Jazz should target Chick-fil-A for a sponsorship based on having the highest composite index of 543',
            'sub_bullet': 'The composite index indicates a brand with significant likelihood for more fans to be spending more frequently, and at a higher spend per fan vs. other brands'
        }
    }

    # Complete analysis results as would be returned by CategoryAnalyzer
    analysis_results = {
        'category_key': 'restaurants',
        'display_name': 'Restaurants',
        'slide_title': 'Dining & QSR Analysis',
        'category_metrics': category_metrics,
        'subcategory_stats': subcategory_stats,
        'insights': insights,
        'merchant_stats': (merchant_df, top_merchants),
        'merchant_insights': merchant_insights,
        'recommendation': recommendation,
        'validation_report': None
    }

    return analysis_results


def create_mock_team_config():
    """Create mock team configuration"""
    return {
        'team_name': 'Utah Jazz',
        'team_name_short': 'Jazz',
        'league': 'NBA',
        'primary_color': '002B5C',  # Jazz navy blue
        'secondary_color': 'F9A01B',  # Jazz gold
        'logo_path': None  # Would normally point to team logo
    }


def load_sil_template():
    """Load the SIL template presentation"""
    # Try different possible paths for the template
    template_paths = [
        Path("templates/sil_combined_template.pptx"),
        Path("../templates/sil_combined_template.pptx"),
        Path("../../templates/sil_combined_template.pptx"),
        Path.cwd() / "templates" / "sil_combined_template.pptx",
    ]

    for template_path in template_paths:
        if template_path.exists():
            logger.info(f"Found SIL template at: {template_path}")
            return Presentation(str(template_path))

    logger.warning("SIL template not found, creating blank presentation")
    logger.warning("Looked in: " + ", ".join(str(p) for p in template_paths))
    return None


def test_category_slide_generation():
    """Test generating category analysis slides"""
    logger.info("Starting category slide generation test...")

    # Create mock data
    analysis_results = create_mock_category_data()
    team_config = create_mock_team_config()

    # Load SIL template
    template_presentation = load_sil_template()

    # Initialize slide generator with template
    generator = CategorySlide(presentation=template_presentation)

    # Log available layouts
    if template_presentation:
        logger.info(f"Template has {len(generator.presentation.slide_layouts)} layouts")
        for i, layout in enumerate(generator.presentation.slide_layouts):
            logger.info(f"  Layout {i}: {layout.name}")

    # Generate the category analysis slide (first slide)
    logger.info("Generating category analysis slide...")
    presentation = generator.generate(analysis_results, team_config)

    # Generate the brand analysis slide (second slide)
    logger.info("Generating brand analysis slide...")
    presentation = generator.generate_brand_slide(analysis_results, team_config)

    # Save the presentation
    output_path = Path("test_category_slides.pptx")
    presentation.save(output_path)
    logger.info(f"✅ Test presentation saved to: {output_path.absolute()}")

    # Print summary
    logger.info("\nSlide Generation Summary:")
    logger.info(f"- Category: {analysis_results['display_name']}")
    logger.info(f"- Team: {team_config['team_name']}")
    logger.info(f"- Total slides generated: {len(presentation.slides)}")
    logger.info(f"- Category metrics: {analysis_results['category_metrics'].percent_fans * 100:.0f}% of fans spend")
    logger.info(f"- Top merchant: {analysis_results['merchant_stats'][0].iloc[0]['Brand']}")
    logger.info(f"- Recommendation: {analysis_results['recommendation']['merchant']}")

    return presentation


def test_multiple_categories():
    """Test generating slides for multiple categories"""
    logger.info("\nTesting multiple category generation...")

    # Categories to test
    categories = [
        {
            'display_name': 'Restaurants',
            'slide_title': 'Dining & QSR Analysis',
            'top_brands': ['Chick-fil-A', 'In-N-Out Burger', 'Chipotle', 'Starbucks', "McDonald's"]
        },
        {
            'display_name': 'Athleisure',
            'slide_title': 'Athleisure & Sporting Goods',
            'top_brands': ['Nike', 'Adidas', 'Under Armour', 'Lululemon', 'Dick\'s Sporting Goods']
        },
        {
            'display_name': 'Finance',
            'slide_title': 'Financial Services',
            'top_brands': ['Chase', 'Bank of America', 'Wells Fargo', 'American Express', 'Capital One']
        }
    ]

    team_config = create_mock_team_config()

    # Load SIL template
    template_presentation = load_sil_template()
    presentation = template_presentation

    for i, cat_info in enumerate(categories):
        logger.info(f"\nGenerating slides for {cat_info['display_name']}...")

        # Create mock data for this category
        analysis_results = create_mock_category_data()
        analysis_results['display_name'] = cat_info['display_name']
        analysis_results['slide_title'] = cat_info['slide_title']

        # Update merchant data with appropriate brands
        merchant_data = []
        for j, brand in enumerate(cat_info['top_brands']):
            merchant_data.append({
                'Rank': j + 1,
                'Brand': brand,
                'Percent of Fans Who Spend': f"{45 - j * 2:.1f}%",
                'Likelihood to spend (vs. Local Gen Pop)': f"{150 - j * 10}% More",  # UPDATED column name
                'Purchases Per Fan (vs. Gen Pop)': f"{110 - j * 8}% More"
            })

        merchant_df = pd.DataFrame(merchant_data)
        analysis_results['merchant_stats'] = (merchant_df, cat_info['top_brands'])
        analysis_results['recommendation']['merchant'] = cat_info['top_brands'][0]

        # Generate slides
        generator = CategorySlide(presentation)
        presentation = generator.generate(analysis_results, team_config)
        presentation = generator.generate_brand_slide(analysis_results, team_config)

    # Save the multi-category presentation
    output_path = Path("test_multiple_categories.pptx")
    presentation.save(output_path)
    logger.info(f"\n✅ Multi-category presentation saved to: {output_path.absolute()}")
    logger.info(f"Total slides: {len(presentation.slides)}")

    return presentation


def test_edge_cases():
    """Test edge cases and error handling"""
    logger.info("\nTesting edge cases...")

    # Load SIL template
    template_presentation = load_sil_template()

    # Test with empty subcategory data
    analysis_results = create_mock_category_data()
    analysis_results['subcategory_stats'] = pd.DataFrame()  # Empty dataframe

    team_config = create_mock_team_config()
    generator = CategorySlide(presentation=template_presentation)

    try:
        presentation = generator.generate(analysis_results, team_config)
        logger.info("✅ Handled empty subcategory data successfully")
    except Exception as e:
        logger.error(f"❌ Failed with empty subcategory data: {e}")

    # Test with very long brand names
    analysis_results = create_mock_category_data()
    merchant_df, _ = analysis_results['merchant_stats']
    merchant_df.loc[0, 'Brand'] = "Very Long Restaurant Name That Might Cause Layout Issues"

    try:
        presentation = generator.generate(analysis_results, team_config)
        presentation = generator.generate_brand_slide(analysis_results, team_config)
        logger.info("✅ Handled long brand names successfully")
    except Exception as e:
        logger.error(f"❌ Failed with long brand names: {e}")

    # Save edge case test
    output_path = Path("test_edge_cases.pptx")
    presentation.save(output_path)
    logger.info(f"✅ Edge case presentation saved to: {output_path.absolute()}")


if __name__ == "__main__":
    # Run the tests
    try:
        # Test 1: Basic category slide generation
        test_category_slide_generation()

        # Test 2: Multiple categories
        test_multiple_categories()

        # Test 3: Edge cases
        test_edge_cases()

        logger.info("\n🎉 All tests completed successfully!")
        logger.info("\nGenerated files:")
        logger.info("- test_category_slides.pptx (2 slides for Restaurants)")
        logger.info("- test_multiple_categories.pptx (6 slides for 3 categories)")
        logger.info("- test_edge_cases.pptx (edge case testing)")

    except Exception as e:
        logger.error(f"\n❌ Test failed with error: {e}")
        import traceback

        traceback.print_exc()