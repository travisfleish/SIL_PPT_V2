#!/usr/bin/env python3
"""
fixed_template_test.py
Test the demographic overview slide with proper template loading
This mimics how your pptx_builder loads the template first
"""

import sys
import os
from pathlib import Path
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches

# Add project root to path - get the actual project root
current_file = Path(__file__).resolve()
# If running from slide_generators/tests/, go up 2 levels to get to project root
if current_file.parent.name == 'tests' and current_file.parent.parent.name == 'slide_generators':
    project_root = current_file.parent.parent.parent  # Go up 3 levels: tests -> slide_generators -> project_root
else:
    project_root = current_file.parent.parent  # Default: go up 2 levels

sys.path.insert(0, str(project_root))

print(f"Running from: {current_file}")
print(f"Project root: {project_root}")
print(f"Python path includes: {project_root}")


def test_with_template():
    """Test the slide generator with proper template loading"""

    print("🔧 Testing demographic overview slide with TEMPLATE LOADING...")
    print("=" * 60)

    try:
        # 1. Load the template first (like your pptx_builder does)
        TEMPLATE_PATH = project_root / 'templates' / 'sil_combined_template.pptx'

        if not TEMPLATE_PATH.exists():
            print(f"❌ Template not found at: {TEMPLATE_PATH}")

            # Try alternative paths
            alternative_paths = [
                Path.cwd() / 'templates' / 'sil_combined_template.pptx',
                Path(__file__).resolve().parent.parent.parent / 'templates' / 'sil_combined_template.pptx'
            ]

            for alt_path in alternative_paths:
                print(f"🔍 Trying: {alt_path}")
                if alt_path.exists():
                    TEMPLATE_PATH = alt_path
                    print(f"✅ Found template at: {TEMPLATE_PATH}")
                    break
            else:
                print("❌ Could not find template file")
                print("🔧 Make sure the template file exists at one of these locations:")
                print(f"   - {project_root / 'templates' / 'sil_combined_template.pptx'}")
                for path in alternative_paths:
                    print(f"   - {path}")
                return False

        print(f"✅ Found template: {TEMPLATE_PATH}")

        # Load the template presentation
        presentation = Presentation(str(TEMPLATE_PATH))
        print(f"✅ Loaded template with {len(presentation.slide_layouts)} layouts")

        # Log available layouts (for debugging)
        for i, layout in enumerate(presentation.slide_layouts):
            print(f"   Layout {i}: {layout.name}")

        # Set 16:9 dimensions
        presentation.slide_width = Inches(13.333)
        presentation.slide_height = Inches(7.5)

        # 2. Import the slide generator
        from slide_generators.demographic_overview_slide import DemographicOverviewSlide
        from utils.team_config_manager import TeamConfigManager

        # 3. Get team config
        config_manager = TeamConfigManager()
        team_config = config_manager.get_team_config('utah_jazz')

        # 4. Create sample AI insight
        ai_insight = ("Utah Jazz fans are significantly younger, with 79% being Millennials/Gen X "
                      "compared to 45% in the Utah general population, have higher household income "
                      "with 44% earning $100K+ compared to 29% in the Utah general population, "
                      "and are largely working professionals at 50% compared to 34% in the Utah "
                      "general population.")

        print(f"✅ Team: {team_config['team_name']}")
        print(f"✅ AI insight: {ai_insight[:80]}...")

        # 5. Create the slide generator with the template-loaded presentation
        slide_generator = DemographicOverviewSlide(presentation)  # Pass the template!

        # 6. Generate the slide
        slide_generator.generate(
            team_config=team_config,
            ai_insights=ai_insight
        )

        # 7. Save the result
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        output_file = f"demographic_overview_WITH_TEMPLATE_{timestamp}.pptx"
        presentation.save(output_file)

        print(f"✅ SUCCESS! Slide saved as: {output_file}")
        print(f"📊 Presentation now has {len(presentation.slides)} slides")

        # 8. Validate the slide
        if len(presentation.slides) > 0:
            last_slide = presentation.slides[-1]  # Our new slide should be the last one
            shape_count = len(last_slide.shapes)
            print(f"✅ New slide has {shape_count} shapes")

            # Check for text content
            text_shapes = []
            for shape in last_slide.shapes:
                if hasattr(shape, 'text_frame') and shape.has_text_frame:
                    if shape.text_frame.text.strip():
                        text_shapes.append(shape.text_frame.text.strip()[:50])

            print(f"✅ Found {len(text_shapes)} text elements:")
            for i, text in enumerate(text_shapes):
                print(f"   {i + 1}. {text}...")

        print("\n🎉 Template-based test completed successfully!")
        print("📋 The slide should now have:")
        print("   ✓ Blue background from layout #11")
        print("   ✓ Header bar with team name")
        print("   ✓ Main title and AI insights")
        print("   ✓ Image placeholder")

        return True

    except Exception as e:
        print(f"❌ Error: {e}")
        import traceback
        traceback.print_exc()
        return False


def debug_template_layouts():
    """Debug function to examine template layouts"""
    print("\n🔍 DEBUGGING TEMPLATE LAYOUTS...")

    try:
        # Use the corrected project root path
        TEMPLATE_PATH = project_root / 'templates' / 'sil_combined_template.pptx'
        print(f"Looking for template at: {TEMPLATE_PATH}")

        if not TEMPLATE_PATH.exists():
            print(f"❌ Template not found at: {TEMPLATE_PATH}")
            # Try to find it in other common locations
            possible_paths = [
                Path.cwd() / 'templates' / 'sil_combined_template.pptx',
                Path(__file__).parent.parent.parent / 'templates' / 'sil_combined_template.pptx',
            ]
            for path in possible_paths:
                if path.exists():
                    print(f"✅ Found template at: {path}")
                    TEMPLATE_PATH = path
                    break
            else:
                print("❌ Could not find template in any expected location")
                return

        presentation = Presentation(str(TEMPLATE_PATH))

        print(f"Template has {len(presentation.slide_layouts)} layouts:")
        for i, layout in enumerate(presentation.slide_layouts):
            layout_name = getattr(layout, 'name', 'Unknown')
            print(f"   {i:2d}: {layout_name}")

            # Try to identify which ones might be blue/title layouts
            if i == 11:
                print(f"       ^ This should be the SIL blue layout")

    except Exception as e:
        print(f"Error debugging layouts: {e}")


if __name__ == "__main__":
    # Debug layouts first
    debug_template_layouts()

    # Run the corrected test
    success = test_with_template()

    if success:
        print("\n✅ SOLUTION CONFIRMED!")
        print("🔑 Key insight: The slide generator needs the template-loaded presentation")
        print("📝 Integration steps:")
        print("   1. Your pptx_builder already loads the template correctly")
        print("   2. Pass that presentation to DemographicOverviewSlide(presentation)")
        print("   3. The slide will have access to layout #11 (SIL blue)")
    else:
        print("\n❌ Still having issues")
        print("🔧 Check that:")
        print("   1. Template file exists in templates/sil_combined_template.pptx")
        print("   2. Layout #11 exists in the template")
        print("   3. All required modules are importable")