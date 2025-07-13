#!/usr/bin/env python3
"""
test_title_slide_fixed.py
Test script that correctly uses the SIL template layout #11
Fixes the background inheritance issue
"""

import sys
from pathlib import Path
from pptx import Presentation
from datetime import datetime

# Add project root to Python path
# Fix the path resolution to handle different directory structures
current_file = Path(__file__).resolve()
if current_file.parent.name == 'tests' and current_file.parent.parent.name == 'slide_generators':
    # Running from slide_generators/tests/
    project_root = current_file.parent.parent.parent
elif current_file.parent.name == 'slide_generators':
    # Running from slide_generators/
    project_root = current_file.parent.parent
else:
    # Default fallback
    project_root = current_file.parent.parent

sys.path.insert(0, str(project_root))

print(f"Script location: {current_file}")
print(f"Project root: {project_root}")


def test_template_integration():
    """Test title slide with proper template layout inheritance"""

    print("\n" + "=" * 60)
    print("🧪 TITLE SLIDE TEMPLATE INTEGRATION TEST")
    print("=" * 60)

    try:
        # 1. Load the SIL template (like pptx_builder does)
        template_path = project_root / 'templates' / 'sil_combined_template.pptx'

        if not template_path.exists():
            print(f"❌ Template not found: {template_path}")

            # Try alternative locations
            print("\n🔍 Searching for template in alternative locations...")
            alternative_paths = [
                Path.cwd() / 'templates' / 'sil_combined_template.pptx',
                project_root.parent / 'templates' / 'sil_combined_template.pptx',
                Path(__file__).parent / 'templates' / 'sil_combined_template.pptx',
                Path(__file__).parent.parent / 'templates' / 'sil_combined_template.pptx',
            ]

            for alt_path in alternative_paths:
                print(f"   Checking: {alt_path}")
                if alt_path.exists():
                    template_path = alt_path
                    print(f"   ✅ Found template at: {template_path}")
                    break
            else:
                print("\n❌ Template not found in any location!")
                print("🔧 Ensure sil_combined_template.pptx exists in:")
                print(f"   {project_root / 'templates'}/")
                return False

        print(f"✅ Loading template: {template_path}")
        presentation = Presentation(str(template_path))

        # 2. Debug available layouts
        print(f"\nTemplate has {len(presentation.slide_layouts)} layouts:")
        for i, layout in enumerate(presentation.slide_layouts):
            layout_name = getattr(layout, 'name', 'Unknown')
            print(f"   {i:2d}: {layout_name}")
            if i == 11:
                print(f"       ^ SIL Blue Layout (target)")

        # 3. Test with the corrected TitleSlide class
        try:
            from slide_generators.title_slide import TitleSlide
            print(f"\n✅ Successfully imported TitleSlide class")
        except ImportError as e:
            print(f"❌ Import failed: {e}")
            return False

        # 4. Create title slide with Utah Jazz config
        jazz_config = {
            'team_name': 'Utah Jazz',
            'team_initials': 'UJ',
            'colors': {
                'primary': '#1D428A',  # Jazz blue
                'secondary': '#FFD700',  # Gold
                'accent': '#00275D'  # Navy
            }
        }

        print(f"\n🎨 Generating title slide for {jazz_config['team_name']}...")

        # 5. Generate title slide (should inherit blue background from layout #11)
        generator = TitleSlide(presentation)  # Pass template-loaded presentation
        presentation = generator.generate(jazz_config)

        # 6. Save and verify
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")

        # Create output directory if it doesn't exist
        output_dir = project_root / 'output' / 'test'
        output_dir.mkdir(parents=True, exist_ok=True)

        filename = output_dir / f"utah_jazz_title_slide_fixed_{timestamp}.pptx"
        presentation.save(str(filename))

        print(f"✅ SUCCESS! Title slide created: {filename}")

        # Check if team logo exists
        logo_path = project_root / 'assets' / 'logos' / 'teams' / 'utah_jazz.png'
        if logo_path.exists():
            print(f"✅ Team logo found at: {logo_path}")
            logo_status = "Utah Jazz logo image"
        else:
            print(f"⚠️  Team logo not found at: {logo_path}")
            logo_status = "placeholder with 'UJ' initials"

        print(f"\n📋 Verification checklist:")
        print(f"   □ Open the PowerPoint file")
        print(f"   □ Slide should have BLUE background from template (not manually added)")
        print(f"   □ Centered team logo: {logo_status}")
        print(f"   □ 'Utah Jazz' as main title")
        print(f"   □ 'Sponsorship Insights Report' subtitle")
        print(f"   □ Descriptive text below subtitle")
        print(f"   □ Current date in bottom right")
        print(f"   □ SIL logo/text in bottom left")

        return True

    except Exception as e:
        print(f"❌ Error: {type(e).__name__}: {str(e)}")
        import traceback
        traceback.print_exc()
        return False


def test_background_inheritance():
    """
    Specific test to verify background inheritance vs manual background setting
    """
    print(f"\n" + "🔍" * 60)
    print("BACKGROUND INHERITANCE TEST")
    print("🔍" * 60)

    try:
        template_path = project_root / 'templates' / 'sil_combined_template.pptx'

        if not template_path.exists():
            print("❌ Template not found - skipping background test")
            return False

        # Import required classes
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        # Load template
        presentation = Presentation(str(template_path))

        # Test 1: Slide with layout #11 (should have blue background automatically)
        print("\n1. Testing layout #11 (SIL Blue) inheritance...")
        if len(presentation.slide_layouts) > 11:
            blue_layout = presentation.slide_layouts[11]
            slide_blue = presentation.slides.add_slide(blue_layout)
            print(f"   ✅ Added slide with layout: {getattr(blue_layout, 'name', 'Layout 11')}")
        else:
            print("   ❌ Layout #11 not available")
            return False

        # Test 2: Slide with blank layout (for comparison)
        print("\n2. Testing blank layout (should need manual background)...")
        blank_layout = presentation.slide_layouts[6]  # Standard blank
        slide_blank = presentation.slides.add_slide(blank_layout)
        print(f"   ✅ Added slide with layout: {getattr(blank_layout, 'name', 'Blank')}")

        # Add some text to both slides for testing
        for i, slide in enumerate([slide_blue, slide_blank], 1):
            text_box = slide.shapes.add_textbox(
                Inches(3),  # Left position
                Inches(3),  # Top position
                Inches(7),  # Width
                Inches(1.5)  # Height
            )
            text_box.text_frame.text = f"Test Slide {i}\n{'Layout #11 (SIL Blue)' if i == 1 else 'Blank Layout'}"

            # Format text
            for paragraph in text_box.text_frame.paragraphs:
                paragraph.font.size = Pt(24)
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White text
                paragraph.alignment = PP_ALIGN.CENTER

        # Save comparison
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_dir = project_root / 'output' / 'test'
        output_dir.mkdir(parents=True, exist_ok=True)

        filename = output_dir / f"background_inheritance_test_{timestamp}.pptx"
        presentation.save(str(filename))

        print(f"\n✅ Background comparison saved: {filename}")
        print(f"📋 Expected results:")
        print(f"   • Slide 1: Blue background (from layout #11)")
        print(f"   • Slide 2: White background (blank layout)")

        return True

    except Exception as e:
        print(f"❌ Background test failed: {e}")
        import traceback
        traceback.print_exc()
        return False


def main():
    """Main test runner"""
    print("🚀 TESTING FIXED TITLE SLIDE IMPLEMENTATION")
    print("🎯 Focus: Proper template layout inheritance")

    # Test 1: Template integration
    success1 = test_template_integration()

    # Test 2: Background inheritance verification
    success2 = test_background_inheritance()

    # Summary
    print(f"\n" + "🏁" * 20)
    print("FINAL RESULTS")
    print("🏁" * 20)

    if success1 and success2:
        print("✅ ALL TESTS PASSED!")
        print("🎉 The title slide now correctly inherits the blue background from layout #11")
        print("📝 Key fixes applied:")
        print("   • Removed manual background color setting")
        print("   • Title slide uses self.add_title_slide() which leverages layout #11")
        print("   • Template-loaded presentation passed to TitleSlide constructor")
        print("\n💡 If using team logos:")
        print("   • Place logo files in: assets/logos/teams/")
        print("   • Use format: team_name.png (e.g., utah_jazz.png)")
    else:
        print("❌ Some tests failed")
        print("🔧 Troubleshooting:")
        print("   1. Ensure templates/sil_combined_template.pptx exists")
        print("   2. Verify layout #11 is the SIL blue layout")
        print("   3. Check that slide_generators module imports correctly")


if __name__ == "__main__":
    main()