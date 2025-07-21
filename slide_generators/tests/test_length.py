"""
Test script to validate dynamic spacing in category insights
Creates test slides with varying insight lengths to verify NBA section positioning
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import sys
import os

# Add the parent directory to the path so we can import your modules
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from slide_generators.category_slide import CategorySlide

def create_test_presentation():
    """Create test presentation with various insight lengths"""

    # Create a new presentation
    prs = Presentation()

    # Initialize the CategorySlide generator
    slide_generator = CategorySlide(prs)

    # Test team configuration
    team_config = {
        'team_name': 'Carolina Panthers',
        'team_name_short': 'Panthers',
        'league': 'NFL'
    }

    # Test cases with different insight lengths - NOW INCLUDING NBA INSIGHTS
    test_cases = [
        {
            "name": "Short Insights Test",
            "category_name": "Restaurants",
            "insights": [
                "Panthers fans spend 32% more at QSR",
                "Average of 5 visits per month",
                "Spend $125 per fan annually",
                "Panthers fans are 45% more likely to visit fast food compared to NFL average",
                "Carolina Panthers fans prefer chain restaurants 73% more than compared to NBA fans nationwide",
                "Panthers fans spend 52% more on dining out when compared to NBA average"
            ],
            "is_emerging": False
        },
        {
            "name": "Medium Insights Test",
            "category_name": "Finance",
            "insights": [
                "56% of Panthers fans use mobile payment apps, significantly higher than the local gen pop average of 32%",
                "Panthers fans make an average of 25 transactions per month using digital payment platforms",
                "The average Panthers fan spends $3,767 annually through mobile payment apps, demonstrating high digital finance engagement",
                "Digital wallet adoption reaches 82% among season ticket holders",
                "Panthers fans are 98% more likely to use buy-now-pay-later services compared to NBA average",
                "Mobile payment usage is 145% higher among Panthers fans when compared to NBA fans overall"
            ],
            "is_emerging": False
        },
        {
            "name": "Long Insights Test",
            "category_name": "Travel",
            "insights": [
                "Carolina Panthers fans demonstrate exceptional travel spending patterns with 73% of fans booking travel accommodations at least twice per year, significantly exceeding the local gen pop average of 41% who travel with similar frequency",
                "The most active Panthers fans make an average of 8.5 travel-related purchases annually including flights, hotels, rental cars, and travel experiences, showing much higher engagement than typical sports fans",
                "Panthers fans spend an average of $4,892 per year on travel-related expenses, which is 62% higher than the local gen pop average annual travel spend of $3,021",
                "Premium travel services see exceptional adoption with 34% of Panthers fans holding airline status",
                "When compared to NBA fans nationwide, Panthers fans are 125% more likely to book premium travel accommodations and 87% more likely to purchase travel insurance",
                "International travel is 92% more common among Panthers fans compared to NBA average"
            ],
            "is_emerging": False
        },
        {
            "name": "Emerging Category Test",
            "category_name": "Subscription Services",
            "insights": [
                "Panthers fans show growing adoption of subscription services with 42% actively using 3 or more streaming platforms",
                "Monthly subscription spending averages $127 per fan across all service types including entertainment, fitness, and meal delivery",
                "The emerging nature of this category presents significant sponsorship opportunities as fan engagement continues to grow rapidly",
                "Digital subscription adoption has grown 47% year-over-year among the fanbase",
                "Panthers fans are 156% more likely to try new subscription services compared to NBA average, indicating high potential for brand partnerships",
                "Subscription stacking is 89% more prevalent among Panthers fans compared to NBA fans"
            ],
            "is_emerging": True
        },
        {
            "name": "Very Long Insights Test (Max Height)",
            "category_name": "Athleisure",
            "insights": [
                "Carolina Panthers fans represent one of the most engaged athleisure consumer segments in professional sports, with 89% of fans reporting at least one athleisure purchase in the past 90 days compared to just 52% of the local gen pop, demonstrating the massive market opportunity for athletic apparel brands looking to connect with this passionate fanbase",
                "The purchasing frequency data reveals that Panthers fans make an average of 14.7 athleisure-related purchases per year across categories including performance apparel, athletic footwear, accessories, and fitness equipment, which is nearly double the 7.8 annual purchases made by the average local consumer",
                "From a spending perspective, Panthers fans invest an average of $1,247 annually in athleisure products, representing a 94% premium over the local gen pop average of $642, with particular strength in premium brands and limited edition team-affiliated merchandise that commands higher price points",
                "Brand loyalty metrics show Panthers fans maintain relationships with an average of 4.3 distinct athleisure brands, significantly higher than the 2.6 brand average among local consumers",
                "The index analysis shows Panthers fans are 372% more likely to purchase from emerging direct-to-consumer athleisure brands compared to NBA fans overall, suggesting an openness to new brand partnerships and innovative marketing approaches that could benefit forward-thinking sponsors",
                "Premium athleisure spending is 428% higher among Panthers fans compared to NBA average, with particular strength in performance footwear and technical apparel categories"
            ],
            "is_emerging": False
        },
        {
            "name": "Minimal Insights with NBA Comparison",
            "category_name": "Auto",
            "insights": [
                "42% of fans own trucks",
                "Average 2.3 vehicles per household",
                "Panthers fans are 67% more likely to purchase American brands compared to NBA fans",
                "Truck ownership is 89% higher compared to NBA average"
            ],
            "is_emerging": False
        }
    ]

    # Create slides for each test case
    for i, test_case in enumerate(test_cases):
        # Add a blank slide
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Add test case title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = f"Test {i+1}: {test_case['name']}"
        p.font.size = Pt(18)
        p.font.bold = True

        # Create mock analysis results
        analysis_results = {
            'display_name': test_case['category_name'],
            'slide_title': f"Category Analysis: {test_case['category_name']}",
            'insights': test_case['insights'],
            'is_emerging': test_case['is_emerging'],
            'category_metrics': MockCategoryMetrics(),
            'subcategory_stats': None  # We'll skip subcategory table for this test
        }

        # Use the actual _add_category_insights method
        slide_generator._add_category_insights(
            slide,
            analysis_results,
            team_config['team_name_short'],
            team_config,
            is_emerging=test_case['is_emerging']
        )

        # Add visual guides to show spacing
        add_spacing_guides(slide, test_case['is_emerging'])

        # Add measurement annotations
        add_measurement_annotations(slide, test_case['is_emerging'])

        # Add debug information
        regular_insights = [i for i in test_case['insights'] if "NBA" not in i or "compared to" not in i]
        nba_insights = [i for i in test_case['insights'] if "NBA" in i and "compared to" in i]

        debug_text = f"Regular insights: {len(regular_insights)} bullets\n"
        debug_text += f"NBA insights: {len(nba_insights)} bullets\n"
        debug_text += f"Avg chars (regular): {sum(len(i) for i in regular_insights) // len(regular_insights) if regular_insights else 0}\n"
        debug_text += f"Total chars: {sum(len(i) for i in test_case['insights'])}"

        debug_box = slide.shapes.add_textbox(Inches(10), Inches(1), Inches(3), Inches(2))
        debug_frame = debug_box.text_frame
        debug_frame.text = debug_text
        debug_frame.paragraphs[0].font.size = Pt(10)
        debug_frame.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)

    # Save the presentation
    filename = "category_insights_spacing_test.pptx"
    prs.save(filename)
    print(f"Test presentation saved as '{filename}'")
    print("\nTest cases created:")
    for i, test in enumerate(test_cases):
        regular = len([ins for ins in test['insights'] if "NBA" not in ins or "compared to" not in ins])
        nba = len([ins for ins in test['insights'] if "NBA" in ins and "compared to" in ins])
        print(f"  Slide {i+1}: {test['name']} - {regular} regular + {nba} NBA insights")
    print("\nCheck each slide to verify:")
    print("  1. No overlap between insights box and NBA section")
    print("  2. Consistent 0.3\" spacing between sections")
    print("  3. NBA section appears when NBA insights are present")
    print("  4. Dynamic positioning based on insights box height")


def add_spacing_guides(slide, is_emerging):
    """Add visual guides to help verify spacing"""
    # Add vertical ruler lines at key positions
    positions = [
        (0.5, "0.5\" (left margin)"),
        (5.2, "5.2\" (insights right edge)"),
        (6.2, "6.2\" (table left edge)")
    ]

    for x_pos, label in positions:
        line = slide.shapes.add_connector(
            1,  # Straight line
            Inches(x_pos), Inches(1),
            Inches(x_pos), Inches(7)
        )
        line.line.color.rgb = RGBColor(200, 200, 200)
        line.line.width = Pt(0.5)

        # Add label
        label_box = slide.shapes.add_textbox(
            Inches(x_pos + 0.05), Inches(0.5),
            Inches(1.5), Inches(0.3)
        )
        label_box.text_frame.text = label
        p = label_box.text_frame.paragraphs[0]
        p.font.size = Pt(8)
        p.font.color.rgb = RGBColor(150, 150, 150)


def add_measurement_annotations(slide, is_emerging):
    """Add annotations to show dynamic spacing measurements"""
    # Starting Y positions
    insights_y = Inches(3.6) if is_emerging else Inches(3.2)

    # Add horizontal line to show where insights box starts
    start_line = slide.shapes.add_connector(
        1,  # Straight line
        Inches(5.3), insights_y,
        Inches(5.8), insights_y
    )
    start_line.line.color.rgb = RGBColor(255, 0, 0)
    start_line.line.width = Pt(1)

    # Add label
    label_box = slide.shapes.add_textbox(
        Inches(5.9), insights_y - Inches(0.1),
        Inches(2), Inches(0.3)
    )
    label_box.text_frame.text = "Insights Start"
    p = label_box.text_frame.paragraphs[0]
    p.font.size = Pt(8)
    p.font.color.rgb = RGBColor(255, 0, 0)


class MockCategoryMetrics:
    """Mock category metrics for testing"""
    def format_percent_fans(self):
        return "42%"

    def format_likelihood(self):
        return "125% More"

    def format_purchases(self):
        return "8 More"


def run_height_calculation_tests():
    """Test the height calculation function with various text lengths"""
    print("\nTesting height calculation function:")
    print("-" * 50)

    slide_generator = CategorySlide()

    test_texts = [
        ("Short", "• Quick insight\n• Another point\n• Third item"),
        ("Medium", "• This is a medium length insight that will wrap to multiple lines\n• Another medium length point with more detail\n• Third point with additional context"),
        ("Long", "• This is a very long insight that contains extensive detail about fan behavior patterns and will definitely need to wrap across multiple lines to fit within the available width\n• Another extremely detailed point that discusses spending patterns and demographic information\n• Yet another comprehensive insight with statistical analysis"),
    ]

    for name, text in test_texts:
        height = slide_generator._calculate_text_height_estimate(text, 12, 4.5)
        lines = text.count('\n') + 1
        print(f"{name} text:")
        print(f"  Lines: {lines}")
        print(f"  Characters: {len(text)}")
        print(f"  Estimated height: {height:.2f} inches")
        print(f"  With padding (0.4): {height + 0.4:.2f} inches")
        print()


if __name__ == "__main__":
    # Run height calculation tests
    run_height_calculation_tests()

    # Create test presentation
    create_test_presentation()

    print("\nDone! Open 'category_insights_spacing_test.pptx' to verify dynamic spacing.")