# slide_generators/category_slide.py
"""
Generate category analysis slides for PowerPoint presentations
Formats category spending data into professional slides matching SIL template
"""

from pathlib import Path
from typing import Dict, Optional, Any, List, Tuple
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
import pandas as pd
import logging
import re
from PIL import Image, ImageDraw
import io

from .base_slide import BaseSlide
from data_processors.category_analyzer import CategoryAnalyzer, CategoryMetrics
from utils.logo_manager import LogoManager

logger = logging.getLogger(__name__)


# Formatting utility functions
def format_percentage_no_decimal(value):
    """Format percentage without decimals, handle 0% as EQUAL"""
    if isinstance(value, str):
        # Extract number from string like "10.5%"
        try:
            num = float(value.replace('%', '').strip())
        except:
            return value
    else:
        num = float(value)

    # Handle 0% special case
    if num == 0:
        return "EQUAL"

    # Round to nearest whole number
    return f"{int(round(num))}%"


def format_percent_of_fans(value):
    """Format 'Percent of Fans Who Spend' column with <5% for small values"""
    if isinstance(value, str):
        # Extract number from string like "4.5%"
        try:
            num = float(value.replace('%', '').strip())
        except:
            return value
    else:
        num = float(value)

    # If less than 5%, return "<5%"
    if num < 5:
        return "<5%"

    # Otherwise, round to nearest whole number
    return f"{int(round(num))}%"


def format_currency_no_cents(value):
    """Format currency without cents and with commas"""
    if isinstance(value, str):
        # Extract number from string like "$2,536.17"
        try:
            num = float(value.replace('$', '').replace(',', '').strip())
        except:
            return value
    else:
        num = float(value)

    return f"${int(round(num)):,}"


def clean_text_references(text):
    """Remove (Excl. Jazz Fans) and similar references"""
    # Remove various exclusion patterns
    patterns = [
        r'\(Excl\. [^)]*\)',
        r'\(excl\. [^)]*\)',
        r'\(Excluding [^)]*\)',
        r'\(excluding [^)]*\)'
    ]

    for pattern in patterns:
        text = re.sub(pattern, '', text)

    # Clean up extra spaces
    text = re.sub(r'\s+', ' ', text).strip()

    return text


def format_gen_pop_references(text):
    """Convert 'gen pop' to 'local gen pop' without duplication"""
    # First handle cases where "local gen pop" might already exist (fix duplications)
    text = re.sub(r'\blocal\s+local\s+gen\s+pop\b', 'local gen pop', text, flags=re.IGNORECASE)

    # Handle "Local Gen Pop" (capitalized) -> "local gen pop"
    text = re.sub(r'\bLocal\s+Gen\s+Pop\b', 'local gen pop', text)

    # Then handle standalone "gen pop" that doesn't already have "local"
    text = re.sub(r'(?<!local\s)\bgen pop\b', 'local gen pop', text, flags=re.IGNORECASE)

    return text


def clean_subcategory_duplication(text, category_name):
    """Remove category name duplication in subcategory references"""
    # Convert category name to different case variations for matching
    category_variations = [
        category_name,
        category_name.lower(),
        category_name.upper(),
        category_name.title()
    ]

    # Remove "Category - Subcategory" patterns and just keep "Subcategory"
    for cat_var in category_variations:
        # Pattern: "Category - Subcategory" -> "Subcategory"
        pattern = rf'\b{re.escape(cat_var)}\s*-\s*'
        text = re.sub(pattern, '', text, flags=re.IGNORECASE)

    return text.strip()


def remove_unnecessary_comparisons(text):
    """Remove comparison language when no actual comparison is provided"""
    # Remove "vs. the Local Gen Pop" when it's just stating a fact without comparison
    # Look for patterns like "spend $X vs. the Local Gen Pop" without actual comparison

    # Pattern: "spend $amount vs. the Local Gen Pop" -> "spend $amount"
    text = re.sub(r'spend\s+(\$[\d,]+)\s+vs\.\s+the\s+Local\s+Gen\s+Pop', r'spend \1', text, flags=re.IGNORECASE)

    # Pattern: "spend $amount per year vs. the Local Gen Pop" -> "spend $amount per year"
    text = re.sub(r'spend\s+(\$[\d,]+)\s+per\s+year\s+vs\.\s+the\s+Local\s+Gen\s+Pop', r'spend \1 per year',
                  text, flags=re.IGNORECASE)

    return text


def lowercase_fan_references(text):
    """Convert 'Fan' or 'Fans' to lowercase except at start of sentences"""
    # Pattern to match 'fan' or 'fans' not at the start of a sentence
    # Negative lookbehind for sentence start indicators
    text = re.sub(r'(?<![.!?]\s)(?<![.!?]\s\s)(?<!^)(?<!^\s)\b(Fan|Fans)\b',
                  lambda m: m.group(1).lower(), text)

    return text


def process_insight_text_enhanced(insight_text: str, category_name: str = "") -> str:
    """Enhanced process insight text with all formatting rules including subcategory fixes"""
    # Apply all formatting rules
    text = clean_text_references(insight_text)
    text = format_gen_pop_references(text)

    # Fix capitalization
    text = text.replace(' MORE ', ' more ')

    # Clean subcategory duplication if category name provided
    if category_name:
        text = clean_subcategory_duplication(text, category_name)

    # Remove unnecessary comparison language
    text = remove_unnecessary_comparisons(text)

    # Convert Fan/Fans to lowercase (except at sentence start)
    text = lowercase_fan_references(text)

    # Format any embedded percentages and currency
    def format_pct_match(match):
        return format_percentage_no_decimal(match.group(0))

    text = re.sub(r'\d+\.?\d*%', format_pct_match, text)

    # Find and format currency
    def format_curr_match(match):
        return format_currency_no_cents(match.group(0))

    text = re.sub(r'\$[\d,]+\.?\d*', format_curr_match, text)

    return text


def process_insight_text(insight_text: str) -> str:
    """Process insight text with all formatting rules - legacy function for compatibility"""
    return process_insight_text_enhanced(insight_text)


def format_subcategory_text(category_name: str, subcategory_name: str) -> str:
    """Format subcategory text to avoid duplication"""
    # If subcategory already contains the category name, just return subcategory
    if category_name.lower() in subcategory_name.lower():
        return subcategory_name

    # Otherwise return the full name
    return f"{category_name} - {subcategory_name}"


class CategorySlide(BaseSlide):
    """Generate category analysis slides with insights, subcategory stats, and merchant rankings"""

    def __init__(self, presentation: Presentation = None):
        """
        Initialize category slide generator

        Args:
            presentation: Existing presentation to add slide to
        """
        super().__init__(presentation)

        # Initialize LogoManager
        self.logo_manager = LogoManager()

        # Colors for the slide - UPDATED with EQUAL color and emerging category
        self.colors = {
            'header_bg': RGBColor(240, 240, 240),
            'header_border': RGBColor(200, 200, 200),
            'table_header': RGBColor(0, 0, 0),  # Black background for headers
            'table_border': RGBColor(0, 0, 0),
            'positive': RGBColor(0, 176, 80),  # Green
            'negative': RGBColor(255, 0, 0),  # Red
            'equal': RGBColor(184, 134, 11),  # Dark yellow for EQUAL
            'neutral': RGBColor(0, 0, 0),  # Black
            'emerging_bg': RGBColor(217, 217, 217)  # Light gray for emerging category
        }

    def generate(self,
                 analysis_results: Dict[str, Any],
                 team_config: Dict[str, Any],
                 slide_index: Optional[int] = None) -> Presentation:
        """
        Generate the complete category analysis slide

        Args:
            analysis_results: Results from CategoryAnalyzer.analyze_category()
            team_config: Team configuration with colors, names, etc.
            slide_index: Where to insert slide (None = append)

        Returns:
            Updated presentation object
        """
        # Extract team info
        team_name = team_config.get('team_name', 'Team')
        team_short = team_config.get('team_name_short', team_name.split()[-1])

        # Check if this is an emerging category
        is_emerging = analysis_results.get('is_emerging', False)

        # Use the content layout (SIL white layout #12)
        slide = self.add_content_slide()
        logger.info(f"Added category slide for {analysis_results['display_name']} using SIL white layout")

        # Add header
        self._add_header(slide, team_name, analysis_results['slide_title'])

        # Add title with special handling for emerging categories
        if is_emerging:
            self._add_emerging_category_title(slide, analysis_results['display_name'])
        else:
            category_title = f"Category Analysis: {analysis_results['display_name']}"
            self._add_title(slide, category_title)

        # Add category insights (left side) - adjusted position for emerging categories
        self._add_category_insights(slide, analysis_results, team_short, team_config, is_emerging=is_emerging)

        # Add category metrics table (top right) - pass results for category name
        self._add_category_table(slide, analysis_results)

        # Add subcategory table (middle right) - adjusted for 16:9
        self._add_subcategory_table(slide, analysis_results['subcategory_stats'])

        logger.info(f"Generated {analysis_results['display_name']} slide")
        return self.presentation

    def generate_brand_slide(self,
                             analysis_results: Dict[str, Any],
                             team_config: Dict[str, Any],
                             slide_index: Optional[int] = None) -> Presentation:
        """
        Generate the brand analysis slide (second slide for category)

        Args:
            analysis_results: Results from CategoryAnalyzer
            team_config: Team configuration
            slide_index: Where to insert slide

        Returns:
            Updated presentation
        """
        # Extract team info
        team_name = team_config.get('team_name', 'Team')
        team_short = team_config.get('team_name_short', team_name.split()[-1])
        category_name = analysis_results['display_name']

        # Use the content layout (SIL white layout #12)
        slide = self.add_content_slide()
        logger.info(f"Added brand slide for {analysis_results['display_name']} using SIL white layout")

        # UPDATED brand slide title format based on category
        if category_name.upper() == "QSR":
            brand_title = f"Top QSR Brands for {team_name} Fans"
        else:
            brand_title = f"Top {category_name} Brands for {team_name} Fans"

        # Add header - brand slide uses updated title format
        header_title = f"Sponsor Spending Analysis: {category_name} Brands"
        self._add_header(slide, team_name, header_title)

        # Add title with updated format
        self._add_title(slide, brand_title)

        # Add brand logos (numbered circles) - adjusted for 16:9
        self._add_brand_logos(slide, analysis_results['merchant_stats'])

        # Add brand insights (left side) - UPDATED with formatting
        self._add_brand_insights(slide, analysis_results, team_name, team_short, category_name)

        # Add brand table (right side) - adjusted for 16:9
        self._add_brand_table(slide, analysis_results['merchant_stats'])

        # Add sponsorship recommendation - UPDATED with formatting
        self._add_sponsor_recommendation(slide, analysis_results['recommendation'], team_config)

        logger.info(f"Generated {analysis_results['display_name']} brand slide")
        return self.presentation

    def _add_header(self, slide, team_name: str, slide_title: str):
        """Add header with team name and slide title"""
        # Header background (adjusted for 16:9) - matching behaviors slide
        header_rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(13.333), Inches(0.5)  # Full 16:9 width
        )
        header_rect.fill.solid()
        header_rect.fill.fore_color.rgb = RGBColor(240, 240, 240)
        header_rect.line.color.rgb = RGBColor(200, 200, 200)
        header_rect.line.width = Pt(0.5)

        # Team name (left)
        team_text = slide.shapes.add_textbox(
            Inches(0.2), Inches(0.1),
            Inches(3), Inches(0.3)
        )
        team_text.text_frame.text = team_name
        p = team_text.text_frame.paragraphs[0]
        p.font.name = self.default_font  # Red Hat Display
        p.font.size = Pt(14)
        p.font.bold = True

        # Slide indicator (right)
        slide_text = slide.shapes.add_textbox(
            Inches(6.5), Inches(0.1),  # Adjusted for 16:9
            Inches(6.633), Inches(0.3)  # Adjusted width
        )
        slide_text.text_frame.text = slide_title
        p = slide_text.text_frame.paragraphs[0]
        p.font.name = self.default_font  # Red Hat Display
        p.alignment = PP_ALIGN.RIGHT
        p.font.size = Pt(14)

    def _add_title(self, slide, title: str):
        """Add main slide title with appropriate width and automatic line wrapping"""

        # Determine width based on title type
        if "Top" in title and "Brands for" in title:
            # Brand slide - title should end where brand table starts
            width = Inches(5.3)  # Table starts at 5.833, so stop just before
        elif "Category Analysis:" in title:
            # Category slide - title should end where category table starts
            width = Inches(5.3)  # Table starts at 6.2, so stop just before
        else:
            # Default width for other slides
            width = Inches(12.333)

        # Make text box taller to accommodate potential second line
        title_box = slide.shapes.add_textbox(
            Inches(0.6), Inches(1.3),
            width, Inches(1.8)  # Increased height from 1.0 to 1.8 for two lines
        )

        # Configure text frame for proper text wrapping
        text_frame = title_box.text_frame
        text_frame.word_wrap = True  # This enables automatic line wrapping
        text_frame.auto_size = MSO_AUTO_SIZE.NONE  # Prevent box from auto-resizing

        # Remove margins to maximize text space
        text_frame.margin_left = Inches(0)
        text_frame.margin_right = Inches(0)
        text_frame.margin_top = Inches(0)
        text_frame.margin_bottom = Inches(0)

        # Split title for category slides
        if "Category Analysis:" in title:
            category_name = title.replace("Category Analysis: ", "")
            text_frame.text = f"Category Analysis:\n{category_name}"
        else:
            text_frame.text = title

        # Format all paragraphs with working bold+italic font
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = "Red Hat Display"  # Use base font name
                run.font.size = Pt(26)
                run.font.bold = True
                run.font.italic = True  # Now this works!
                run.font.color.rgb = RGBColor(0, 0, 0)
            paragraph.line_spacing = 1.0  # Adjust line spacing if needed

    def _add_emerging_category_title(self, slide, category_name: str):
        """Add special title formatting for emerging categories"""

        # Main title: "Emerging Category:"
        title_box = slide.shapes.add_textbox(
            Inches(0.6), Inches(1.3),
            Inches(5.3), Inches(0.5)
        )

        text_frame = title_box.text_frame
        text_frame.text = "Emerging Category:"

        p = text_frame.paragraphs[0]
        p.font.name = "Red Hat Display"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.italic = True
        p.font.color.rgb = RGBColor(0, 0, 0)

        # Category name on second line
        category_box = slide.shapes.add_textbox(
            Inches(0.6), Inches(1.8),
            Inches(5.3), Inches(0.6)
        )

        text_frame = category_box.text_frame
        text_frame.text = category_name

        p = text_frame.paragraphs[0]
        p.font.name = "Red Hat Display"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.italic = True
        p.font.color.rgb = RGBColor(0, 0, 0)

        # Explanatory subtext
        subtext_box = slide.shapes.add_textbox(
            Inches(0.6), Inches(2.4),
            Inches(5.3), Inches(0.6)
        )

        text_frame = subtext_box.text_frame
        text_frame.word_wrap = True
        text_frame.text = (
            "The Emerging Category is where at least 10% of your fans are "
            "spending, but where there isn't one clear brand leader, and the "
            "category has room to grow"
        )

        p = text_frame.paragraphs[0]
        p.font.name = "Red Hat Display"
        p.font.size = Pt(12)
        p.font.italic = True
        p.font.color.rgb = RGBColor(0, 0, 0)

    def _add_category_insights(self, slide, results: Dict[str, Any], team_short: str,
                               team_config: Dict[str, Any], is_emerging: bool = False):
        """Add category insights section - UPDATED WITH ALL FORMATTING FIXES"""

        # Adjust vertical position based on whether it's an emerging category
        title_y = Inches(3.2) if is_emerging else Inches(2.8)
        insights_y = Inches(3.6) if is_emerging else Inches(3.2)
        nba_label_y = Inches(5.8) if is_emerging else Inches(5.4)
        nba_box_y = Inches(6.2) if is_emerging else Inches(5.8)

        # Insights title
        insights_title = slide.shapes.add_textbox(
            Inches(0.5), title_y,
            Inches(4), Inches(0.3)
        )
        insights_title.text_frame.text = "Category Insights:"
        p = insights_title.text_frame.paragraphs[0]
        p.font.name = self.default_font
        p.font.size = Pt(14)
        p.font.bold = True

        # Get category name for subcategory cleaning
        category_name = results.get('display_name', '')

        # Separate insights into regular and NBA comparison
        regular_insights = []
        nba_insights = []

        for insight in results['insights']:
            # Clean and format the insight text with enhanced processing
            cleaned_insight = process_insight_text_enhanced(insight, category_name)

            if "NBA" in cleaned_insight and "compared to" in cleaned_insight:
                nba_insights.append(cleaned_insight)
            else:
                regular_insights.append(cleaned_insight)

        # Regular insights box
        insights_box = slide.shapes.add_textbox(
            Inches(0.7), insights_y,
            Inches(4.5), Inches(3.0)
        )

        text_frame = insights_box.text_frame
        text_frame.word_wrap = True

        # Add regular insights with BULLETS instead of numbers
        for i, insight in enumerate(regular_insights[:4]):
            p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
            p.text = f"• {insight}"  # Changed from f"{i+1}. {insight}"
            p.font.name = self.default_font
            p.font.size = Pt(12)
            p.line_spacing = 1.2

        # Add NBA comparison section if there are NBA insights
        if nba_insights:
            # NBA comparison label
            nba_label = slide.shapes.add_textbox(
                Inches(0.5), nba_label_y,
                Inches(4), Inches(0.3)
            )
            nba_label.text_frame.text = f"{team_short} Fans vs. {team_config.get('league', 'NBA')} Fans:"
            p = nba_label.text_frame.paragraphs[0]
            p.font.name = self.default_font
            p.font.size = Pt(14)
            p.font.bold = True

            # NBA insights box
            nba_box = slide.shapes.add_textbox(
                Inches(0.7), nba_box_y,
                Inches(4.5), Inches(1.2)
            )

            nba_text_frame = nba_box.text_frame
            nba_text_frame.word_wrap = True

            # Add NBA insights with BULLETS instead of numbers
            for i, insight in enumerate(nba_insights[:2]):
                p = nba_text_frame.add_paragraph() if i > 0 else nba_text_frame.paragraphs[0]
                p.text = f"• {insight}"  # Changed from numbered
                p.font.name = self.default_font
                p.font.size = Pt(12)
                p.line_spacing = 1.2

    def _add_category_table(self, slide, results: Dict[str, Any]):
        """Add category metrics table (adjusted for 16:9) with formatting fixes"""
        # Extract metrics from results
        metrics = results['category_metrics']
        is_emerging = results.get('is_emerging', False)

        # Table position - adjusted to prevent bleeding and moved down
        left = Inches(6.2)  # Moved left to fit better
        top = Inches(1.4)  # Moved down from 1.5
        width = Inches(6.8)  # Adjusted width to fit within slide
        height = Inches(0.8)  # Reduced height for better proportions

        # Adjust row height here
        row_height = Inches(0.7)  # Height per row (increase this value for taller rows)
        num_rows = 2  # Header + 1 data row
        height = row_height * num_rows  # Total table height

        # Create table
        table_shape = slide.shapes.add_table(2, 4, left, top, width, height)
        table = table_shape.table

        # Set column widths - matching subcategory table
        table.columns[0].width = Inches(1.6)  # Category
        table.columns[1].width = Inches(1.5)  # Percent of Fans
        table.columns[2].width = Inches(1.7)  # How likely
        table.columns[3].width = Inches(1.6)  # Purchases

        # Header row - UPDATED with consistent header
        headers = ['Category', 'Percent of Fans\nWho Spend', 'Likelihood to Spend\n(vs. Local Gen Pop)',
                   'Purchases Per Fan\n(vs. Local Gen Pop)']

        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            self._format_header_cell(cell)

        # Data row - extract category name properly and apply formatting
        category_name = results.get('display_name', 'Category')
        table.cell(1, 0).text = category_name

        # UPDATED: Use format_percent_of_fans for "Percent of Fans Who Spend"
        percent_value = metrics.format_percent_fans()
        table.cell(1, 1).text = format_percent_of_fans(percent_value)

        table.cell(1, 2).text = self._format_metric_value(metrics.format_likelihood())
        table.cell(1, 3).text = self._format_metric_value(metrics.format_purchases())

        # Format data cells - special handling for emerging categories
        for i in range(4):
            cell = table.cell(1, i)
            self._format_data_cell(cell)

            # Add gray background for emerging category data row
            if is_emerging and i == 0:  # Only for the category name cell
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(217, 217, 217)  # Light gray background

    def _add_subcategory_table(self, slide, subcategory_stats: pd.DataFrame):
        """Add subcategory statistics table (adjusted for 16:9) with formatting fixes"""
        if subcategory_stats.empty:
            return

        # Table position - adjusted to prevent bleeding and extend down
        left = Inches(6.2)  # Moved left to fit better
        top = Inches(3.0)  # Adjusted vertical position (was 2.7)
        width = Inches(6.8)  # Adjusted width to fit within slide

        # Create table with header + data rows
        rows = min(len(subcategory_stats), 5) + 1  # Max 5 subcategories + header

        # Calculate height to extend down the slide more
        row_height = Inches(.7)  # Increased from 0.3
        table_height = row_height * rows

        table_shape = slide.shapes.add_table(rows, 4, left, top, width, table_height)
        table = table_shape.table

        # Set column widths - better distribution
        table.columns[0].width = Inches(1.6)  # Sub-Category
        table.columns[1].width = Inches(1.5)  # Percent of Fans
        table.columns[2].width = Inches(1.7)  # How likely
        table.columns[3].width = Inches(1.6)  # Purchases

        # Headers - UPDATED with consistent header
        headers = ['Sub-Category', 'Percent of Fans\nWho Spend', 'Likelihood to Spend\n(vs. Local Gen Pop)',
                   'Purchases Per Fan\n(vs. Local Gen Pop)']

        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            self._format_header_cell(cell)

        # Data rows with formatting
        for row_idx, (_, row) in enumerate(subcategory_stats.iterrows(), 1):
            if row_idx >= rows:
                break

            # Apply subcategory text formatting to avoid duplication
            subcategory_text = row['Subcategory']
            table.cell(row_idx, 0).text = subcategory_text

            # UPDATED: Use format_percent_of_fans for "Percent of Fans Who Spend"
            percent_fans = row['Percent of Fans Who Spend']
            table.cell(row_idx, 1).text = format_percent_of_fans(percent_fans)

            table.cell(row_idx, 2).text = self._format_metric_value(
                row.get('Likelihood to spend (vs. Local Gen Pop)', row.get('Likelihood to spend vs. gen pop', '')))
            table.cell(row_idx, 3).text = self._format_metric_value(row.get('Purchases Per Fan (vs. Gen Pop)', ''))

            # Format cells
            for col in range(4):
                self._format_data_cell(table.cell(row_idx, col))

        # ADD EXPLANATORY TEXT BELOW THE TABLE
        # Calculate position below the table
        explanation_top = top + table_height + Inches(0.1)  # Add small gap after table

        # Add explanatory text
        explanation_box = slide.shapes.add_textbox(
            left, explanation_top,
            width, Inches(0.3)
        )

        text_frame = explanation_box.text_frame
        text_frame.word_wrap = True

        # Set the explanatory text
        text_frame.text = "Subcategories shown in descending order by composite index"

        # Format the text
        p = text_frame.paragraphs[0]
        p.font.name = self.default_font  # Red Hat Display
        p.font.size = Pt(10)
        p.font.italic = True
        p.font.color.rgb = RGBColor(100, 100, 100)  # Gray color for de-emphasis
        p.alignment = PP_ALIGN.LEFT  # Left align to match table

    def _has_colored_background(self, image: Image.Image, threshold: int = 240) -> bool:
        """
        Check if an image has a colored (non-white) background

        Args:
            image: PIL Image to check
            threshold: RGB value threshold below which we consider it colored (default 240)

        Returns:
            True if image has colored background, False if white/transparent
        """
        # Convert to RGBA if not already
        if image.mode != 'RGBA':
            image = image.convert('RGBA')

        width, height = image.size

        # Sample more points around the edges to better detect background
        sample_points = []

        # Add corners
        sample_points.extend([
            (0, 0), (width - 1, 0), (0, height - 1), (width - 1, height - 1)
        ])

        # Add edge midpoints
        sample_points.extend([
            (width // 2, 0), (width // 2, height - 1),  # top and bottom middle
            (0, height // 2), (width - 1, height // 2)  # left and right middle
        ])

        # Add points slightly inward from edges (to avoid anti-aliasing artifacts)
        edge_offset = min(5, width // 10, height // 10)
        sample_points.extend([
            (edge_offset, edge_offset),
            (width - edge_offset - 1, edge_offset),
            (edge_offset, height - edge_offset - 1),
            (width - edge_offset - 1, height - edge_offset - 1)
        ])

        colored_pixels = 0
        total_opaque_pixels = 0

        for x, y in sample_points:
            try:
                r, g, b, a = image.getpixel((x, y))

                # Skip transparent pixels
                if a < 128:
                    continue

                total_opaque_pixels += 1

                # Check if this pixel is colored (not white/gray)
                # A colored pixel has significant variation in RGB values or low values
                if (r < threshold or g < threshold or b < threshold) or \
                        (max(r, g, b) - min(r, g, b) > 30):
                    colored_pixels += 1

            except:
                continue

        # If more than 50% of opaque edge pixels are colored, consider it a colored background
        if total_opaque_pixels > 0:
            return colored_pixels / total_opaque_pixels > 0.5

        return False

    def _add_brand_logos(self, slide, merchant_stats: Tuple[pd.DataFrame, List[str]]):
        """Add brand logos aligned with the brand table below"""
        merchant_df, top_merchants = merchant_stats

        if merchant_df.empty:
            return

        # Align with brand table dimensions - MATCH THE TABLE POSITION AND WIDTH
        table_left = Inches(6.6)  # Updated to match table left position
        table_width = Inches(6.0)  # Updated to match table width

        # Calculate logo positions
        num_logos = min(5, len(merchant_df))
        logo_size = (120, 120)  # Size in pixels for logo processing
        display_size = Inches(1.0)  # Logo display size

        # Calculate spacing between logos
        # Account for logo width when calculating available space
        available_space = table_width - (display_size * num_logos)
        spacing_between = available_space / (num_logos - 1) if num_logos > 1 else 0

        # Position for logos - above the table
        y = Inches(1.2)  # Shifted down from 1.0 to 1.5

        # Add logos for top 5 brands
        for i in range(num_logos):
            # Calculate x position for each logo
            x = table_left + (i * (display_size + spacing_between))
            merchant_name = merchant_df.iloc[i]['Brand']

            # Try to get logo from LogoManager
            logo_image = self.logo_manager.get_logo(merchant_name, size=logo_size)

            # If no logo found, create fallback with initials
            if logo_image is None:
                logger.info(f"No logo found for {merchant_name}, creating fallback")
                logo_image = self.logo_manager.create_fallback_logo(
                    merchant_name,
                    size=logo_size,
                    bg_color='white',
                    text_color='#888888'
                )

            # Check if logo has colored background
            has_colored_bg = self._has_colored_background(logo_image)

            # Log the detection result for debugging
            logger.debug(f"{merchant_name} - Colored background detected: {has_colored_bg}")

            if has_colored_bg:
                # For colored background logos, create a circular mask to ensure clean edges
                masked_logo = Image.new('RGBA', logo_size, (0, 0, 0, 0))

                # Create circular mask
                mask = Image.new('L', logo_size, 0)
                draw = ImageDraw.Draw(mask)
                draw.ellipse([0, 0, logo_size[0] - 1, logo_size[1] - 1], fill=255)

                # Apply mask to logo
                masked_logo.paste(logo_image, (0, 0))
                masked_logo.putalpha(mask)

                # Convert masked logo to bytes
                image_stream = io.BytesIO()
                masked_logo.save(image_stream, format='PNG')
                image_stream.seek(0)

                # Add the logo directly without circle background
                try:
                    pic = slide.shapes.add_picture(
                        image_stream,
                        x, y,
                        display_size, display_size
                    )
                except Exception as e:
                    logger.error(f"Failed to add logo for {merchant_name}: {e}")
                    self._add_numbered_circle_fallback(slide, i + 1, x, y, display_size)

            else:
                # For white/transparent background logos, add circle border
                circle = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    x, y,
                    display_size, display_size
                )
                circle.fill.solid()
                circle.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background
                circle.line.color.rgb = RGBColor(200, 200, 200)  # Light gray border
                circle.line.width = Pt(1)

                # Make the logo slightly smaller to fit within the circle border
                logo_display_size = Inches(0.9)  # Slightly smaller than the circle
                offset = (display_size - logo_display_size) / 2  # Center the logo

                # Convert PIL Image to bytes for PowerPoint
                image_stream = io.BytesIO()
                logo_image.save(image_stream, format='PNG')
                image_stream.seek(0)

                # Add the logo centered within the circle
                try:
                    pic = slide.shapes.add_picture(
                        image_stream,
                        x + offset,
                        y + offset,
                        logo_display_size, logo_display_size
                    )

                    # Ensure the logo is on top of the circle
                    slide.shapes._spTree.remove(pic._element)
                    slide.shapes._spTree.append(pic._element)

                except Exception as e:
                    logger.error(f"Failed to add logo for {merchant_name}: {e}")
                    self._add_numbered_circle_fallback(slide, i + 1, x, y, display_size)

            # Add ranking number below each logo
            number_box = slide.shapes.add_textbox(
                x, y + display_size + Inches(0.05),
                display_size, Inches(0.3)
            )
            number_box.text_frame.text = str(i + 1)
            p = number_box.text_frame.paragraphs[0]
            p.font.name = self.default_font
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = RGBColor(100, 100, 100)

    def _add_numbered_circle_fallback(self, slide, number: int, x: float, y: float, size: float):
        """Add numbered circle as ultimate fallback if logo processing fails"""
        # Circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            x, y,
            size, size
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = RGBColor(255, 255, 255)
        circle.line.color.rgb = RGBColor(200, 200, 200)
        circle.line.width = Pt(2)

        # Number
        text_box = slide.shapes.add_textbox(
            x, y + (size - Inches(0.4)) / 2,
            size, Inches(0.4)
        )
        text_box.text_frame.text = str(number)
        p = text_box.text_frame.paragraphs[0]
        p.font.name = self.default_font
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(48)
        p.font.color.rgb = RGBColor(150, 150, 150)

    def _add_brand_insights(self, slide, results: Dict[str, Any], team_name: str, team_short: str, category_name: str, team_config: Dict[str, Any]):
        """Add brand-specific insights with updated formatting to match reference"""
        # Top Brand Insights section
        insights_title = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.6),
            Inches(4), Inches(0.3)
        )
        insights_title.text_frame.text = "Top Brand Insights"
        p = insights_title.text_frame.paragraphs[0]
        p.font.name = self.default_font  # Red Hat Display
        p.font.size = Pt(14)
        p.font.bold = True

        # Insights box - extended width to be closer to logos/table
        insights_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(3.0),
            Inches(5.5), Inches(1.6)  # Extended width from 4.5 to 5.5
        )

        text_frame = insights_box.text_frame
        text_frame.word_wrap = True

        # Get merchant insights
        merchant_insights = results.get('merchant_insights', [])
        merchant_df, _ = results.get('merchant_stats', (pd.DataFrame(), []))

        # Format insights based on reference image pattern
        for i, insight in enumerate(merchant_insights[:4]):
            p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
            p.line_spacing = 1.2

            # Determine which type of insight this is and format accordingly
            if i == 0:  # First insight - Highest % of Fans
                # Add the label in bold
                run1 = p.add_run()
                run1.text = f"• Highest % of Fans: "
                run1.font.name = self.default_font
                run1.font.size = Pt(12)  # Increased from 10 to match Top Brand Target
                run1.font.bold = True

                # Add the rest of the insight
                run2 = p.add_run()
                # Extract percentage and brand from data if available
                if not merchant_df.empty:
                    top_brand = merchant_df.iloc[0]['Brand']
                    percent = merchant_df.iloc[0]['Percent of Fans Who Spend']
                    # UPDATED: Use format_percent_of_fans here too
                    formatted_percent = format_percent_of_fans(percent)
                    run2.text = f"{formatted_percent} of {team_name} fans spend at {top_brand}"
                else:
                    # Fallback to formatting the existing insight
                    formatted_insight = process_insight_text(insight)
                    run2.text = formatted_insight.replace("• ", "")
                run2.font.name = self.default_font
                run2.font.size = Pt(12)  # Match Top Brand Target font size
                run2.font.bold = False

            elif i == 1:  # Second insight - Most Purchases per Fan
                # Add the label in bold
                run1 = p.add_run()
                run1.text = "• Most Purchases per Fan: "
                run1.font.name = self.default_font
                run1.font.size = Pt(12)
                run1.font.bold = True

                # Add the rest of the insight
                run2 = p.add_run()
                # Extract and format the purchases per year info
                if 'purchases per year' in insight.lower():
                    formatted_insight = self._format_insight_two(insight, team_short, category_name)
                    # Remove the team name from start if it's there
                    if formatted_insight.startswith(f"{team_short} fans"):
                        formatted_insight = formatted_insight.replace(f"{team_short} fans", f"{team_short} fans", 1)
                    run2.text = formatted_insight
                else:
                    formatted_insight = process_insight_text(insight)
                    run2.text = formatted_insight.replace("• ", "")
                run2.font.name = self.default_font
                run2.font.size = Pt(12)
                run2.font.bold = False

            elif i == 2:  # Third insight - Highest Spend per Fan
                # Add the label in bold
                run1 = p.add_run()
                run1.text = "• Highest Spend per Fan: "
                run1.font.name = self.default_font
                run1.font.size = Pt(12)
                run1.font.bold = True

                # Add the rest of the insight
                run2 = p.add_run()
                # Format the spend per fan info
                formatted_insight = process_insight_text(insight)
                # Extract just the relevant part after the bullet
                insight_text = formatted_insight.replace("• ", "")
                # Add [Team] notation if team name is mentioned
                if team_name in insight_text:
                    insight_text = insight_text.replace(f"{team_name} fans", f"{team_short} fans")
                run2.text = insight_text
                run2.font.name = self.default_font
                run2.font.size = Pt(12)
                run2.font.bold = False

            elif i == 3:  # Fourth insight - NBA/League comparison
                # Add the label in bold
                run1 = p.add_run()
                run1.text = f"• Highest % of Fans Index vs {team_config.get('league', 'NBA')}: "
                run1.font.name = self.default_font
                run1.font.size = Pt(12)
                run1.font.bold = True

                # Add the rest of the insight
                run2 = p.add_run()
                # Format the NBA comparison
                formatted_insight = process_insight_text(insight)
                insight_text = formatted_insight.replace("• ", "")
                # Simplify the comparison text
                if "more likely" in insight_text.lower():
                    # Extract the percentage and brand info
                    import re
                    match = re.search(r'(\d+)%\s+more\s+likely.*?(?:on|at)\s+([^"]+?)(?:\s+than|\s+when)', insight_text,
                                      re.IGNORECASE)
                    if match:
                        percent = match.group(1)
                        brand = match.group(2).strip()
                        league = team_config.get('league', 'NBA')
                        run2.text = f"{team_short} fans are {percent}% more likely to spend on {brand} than the average {league} fan"
                    else:
                        run2.text = insight_text
                else:
                    run2.text = insight_text
                run2.font.name = self.default_font
                run2.font.size = Pt(12)
                run2.font.bold = False

    def _add_brand_table(self, slide, merchant_stats: Tuple[pd.DataFrame, List[str]]):
        """Add brand ranking table - adjusted for 16:9 with reduced width"""
        merchant_df, _ = merchant_stats

        if merchant_df.empty:
            return

        # Create table - adjusted position and size for 16:9
        rows = min(len(merchant_df) + 1, 6)  # Header + max 5 brands
        cols = 5
        left = Inches(6.6)  # Keep same left position
        top = Inches(3.0)  # Shifted down from 2.8 to 3.3
        width = Inches(6.0)  # REDUCED from 7.0 to 6.0

        # Calculate height with proper row spacing
        row_height = Inches(0.6)  # or whatever you're using
        height = row_height * rows

        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table

        # ADJUSTED column widths proportionally for smaller table
        table.columns[0].width = Inches(1.0)  # Rank (reduced from 1.3)
        table.columns[1].width = Inches(1.4)  # Brand (reduced from 1.6)
        table.columns[2].width = Inches(1.2)  # Percent of Fans (reduced from 1.5)
        table.columns[3].width = Inches(1.2)  # How likely (reduced from 1.7)
        table.columns[4].width = Inches(1.2)  # Purchases (reduced from 1.6)

        # Headers - UPDATED with consistent header
        headers = ['Rank (by\npercent of\nfans who\nspend)', 'Brand',
                   'Percent of\nFans Who\nSpend',
                   'Likelihood to Spend\n(vs. Local Gen Pop)',
                   'Purchases Per Fan\n(vs. Local Gen Pop)']

        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            self._format_header_cell(cell, small=True)

        # Data rows with formatting
        for row_idx, (_, row) in enumerate(merchant_df.iterrows(), 1):
            if row_idx >= rows:
                break

            table.cell(row_idx, 0).text = str(row['Rank'])
            table.cell(row_idx, 1).text = row['Brand']

            # UPDATED: Use format_percent_of_fans for "Percent of Fans Who Spend"
            table.cell(row_idx, 2).text = format_percent_of_fans(row['Percent of Fans Who Spend'])

            table.cell(row_idx, 3).text = self._format_metric_value(row.get('Likelihood to spend (vs. Local Gen Pop)',
                                                                            row.get(
                                                                                'How likely fans are to spend vs. gen pop',
                                                                                '')))
            table.cell(row_idx, 4).text = self._format_metric_value(row.get('Purchases Per Fan (vs. Gen Pop)', ''))

            # Format cells
            for col in range(5):
                self._format_data_cell(table.cell(row_idx, col), small=True)

    def _add_sponsor_recommendation(self, slide, recommendation: Optional[Dict[str, Any]], team_config: Dict[str, Any]):
        """Add Top Brand Target and sponsorship recommendation with updated formatting"""
        if not recommendation:
            return

        # Top Brand Target header with colon
        target_title = slide.shapes.add_textbox(
            Inches(0.5), Inches(5.0),
            Inches(2.0), Inches(0.3)  # Wider to accommodate full text
        )
        target_title.text_frame.text = "Hot Brand Target:"
        p = target_title.text_frame.paragraphs[0]
        p.font.name = self.default_font  # Red Hat Display
        p.font.size = Pt(14)
        p.font.bold = True

        # Add brand logo next to "Hot Brand Target:" text
        merchant_name = recommendation.get('merchant', 'Brand')
        logo_size = Inches(0.5)  # Small logo size
        logo_x = Inches(2.4)  # Shifted right to avoid text overlap
        logo_y = Inches(4.925)  # Vertically centered with text

        # Try to get logo from LogoManager
        logo_image = self.logo_manager.get_logo(merchant_name, size=(60, 60))

        if logo_image:
            # Check if logo has colored background
            has_colored_bg = self._has_colored_background(logo_image)

            # Convert PIL Image to bytes for PowerPoint
            image_stream = io.BytesIO()
            logo_image.save(image_stream, format='PNG')
            image_stream.seek(0)

            try:
                if has_colored_bg:
                    # For colored background logos, add directly
                    pic = slide.shapes.add_picture(
                        image_stream,
                        logo_x, logo_y,
                        logo_size, logo_size
                    )
                else:
                    # For white/transparent background logos, add with circle border
                    circle = slide.shapes.add_shape(
                        MSO_SHAPE.OVAL,
                        logo_x, logo_y,
                        logo_size, logo_size
                    )
                    circle.fill.solid()
                    circle.fill.fore_color.rgb = RGBColor(255, 255, 255)
                    circle.line.color.rgb = RGBColor(200, 200, 200)
                    circle.line.width = Pt(0.5)

                    # Reset stream position
                    image_stream.seek(0)

                    # Add logo slightly smaller to fit in circle
                    logo_display_size = Inches(0.45)
                    offset = (logo_size - logo_display_size) / 2
                    pic = slide.shapes.add_picture(
                        image_stream,
                        logo_x + offset,
                        logo_y + offset,
                        logo_display_size, logo_display_size
                    )

                    # Ensure logo is on top
                    slide.shapes._spTree.remove(pic._element)
                    slide.shapes._spTree.append(pic._element)

            except Exception as e:
                logger.error(f"Failed to add logo for {merchant_name}: {e}")
        else:
            logger.info(f"No logo found for Hot Brand Target: {merchant_name}")

        # Recommendation content - moved down to create space after logo
        rec_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(5.5),  # Moved down from 5.5 to create spacing
            Inches(5.5), Inches(1.2)  # Extended width to match insights box
        )

        text_frame = rec_box.text_frame
        text_frame.word_wrap = True

        # UPDATED: Standardized recommendation format with composite index
        team_name = team_config.get('team_name', 'Team')
        composite_index_raw = recommendation.get('composite_index', 0)
        try:
            composite_index = int(round(float(composite_index_raw)))
        except (ValueError, TypeError):
            composite_index = 'N/A'

        # First bullet - main recommendation with BOLD brand name
        p1 = text_frame.paragraphs[0]

        # Add the beginning of the sentence
        run1 = p1.add_run()
        run1.text = f"• The {team_name} should target "
        run1.font.name = self.default_font
        run1.font.size = Pt(12)
        run1.font.bold = False

        # Add the brand name in bold
        run2 = p1.add_run()
        run2.text = merchant_name
        run2.font.name = self.default_font
        run2.font.size = Pt(12)
        run2.font.bold = True

        # Add the rest of the sentence
        run3 = p1.add_run()
        run3.text = f" for a sponsorship based on having the highest composite index of {composite_index}"
        run3.font.name = self.default_font
        run3.font.size = Pt(12)
        run3.font.bold = False

        p1.line_spacing = 1.2

        # Second bullet - standardized explanation with proper indentation
        p2 = text_frame.add_paragraph()
        p2.text = "• The composite index indicates a brand with significant likelihood for more fans to be spending more frequently, and at a higher spend per fan vs. other brands"
        p2.font.name = self.default_font
        p2.font.size = Pt(12)
        p2.line_spacing = 1.2
        # Set both left indent and first line indent to ensure entire paragraph is indented
        p2.left_indent = Inches(0.25)  # Indent entire paragraph
        p2.first_line_indent = Inches(0)  # No additional indent for first line

    def _format_header_cell(self, cell, small: bool = False):
        """Format table header cell"""
        cell.fill.solid()
        cell.fill.fore_color.rgb = self.colors['table_header']  # Black background

        # Format text
        text_frame = cell.text_frame
        text_frame.margin_left = Inches(0.05)
        text_frame.margin_right = Inches(0.05)
        text_frame.margin_top = Inches(0.03)
        text_frame.margin_bottom = Inches(0.03)
        text_frame.word_wrap = True

        # Format all paragraphs in the cell (for multi-line headers)
        for paragraph in text_frame.paragraphs:
            paragraph.font.name = self.default_font  # Red Hat Display
            paragraph.font.size = Pt(8) if small else Pt(10)  # Smaller font for compact tables
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White text
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.line_spacing = 1.0  # Tighter line spacing for headers

        # Vertical alignment
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    def _format_data_cell(self, cell, small: bool = False):
        """Format table data cell with updated color handling"""
        # Format text
        text_frame = cell.text_frame
        text_frame.margin_left = Inches(0.05)
        text_frame.margin_right = Inches(0.05)
        text_frame.margin_top = Inches(0.03)
        text_frame.margin_bottom = Inches(0.03)
        text_frame.word_wrap = True

        # Format all paragraphs in the cell
        for paragraph in text_frame.paragraphs:
            paragraph.font.name = self.default_font  # Red Hat Display
            paragraph.font.size = Pt(12) if small else Pt(12)  # Consistent data font size
            paragraph.alignment = PP_ALIGN.CENTER

            # UPDATED color coding with EQUAL handling and bold for green
            if 'More' in cell.text or 'more' in cell.text:
                paragraph.font.color.rgb = self.colors['positive']  # Green
                paragraph.font.bold = True  # ADD THIS LINE to make green text bold
            elif 'Less' in cell.text or 'less' in cell.text or 'fewer' in cell.text:
                paragraph.font.color.rgb = self.colors['negative']  # Red
                paragraph.font.bold = False  # Keep red text not bold
            elif 'EQUAL' in cell.text or 'Equal' in cell.text:
                paragraph.font.color.rgb = self.colors['equal']  # Dark yellow
                paragraph.font.bold = False  # Keep yellow text not bold
            else:
                paragraph.font.bold = False  # Default: not bold

        # Vertical alignment
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    def _format_metric_value(self, value: str) -> str:
        """Apply formatting to metric values (percentages, currency, etc.)"""
        if not value or pd.isna(value):
            return "N/A"

        # Clean text references first
        formatted_value = clean_text_references(str(value))
        formatted_value = format_gen_pop_references(formatted_value)

        # Handle percentage formatting
        if '%' in formatted_value:
            formatted_value = format_percentage_no_decimal(formatted_value)

        # Handle currency formatting
        if '$' in formatted_value:
            formatted_value = format_currency_no_cents(formatted_value)

        return formatted_value

    def _format_insight_two(self, insight: str, team_short: str, category_name: str) -> str:
        """Format insight #2 with standardized format for Most Purchases per Fan"""
        # Extract relevant information using regex

        # Try to extract brand name and number of purchases
        brand_match = re.search(r'at\s+([A-Za-z0-9\s&\'-]+?)(?:\s+than|\s+compared|$)', insight, re.IGNORECASE)
        purchases_match = re.search(r'(\d+)\s+purchases?\s+per\s+year', insight, re.IGNORECASE)

        if brand_match and purchases_match:
            brand = brand_match.group(1).strip()
            purchases = purchases_match.group(1)

            # Remove trailing punctuation from brand name if any
            brand = brand.rstrip('.,;:')

            return f"{team_short} fans make an average of {purchases} purchases per year at {brand}"

        # Fallback to original insight with general formatting
        return process_insight_text(insight)

    def check_missing_logos(self, categories_data: Dict[str, Any]) -> Dict[str, List[str]]:
        """
        Check which logos are missing for all categories

        Args:
            categories_data: Dictionary of category analysis results

        Returns:
            Dictionary mapping category to list of missing logos
        """
        missing_logos = {}

        for category_name, analysis_results in categories_data.items():
            merchant_df, _ = analysis_results.get('merchant_stats', (pd.DataFrame(), []))
            if not merchant_df.empty:
                # Get top 5 merchants
                top_merchants = merchant_df.head(5)['Brand'].tolist()

                # Check which logos are missing
                report = self.logo_manager.add_missing_logos_report(top_merchants)
                missing = [merchant for merchant, has_logo in report.items() if not has_logo]

                if missing:
                    missing_logos[category_name] = missing
                    logger.warning(f"Missing logos for {category_name}: {missing}")

        return missing_logos


# Convenience functions
def create_category_slides(analysis_results: Dict[str, Any],
                           team_config: Dict[str, Any],
                           presentation: Optional[Presentation] = None) -> Presentation:
    """
    Create both category slides (analysis + brands) for a category

    Args:
        analysis_results: Results from CategoryAnalyzer
        team_config: Team configuration
        presentation: Existing presentation (creates new if None)

    Returns:
        Presentation with both slides added
    """
    generator = CategorySlide(presentation)

    # Generate category analysis slide
    presentation = generator.generate(analysis_results, team_config)

    # Generate brand analysis slide
    presentation = generator.generate_brand_slide(analysis_results, team_config)

    return presentation