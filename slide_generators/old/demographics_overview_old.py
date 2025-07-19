# slide_generators/demographic_overview_slide.py
"""
Generate demographic overview slide with AI-generated insights
Uses the blue SIL layout (#11) with team branding and AI insights text
FIXED VERSION: No RGBColor transparency errors, works with template layouts
Updated to use Overpass font for insights text
Enhanced with dynamic fan image loading from assets
"""

from pathlib import Path
from typing import Dict, Optional, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import logging

from .base_slide import BaseSlide

logger = logging.getLogger(__name__)

# Default fonts
DEFAULT_FONT_FAMILY = "Red Hat Display"
INSIGHTS_FONT_FAMILY = "Overpass"  # New font for insights


class DemographicOverviewSlide(BaseSlide):
    """Generate demographic overview slide with AI insights on blue background"""

    def __init__(self, presentation: Presentation = None):
        """
        Initialize demographic overview slide generator

        Args:
            presentation: Existing presentation to add slide to (REQUIRED for template access)
        """
        super().__init__(presentation)

        # Ensure we have a presentation with template layouts
        if presentation is None:
            raise ValueError("DemographicOverviewSlide requires an existing presentation with template layouts")

    def generate(self,
                 team_config: Dict[str, Any],
                 ai_insights: str,
                 slide_index: Optional[int] = None) -> Presentation:
        """
        Generate the demographic overview slide with AI insights

        Args:
            team_config: Team configuration with name and colors
            ai_insights: AI-generated demographic insights text
            slide_index: Where to insert slide (None = append)

        Returns:
            Updated presentation object
        """
        # Extract team info
        team_name = team_config.get('team_name', 'Team')
        team_colors = team_config.get('colors', {})

        # Use the blue layout (#11) - this should work since we have the template loaded
        slide = self.add_title_slide()  # This uses layout #11 (blue background)
        logger.info("Added demographic overview slide using SIL blue layout")

        # Add header bar at the top
        self._add_header_bar(slide, team_name)

        # Add main title
        self._add_main_title(slide, team_name)

        # Add AI insights text
        self._add_insights_text(slide, ai_insights)

        # Add fan image (replaces placeholder)
        self._add_fan_image(slide, team_name)

        logger.info(f"Generated demographic overview slide for {team_name}")
        return self.presentation

    def _add_header_bar(self, slide, team_name: str):
        """Add the standard header bar at the top of the slide"""
        # Header bar dimensions (full width at top)
        header_rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(13.333), Inches(0.5)  # Full 16:9 width
        )

        # Style the header bar with light gray background
        header_rect.fill.solid()
        header_rect.fill.fore_color.rgb = RGBColor(240, 240, 240)  # Light gray
        header_rect.line.color.rgb = RGBColor(200, 200, 200)
        header_rect.line.width = Pt(0.5)

        # Add team name on the left side of header
        team_text = slide.shapes.add_textbox(
            Inches(0.2), Inches(0.1),
            Inches(3), Inches(0.3)
        )
        team_text.text_frame.text = team_name
        team_p = team_text.text_frame.paragraphs[0]
        team_p.font.name = DEFAULT_FONT_FAMILY
        team_p.font.size = Pt(14)
        team_p.font.bold = True
        team_p.font.color.rgb = RGBColor(0, 0, 0)  # Dark gray

        # Add page title on the right side of header
        title_text = slide.shapes.add_textbox(
            Inches(6.5), Inches(0.1),
            Inches(6.633), Inches(0.3)
        )
        title_text.text_frame.text = f"Fan Demographics: How Are {team_name} Fans Unique"
        title_p = title_text.text_frame.paragraphs[0]
        title_p.font.name = DEFAULT_FONT_FAMILY
        title_p.font.size = Pt(14)
        title_p.alignment = PP_ALIGN.RIGHT
        title_p.font.color.rgb = RGBColor(0, 0, 0)  # Dark gray
        # Note: Right side text is NOT bold, which matches other slides

    def _add_main_title(self, slide, team_name: str):
        """Add the main title: Team Name + Fan Demographic Overview"""
        # Title positioned below header bar
        title_text = slide.shapes.add_textbox(
            Inches(1), Inches(2.0),
            Inches(11), Inches(1.2)
        )
        title_frame = title_text.text_frame
        title_frame.text = f"{team_name}\nFan Demographic Overview"
        title_frame.margin_left = Inches(0)
        title_frame.margin_right = Inches(0)
        title_frame.margin_top = Inches(0)
        title_frame.margin_bottom = Inches(0)
        title_frame.word_wrap = True

        # Style the title
        for i, paragraph in enumerate(title_frame.paragraphs):
            paragraph.alignment = PP_ALIGN.LEFT

            for run in paragraph.runs:
                run.font.name = DEFAULT_FONT_FAMILY
                run.font.color.rgb = RGBColor(255, 255, 255)  # White text on blue background
                run.font.bold = True

                # Different sizes for team name vs subtitle
                if i == 0:  # Team name
                    run.font.size = Pt(48)
                else:  # "Fan Demographic Overview"
                    run.font.size = Pt(28)

    def _add_insights_text(self, slide, ai_insights: str):
        """Add the AI-generated insights text with Overpass font"""
        # Position insights text below the title
        insights_text = slide.shapes.add_textbox(
            Inches(1), Inches(3.6),
            Inches(5), Inches(3)  # Left side, leaving room for image
        )
        insights_frame = insights_text.text_frame
        insights_frame.text = ai_insights
        insights_frame.margin_left = Inches(0)
        insights_frame.margin_right = Inches(0.2)
        insights_frame.margin_top = Inches(0)
        insights_frame.margin_bottom = Inches(0)
        insights_frame.word_wrap = True

        # Style the insights text with Overpass font
        insights_p = insights_frame.paragraphs[0]
        insights_p.alignment = PP_ALIGN.LEFT
        insights_p.line_spacing = 1.2  # Slightly more spacing for readability

        for run in insights_p.runs:
            run.font.name = INSIGHTS_FONT_FAMILY  # Using Overpass for insights
            run.font.size = Pt(12)
            run.font.color.rgb = RGBColor(255, 255, 255)  # White text on blue background
            run.font.bold = False

    def _add_fan_image(self, slide, team_name: str):
        """Add the team fan image from assets directory"""
        # Convert team name to filename format (remove spaces)
        team_filename = team_name.replace(" ", "")

        # Build the image path
        # Try multiple ways to find the project root
        current_file = Path(__file__).resolve()

        # Method 1: If we're in slide_generators/, go up one level
        if current_file.parent.name == 'slide_generators':
            project_root = current_file.parent.parent
        # Method 2: If we're deeper (e.g., in a submodule), go up until we find 'assets'
        else:
            project_root = current_file.parent
            while project_root.parent != project_root:  # Not at filesystem root
                if (project_root / 'assets').exists():
                    break
                project_root = project_root.parent

        logger.info(f"Project root resolved to: {project_root}")

        # Try multiple common image extensions without the extension in the base name
        image_extensions = ['.png', '.jpg', '.jpeg', '.PNG', '.JPG', '.JPEG']
        image_path = None

        base_filename = f"{team_filename}_Fanpic"
        fans_dir = project_root / "assets" / "logos" / "fans"

        logger.info(f"Looking for fan image in: {fans_dir}")
        logger.info(f"Base filename: {base_filename}")

        for ext in image_extensions:
            potential_path = fans_dir / f"{base_filename}{ext}"
            logger.debug(f"Checking: {potential_path}")
            if potential_path.exists():
                image_path = potential_path
                logger.info(f"Found fan image at: {image_path}")
                break

        if image_path and image_path.exists():
            try:
                # Add the image as a circular shape
                # Position on the right side, matching the placeholder location
                left = Inches(7.167)
                top = Inches(1.5)
                width = Inches(5)
                height = Inches(5)

                # Add the picture
                picture = slide.shapes.add_picture(
                    str(image_path),
                    left, top, width, height
                )

                # Apply circular crop to the image
                # This creates a circular mask effect
                picture.crop_left = 0
                picture.crop_right = 0
                picture.crop_top = 0
                picture.crop_bottom = 0

                # Note: PowerPoint doesn't directly support circular cropping via python-pptx
                # The image will be square/rectangular, but we can add a circular border

                # Add a white circular border overlay
                border_circle = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    left, top, width, height
                )

                # Make the circle transparent with white border
                border_circle.fill.background()  # Transparent fill
                border_circle.line.color.rgb = RGBColor(255, 255, 255)
                border_circle.line.width = Pt(3)

                # Move border to front
                slide.shapes._spTree.remove(border_circle._element)
                slide.shapes._spTree.append(border_circle._element)

                logger.info(f"Successfully added fan image for {team_name}")

            except Exception as e:
                logger.error(f"Error adding fan image: {str(e)}")
                # Fall back to placeholder if image fails
                self._add_image_placeholder_fallback(slide)
        else:
            logger.warning(f"Fan image not found for {team_name}")
            logger.warning(f"Searched in: {fans_dir}")
            logger.warning(f"Expected filename pattern: {base_filename}.[png/jpg/jpeg]")
            # Add fallback placeholder
            self._add_image_placeholder_fallback(slide)

    def _add_image_placeholder_fallback(self, slide):
        """Add a circular placeholder for team/fan image as fallback"""
        # Create circular placeholder on the right side
        placeholder_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(7.167), Inches(1.5),
            Inches(5), Inches(5)  # Circular shape
        )

        # Style the placeholder - RGBColor only takes RGB values, no alpha
        placeholder_circle.fill.solid()
        placeholder_circle.fill.fore_color.rgb = RGBColor(200, 200, 255)  # Light blue instead of transparent
        placeholder_circle.line.color.rgb = RGBColor(255, 255, 255)
        placeholder_circle.line.width = Pt(3)

        # Add placeholder text inside the circle
        placeholder_text = slide.shapes.add_textbox(
            Inches(8.67), Inches(3.5),
            Inches(2), Inches(1)
        )
        placeholder_frame = placeholder_text.text_frame
        placeholder_frame.text = "Team Image\nPlaceholder"
        placeholder_frame.margin_left = Inches(0)
        placeholder_frame.margin_right = Inches(0)
        placeholder_frame.margin_top = Inches(0)
        placeholder_frame.margin_bottom = Inches(0)

        # Style placeholder text - no alpha channel support
        for paragraph in placeholder_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER

            for run in paragraph.runs:
                run.font.name = DEFAULT_FONT_FAMILY
                run.font.size = Pt(12)
                run.font.color.rgb = RGBColor(100, 100, 150)  # Darker blue instead of transparent
                run.font.italic = True