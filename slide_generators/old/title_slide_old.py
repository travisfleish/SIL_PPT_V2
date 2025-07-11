# slide_generators/title_slide.py
"""
Generate title slide for PowerPoint presentations
Creates the opening slide with team branding and report title
"""

from pathlib import Path
from typing import Dict, Optional, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import logging

from .base_slide import BaseSlide

logger = logging.getLogger(__name__)


class TitleSlide(BaseSlide):
    """Generate the title slide for the presentation"""

    def __init__(self, presentation: Presentation = None):
        """
        Initialize title slide generator

        Args:
            presentation: Existing presentation to add slide to
        """
        super().__init__(presentation)

    def generate(self,
                 team_config: Dict[str, Any],
                 subtitle: str = "Sponsorship Insights Report",
                 slide_index: Optional[int] = None) -> Presentation:
        """
        Generate the title slide

        Args:
            team_config: Team configuration with name and colors
            subtitle: Subtitle for the presentation
            slide_index: Where to insert slide (None = append)

        Returns:
            Updated presentation object
        """
        # Extract team info
        team_name = team_config.get('team_name', 'Team')
        team_colors = team_config.get('colors', {})

        # Use the title layout (SIL blue layout #11)
        slide = self.add_title_slide()
        logger.info("Added title slide using SIL blue layout")

        # The SIL blue layout already has the background, so we don't need to add it
        # Only add background if using fallback layout
        if self.title_layout == self.presentation.slide_layouts[0]:  # Standard title layout
            # Add background color manually for non-SIL templates
            background = slide.background
            fill = background.fill
            fill.solid()

            # Use team primary color or default
            if 'primary' in team_colors:
                primary_color = team_colors['primary']
                if isinstance(primary_color, str):
                    fill.fore_color.rgb = self._hex_to_rgb(primary_color)
                else:
                    fill.fore_color.rgb = primary_color
            else:
                fill.fore_color.rgb = RGBColor(0, 34, 68)  # Default blue

        # Add team logo circle
        logo_x = Inches(6.4)  # Slightly adjusted for 16:9
        logo_y = Inches(0.8)
        logo_size = Inches(1.2)

        # Outer circle (accent color)
        outer_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            logo_x - Inches(0.08),
            logo_y - Inches(0.08),
            logo_size + Inches(0.16),
            logo_size + Inches(0.16)
        )
        outer_circle.fill.solid()

        # Use team accent color or default
        if 'accent' in team_colors:
            accent_color = team_colors['accent']
            if isinstance(accent_color, str):
                outer_circle.fill.fore_color.rgb = self._hex_to_rgb(accent_color)
            else:
                outer_circle.fill.fore_color.rgb = accent_color
        else:
            outer_circle.fill.fore_color.rgb = RGBColor(255, 182, 18)  # Default yellow

        outer_circle.line.fill.background()

        # Inner circle (white background for logo)
        inner_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            logo_x, logo_y,
            logo_size, logo_size
        )
        inner_circle.fill.solid()
        inner_circle.fill.fore_color.rgb = RGBColor(255, 255, 255)
        inner_circle.line.fill.background()

        # Add team initials or logo
        logo_text = slide.shapes.add_textbox(
            logo_x, logo_y,
            logo_size, logo_size
        )
        text_frame = logo_text.text_frame
        text_frame.clear()

        # Get team initials
        team_initials = team_config.get('team_initials',
                                        ''.join([word[0] for word in team_name.split()[:2]]))

        p = text_frame.add_paragraph()
        p.text = team_initials
        p.font.name = self.default_font  # Red Hat Display
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 34, 68)  # Dark blue
        p.alignment = PP_ALIGN.CENTER

        # Center vertically
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Add main title (team name)
        title_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(2.8),  # Adjusted for 16:9
            Inches(10.333), Inches(1.5)  # Wider for 16:9
        )
        title_box.text_frame.text = team_name
        title_box.text_frame.word_wrap = True

        # Format title
        for paragraph in title_box.text_frame.paragraphs:
            paragraph.font.name = self.default_font  # Red Hat Display
            paragraph.font.size = Pt(48)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White text
            paragraph.alignment = PP_ALIGN.CENTER

        # Add subtitle
        subtitle_box = slide.shapes.add_textbox(
            Inches(2), Inches(4.5),  # Adjusted positioning
            Inches(9.333), Inches(1)  # Wider for 16:9
        )
        subtitle_box.text_frame.text = subtitle
        subtitle_box.text_frame.word_wrap = True

        # Format subtitle
        for paragraph in subtitle_box.text_frame.paragraphs:
            paragraph.font.name = self.default_font  # Red Hat Display
            paragraph.font.size = Pt(28)
            paragraph.font.bold = False
            paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White text
            paragraph.alignment = PP_ALIGN.CENTER

        # Add date in bottom right
        from datetime import datetime
        current_date = datetime.now().strftime("%B %Y")

        date_box = slide.shapes.add_textbox(
            Inches(10.5), Inches(6.8),  # Bottom right for 16:9
            Inches(2.5), Inches(0.5)
        )
        date_box.text_frame.text = current_date

        # Format date
        for paragraph in date_box.text_frame.paragraphs:
            paragraph.font.name = self.default_font  # Red Hat Display
            paragraph.font.size = Pt(14)
            paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White text
            paragraph.alignment = PP_ALIGN.RIGHT

        logger.info(f"Generated title slide for {team_name}")
        return self.presentation

    def _hex_to_rgb(self, hex_color: str) -> RGBColor:
        """Convert hex color to RGBColor"""
        return self.hex_to_rgb(hex_color)  # Use parent method


# Convenience function for backward compatibility
def create_title_slide(team_config: Dict[str, Any],
                       subtitle: str = "Sponsorship Insights Report",
                       presentation: Optional[Presentation] = None) -> Presentation:
    """
    Create a title slide with team branding

    Args:
        team_config: Team configuration dictionary
        subtitle: Subtitle text
        presentation: Existing presentation to add to

    Returns:
        Presentation with title slide
    """
    generator = TitleSlide(presentation)
    return generator.generate(team_config, subtitle)