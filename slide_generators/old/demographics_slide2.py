# slide_generators/demographics_slide2.py
"""
Generate second demographics slide (Income, Children, Occupation)
"""

from pathlib import Path
from typing import Dict, Optional, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import logging

from .base_slide import BaseSlide

logger = logging.getLogger(__name__)


class DemographicsSlide2(BaseSlide):
    """Generate the second demographics slide with Income, Children, and Occupation charts"""

    def __init__(self, presentation: Presentation = None):
        """Initialize demographics slide 2 generator"""
        super().__init__(presentation)

    def generate(self,
                 demographic_data: Dict[str, Any],
                 chart_dir: Path,
                 team_config: Dict[str, Any],
                 slide_index: Optional[int] = None) -> Presentation:
        """
        Generate the second demographics slide

        Args:
            demographic_data: Processed demographic data
            chart_dir: Directory containing chart images
            team_config: Team configuration
            slide_index: Where to insert slide (None = append)

        Returns:
            Updated presentation object
        """
        # Extract team info
        team_name = team_config.get('team_name', 'Team')
        team_short = team_config.get('team_name_short', team_name.split()[-1])
        colors = team_config.get('colors', {})

        # Use the content layout
        slide = self.add_content_slide()
        logger.info("Added demographics slide 2 using content layout")

        # Add header
        self._add_header(slide, team_name)

        # Add demographic charts for slide 2
        self._add_slide2_charts(slide, chart_dir)

        logger.info(f"Generated demographics slide 2 for {team_name}")
        return self.presentation

    def _add_header(self, slide, team_name: str):
        """Add header with title"""
        # Header background
        header_rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(13.333), Inches(0.5)
        )
        header_rect.fill.solid()
        header_rect.fill.fore_color.rgb = RGBColor(240, 240, 240)
        header_rect.line.color.rgb = RGBColor(200, 200, 200)
        header_rect.line.width = Pt(0.5)

        # Team name (left)
        team_text = slide.shapes.add_textbox(
            Inches(0.2), Inches(0.1),
            Inches(3), Inches(0.3)
        )
        team_text.text_frame.text = team_name
        p = team_text.text_frame.paragraphs[0]
        p.font.name = self.default_font
        p.font.size = Pt(14)
        p.font.bold = True

        # Slide title (right) - continued from slide 1
        title_text = slide.shapes.add_textbox(
            Inches(6), Inches(0.1),
            Inches(7.133), Inches(0.3)
        )
        title_text.text_frame.text = f"Fan Demographics: How Are {team_name} Fans Unique (continued)"
        p = title_text.text_frame.paragraphs[0]
        p.font.name = self.default_font
        p.font.size = Pt(14)
        p.alignment = PP_ALIGN.RIGHT

    def _add_slide2_charts(self, slide, chart_dir: Path):
        """Add the three charts for slide 2: Income, Children, Occupation"""
        chart_dir = Path(chart_dir)

        # Chart positions for slide 2 - larger layout for 3 charts
        chart_positions = [
            # Top - Income (full width)
            ('income_chart', 1.5, 1.0, 10.5, 2.8),

            # Bottom row - Occupation (left) and Children (right)
            ('occupation_chart', 1.5, 4.2, 5.0, 2.8),
            ('children_chart', 7.0, 4.2, 5.0, 2.8)
        ]

        for chart_name, left, top, width, height in chart_positions:
            # Try both regular and hires versions
            chart_path = chart_dir / f'{chart_name}_hires.png'
            if not chart_path.exists():
                chart_path = chart_dir / f'{chart_name}.png'

            if chart_path.exists():
                try:
                    slide.shapes.add_picture(
                        str(chart_path),
                        Inches(left), Inches(top),
                        width=Inches(width), height=Inches(height)
                    )
                    logger.info(f"Added {chart_name} to slide 2")
                except Exception as e:
                    logger.warning(f"Could not add {chart_name}: {e}")
            else:
                logger.warning(f"Chart not found: {chart_path}")
                # Add placeholder
                placeholder = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(left), Inches(top),
                    Inches(width), Inches(height)
                )
                placeholder.fill.solid()
                placeholder.fill.fore_color.rgb = RGBColor(245, 245, 245)
                placeholder.line.color.rgb = RGBColor(100, 100, 100)
                placeholder.line.width = Pt(1)