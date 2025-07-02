# slide_generators/demographics_slide1.py
"""
Generate first demographics slide (Gender, Ethnicity, Generation)
"""

from pathlib import Path
from typing import Dict, Optional, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import logging

from .base_slide import BaseSlide

logger = logging.getLogger(__name__)


class DemographicsSlide1(BaseSlide):
    """Generate the first demographics slide with Gender, Ethnicity, and Generation charts"""

    def __init__(self, presentation: Presentation = None):
        """Initialize demographics slide 1 generator"""
        super().__init__(presentation)

    def generate(self,
                 demographic_data: Dict[str, Any],
                 chart_dir: Path,
                 team_config: Dict[str, Any],
                 slide_index: Optional[int] = None) -> Presentation:
        """
        Generate the first demographics slide

        Args:
            demographic_data: Processed demographic data
            chart_dir: Directory containing chart images
            team_config: Team configuration
            slide_index: Where to insert slide (None = append)

        Returns:
            Updated presentation object
        """
        # Extract team info
        team_name = team_config.get('team_name', 'Team')
        team_short = team_config.get('team_name_short', team_name.split()[-1])
        colors = team_config.get('colors', {})
        league = team_config.get('league', 'League')

        # Use the content layout
        slide = self.add_content_slide()
        logger.info("Added demographics slide 1 using content layout")

        # Add header
        self._add_header(slide, team_name)

        # Add team logo/circle (left side)
        self._add_team_logo(slide, team_short, colors)

        # Add insight text
        self._add_insight_text(slide, demographic_data.get('key_insights', ''))

        # Add demographic charts for slide 1
        self._add_slide1_charts(slide, chart_dir)

        # Add KEY/legend box
        self._add_legend_box(slide, team_name, team_short, league)

        logger.info(f"Generated demographics slide 1 for {team_name}")
        return self.presentation

    def _add_header(self, slide, team_name: str):
        """Add header with title"""
        # Header background
        header_rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(13.333), Inches(0.5)
        )
        header_rect.fill.solid()
        header_rect.fill.fore_color.rgb = RGBColor(240, 240, 240)
        header_rect.line.color.rgb = RGBColor(200, 200, 200)
        header_rect.line.width = Pt(0.5)

        # Team name (left)
        team_text = slide.shapes.add_textbox(
            Inches(0.2), Inches(0.1),
            Inches(3), Inches(0.3)
        )
        team_text.text_frame.text = team_name
        p = team_text.text_frame.paragraphs[0]
        p.font.name = self.default_font
        p.font.size = Pt(14)
        p.font.bold = True

        # Slide title (right)
        title_text = slide.shapes.add_textbox(
            Inches(6), Inches(0.1),
            Inches(7.133), Inches(0.3)
        )
        title_text.text_frame.text = f"Fan Demographics: How Are {team_name} Fans Unique"
        p = title_text.text_frame.paragraphs[0]
        p.font.name = self.default_font
        p.font.size = Pt(14)
        p.alignment = PP_ALIGN.RIGHT

    def _add_team_logo(self, slide, team_short: str, colors: Dict[str, str]):
        """Add team logo circle"""
        # Create circle shape
        left = Inches(0.5)
        top = Inches(1.5)
        size = Inches(1.8)

        # Outer circle (secondary color)
        outer_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            left - Inches(0.05), top - Inches(0.05),
            size + Inches(0.1), size + Inches(0.1)
        )
        outer_circle.fill.solid()
        outer_circle.fill.fore_color.rgb = self.hex_to_rgb(colors.get('secondary', '#FFB612'))
        outer_circle.line.fill.background()

        # Inner circle (primary color)
        inner_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            left, top,
            size, size
        )
        inner_circle.fill.solid()
        inner_circle.fill.fore_color.rgb = self.hex_to_rgb(colors.get('primary', '#002244'))
        inner_circle.line.fill.background()

        # Team name text
        text_box = slide.shapes.add_textbox(
            left, top + Inches(0.6),
            size, Inches(0.6)
        )
        text_box.text_frame.text = team_short
        p = text_box.text_frame.paragraphs[0]
        p.font.name = self.default_font
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

    def _add_insight_text(self, slide, insight: str):
        """Add insight text below logo"""
        text_box = slide.shapes.add_textbox(
            Inches(0.3), Inches(3.8),
            Inches(2.2), Inches(2.0)
        )

        # Default insight if none provided
        if not insight:
            insight = "Jazz fans are younger, and more likely to be parents who are working professionals versus the Utah gen pop."

        text_box.text_frame.text = insight
        text_box.text_frame.word_wrap = True

        p = text_box.text_frame.paragraphs[0]
        p.font.name = self.default_font
        p.font.size = Pt(14)
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT

    def _add_slide1_charts(self, slide, chart_dir: Path):
        """Add the three charts for slide 1: Gender, Ethnicity, Generation"""
        chart_dir = Path(chart_dir)

        # Chart positions for slide 1 - larger since only 3 charts
        chart_positions = [
            # Top row - Generation (left) and Gender (right)
            ('generation_chart', 3.0, 1.0, 4.5, 2.5),
            ('gender_chart', 8.0, 1.0, 4.5, 2.5),

            # Bottom row - Ethnicity (centered)
            ('ethnicity_chart', 5.5, 4.0, 5.0, 2.5)
        ]

        for chart_name, left, top, width, height in chart_positions:
            # Try both regular and hires versions
            chart_path = chart_dir / f'{chart_name}_hires.png'
            if not chart_path.exists():
                chart_path = chart_dir / f'{chart_name}.png'

            if chart_path.exists():
                try:
                    slide.shapes.add_picture(
                        str(chart_path),
                        Inches(left), Inches(top),
                        width=Inches(width), height=Inches(height)
                    )
                    logger.info(f"Added {chart_name} to slide 1")
                except Exception as e:
                    logger.warning(f"Could not add {chart_name}: {e}")
            else:
                logger.warning(f"Chart not found: {chart_path}")
                # Add placeholder
                placeholder = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(left), Inches(top),
                    Inches(width), Inches(height)
                )
                placeholder.fill.solid()
                placeholder.fill.fore_color.rgb = RGBColor(245, 245, 245)
                placeholder.line.color.rgb = RGBColor(100, 100, 100)
                placeholder.line.width = Pt(1)

    def _add_legend_box(self, slide, team_name: str, team_short: str, league: str):
        """Add KEY legend box"""
        # Box position
        left = Inches(0.3)
        top = Inches(5.8)
        width = Inches(4.5)
        height = Inches(1.2)

        # Create box
        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, top, width, height
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(255, 255, 255)
        box.line.color.rgb = RGBColor(0, 0, 0)
        box.line.width = Pt(1)

        # Add KEY text
        text_frame = box.text_frame
        text_frame.clear()
        text_frame.margin_left = Inches(0.1)
        text_frame.margin_top = Inches(0.1)

        # KEY title
        p = text_frame.add_paragraph()
        p.text = "KEY"
        p.font.name = self.default_font
        p.font.size = Pt(12)
        p.font.bold = True

        # Legend items
        items = [
            f"- {team_name} Fans",
            f"- {team_short} Gen Pop (state level, excluding {team_short} Fans)",
            f"- {league} Fans Total (excluding {team_short} fans)"
        ]

        for item in items:
            p = text_frame.add_paragraph()
            p.text = item
            p.font.name = self.default_font
            p.font.size = Pt(10)
            p.level = 0