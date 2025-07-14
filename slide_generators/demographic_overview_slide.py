# slide_generators/demographic_overview_slide.py
"""
Generate demographic overview slide with AI-generated insights
Uses the blue SIL layout (#11) with team branding and AI insights text
FIXED VERSION: No RGBColor transparency errors, works with template layouts
"""

from pathlib import Path
from typing import Dict, Optional, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import logging

from .base_slide import BaseSlide

logger = logging.getLogger(__name__)

# Default font
DEFAULT_FONT_FAMILY = "Red Hat Display"


class DemographicOverviewSlide(BaseSlide):
    """Generate demographic overview slide with AI insights on blue background"""

    def __init__(self, presentation: Presentation = None):
        """
        Initialize demographic overview slide generator

        Args:
            presentation: Existing presentation to add slide to (REQUIRED for template access)
        """
        super().__init__(presentation)

        # Ensure we have a presentation with template layouts
        if presentation is None:
            raise ValueError("DemographicOverviewSlide requires an existing presentation with template layouts")

    def generate(self,
                 team_config: Dict[str, Any],
                 ai_insights: str,
                 slide_index: Optional[int] = None) -> Presentation:
        """
        Generate the demographic overview slide with AI insights

        Args:
            team_config: Team configuration with name and colors
            ai_insights: AI-generated demographic insights text
            slide_index: Where to insert slide (None = append)

        Returns:
            Updated presentation object
        """
        # Extract team info
        team_name = team_config.get('team_name', 'Team')
        team_colors = team_config.get('colors', {})

        # Use the blue layout (#11) - this should work since we have the template loaded
        slide = self.add_title_slide()  # This uses layout #11 (blue background)
        logger.info("Added demographic overview slide using SIL blue layout")

        # Add header bar at the top
        self._add_header_bar(slide, team_name)

        # Add main title
        self._add_main_title(slide, team_name)

        # Add AI insights text
        self._add_insights_text(slide, ai_insights)

        # Add image placeholder
        self._add_image_placeholder(slide)

        logger.info(f"Generated demographic overview slide for {team_name}")
        return self.presentation

    def _add_header_bar(self, slide, team_name: str):
        """Add the standard header bar at the top of the slide"""
        # Header bar dimensions (full width at top)
        header_rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(13.333), Inches(0.5)  # Full 16:9 width
        )

        # Style the header bar with light gray background
        header_rect.fill.solid()
        header_rect.fill.fore_color.rgb = RGBColor(240, 240, 240)  # Light gray
        header_rect.line.color.rgb = RGBColor(200, 200, 200)
        header_rect.line.width = Pt(0.5)

        # Add team name on the left side of header
        team_text = slide.shapes.add_textbox(
            Inches(0.2), Inches(0.1),
            Inches(3), Inches(0.3)
        )
        team_text.text_frame.text = team_name
        team_p = team_text.text_frame.paragraphs[0]
        team_p.font.name = DEFAULT_FONT_FAMILY
        team_p.font.size = Pt(14)
        team_p.font.bold = True
        team_p.font.color.rgb = RGBColor(0, 0, 0)  # Dark gray

        # Add page title on the right side of header
        title_text = slide.shapes.add_textbox(
            Inches(6.5), Inches(0.1),
            Inches(6.633), Inches(0.3)
        )
        title_text.text_frame.text = f"Fan Demographics: How Are {team_name} Fans Unique"
        title_p = title_text.text_frame.paragraphs[0]
        title_p.font.name = DEFAULT_FONT_FAMILY
        title_p.font.size = Pt(14)
        title_p.alignment = PP_ALIGN.RIGHT
        title_p.font.color.rgb = RGBColor(0, 0, 0)  # Dark gray
        # Note: Right side text is NOT bold, which matches other slides
    def _add_main_title(self, slide, team_name: str):
        """Add the main title: Team Name + Fan Demographic Overview"""
        # Title positioned below header bar
        title_text = slide.shapes.add_textbox(
            Inches(1), Inches(2.0),
            Inches(11), Inches(1.2)
        )
        title_frame = title_text.text_frame
        title_frame.text = f"{team_name}\nFan Demographic Overview"
        title_frame.margin_left = Inches(0)
        title_frame.margin_right = Inches(0)
        title_frame.margin_top = Inches(0)
        title_frame.margin_bottom = Inches(0)
        title_frame.word_wrap = True

        # Style the title
        for i, paragraph in enumerate(title_frame.paragraphs):
            paragraph.alignment = PP_ALIGN.LEFT

            for run in paragraph.runs:
                run.font.name = DEFAULT_FONT_FAMILY
                run.font.color.rgb = RGBColor(255, 255, 255)  # White text on blue background
                run.font.bold = True

                # Different sizes for team name vs subtitle
                if i == 0:  # Team name
                    run.font.size = Pt(48)
                else:  # "Fan Demographic Overview"
                    run.font.size = Pt(28)

    def _add_insights_text(self, slide, ai_insights: str):
        """Add the AI-generated insights text"""
        # Position insights text below the title
        insights_text = slide.shapes.add_textbox(
            Inches(1), Inches(3.6),
            Inches(5), Inches(3)  # Left side, leaving room for image
        )
        insights_frame = insights_text.text_frame
        insights_frame.text = ai_insights
        insights_frame.margin_left = Inches(0)
        insights_frame.margin_right = Inches(0.2)
        insights_frame.margin_top = Inches(0)
        insights_frame.margin_bottom = Inches(0)
        insights_frame.word_wrap = True

        # Style the insights text
        insights_p = insights_frame.paragraphs[0]
        insights_p.alignment = PP_ALIGN.LEFT
        insights_p.line_spacing = 1.2  # Slightly more spacing for readability

        for run in insights_p.runs:
            run.font.name = DEFAULT_FONT_FAMILY
            run.font.size = Pt(12)
            run.font.color.rgb = RGBColor(255, 255, 255)  # White text on blue background
            run.font.bold = False

    def _add_image_placeholder(self, slide):
        """Add a circular placeholder for team/fan image"""
        # Create circular placeholder on the right side
        placeholder_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(7.167), Inches(1.5),
            Inches(5), Inches(5)  # Circular shape
        )

        # Style the placeholder - RGBColor only takes RGB values, no alpha
        placeholder_circle.fill.solid()
        placeholder_circle.fill.fore_color.rgb = RGBColor(200, 200, 255)  # Light blue instead of transparent
        placeholder_circle.line.color.rgb = RGBColor(255, 255, 255)
        placeholder_circle.line.width = Pt(3)

        # Add placeholder text inside the circle
        placeholder_text = slide.shapes.add_textbox(
            Inches(8.67), Inches(3.5),
            Inches(2), Inches(1)
        )
        placeholder_frame = placeholder_text.text_frame
        placeholder_frame.text = "Team Image\nPlaceholder"
        placeholder_frame.margin_left = Inches(0)
        placeholder_frame.margin_right = Inches(0)
        placeholder_frame.margin_top = Inches(0)
        placeholder_frame.margin_bottom = Inches(0)

        # Style placeholder text - no alpha channel support
        for paragraph in placeholder_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER

            for run in paragraph.runs:
                run.font.name = DEFAULT_FONT_FAMILY
                run.font.size = Pt(12)
                run.font.color.rgb = RGBColor(100, 100, 150)  # Darker blue instead of transparent
                run.font.italic = True