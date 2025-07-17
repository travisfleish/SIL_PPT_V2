# slide_generators/behaviors_slide.py
"""
Generate Fan Behaviors slide for PowerPoint presentations
Combines fan wheel and community index chart with insights
UPDATED with 6.5" community chart and centered text alignment
ENHANCED with AI-powered insight generation and dynamic font sizing
"""

from pathlib import Path
from typing import Dict, Optional, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import logging
import os
import pandas as pd
from dotenv import load_dotenv
from openai import OpenAI

from .base_slide import BaseSlide
from data_processors.merchant_ranker import MerchantRanker
from visualizations.fan_wheel import FanWheel  # Updated to use enhanced fan wheel
from visualizations.community_index_chart import CommunityIndexChart

# Load environment variables
load_dotenv()

logger = logging.getLogger(__name__)


class BehaviorsSlide(BaseSlide):
    """Generate the Fan Behaviors slide with fan wheel and community index chart"""

    def __init__(self, presentation: Presentation = None, use_ai_insights: bool = True):
        """
        Initialize behaviors slide generator

        Args:
            presentation: Existing presentation to add slide to (creates new if None)
            use_ai_insights: Whether to use AI for insight generation
        """
        super().__init__(presentation)
        self.use_ai_insights = use_ai_insights and bool(os.getenv('OPENAI_API_KEY'))

        if use_ai_insights and not os.getenv('OPENAI_API_KEY'):
            logger.warning("AI insights requested but no OpenAI API key found. Using template insights.")

        # Initialize OpenAI client if using AI
        if self.use_ai_insights:
            try:
                self.openai_client = OpenAI()
                logger.info("OpenAI client initialized for behavior insights")
            except Exception as e:
                logger.error(f"Failed to initialize OpenAI client: {e}")
                self.use_ai_insights = False

    def generate(self,
                 merchant_ranker: MerchantRanker,
                 team_config: Dict[str, Any],
                 slide_index: Optional[int] = None) -> Presentation:
        """
        Generate the complete behaviors slide

        Args:
            merchant_ranker: MerchantRanker instance with data
            team_config: Team configuration including colors and names
            slide_index: Where to insert slide (None = append)

        Returns:
            Updated presentation object
        """
        # Extract team info
        team_name = team_config.get('team_name', 'Team')
        team_short = team_config.get('team_name_short', team_name.split()[-1])
        colors = team_config.get('colors', {})

        # Create visualizations
        logger.info("Generating fan wheel visualization with logo support...")
        fan_wheel_path = self._create_fan_wheel(merchant_ranker, team_config)

        logger.info("Generating community index chart...")
        chart_path = self._create_community_chart(merchant_ranker, colors)

        # Use the content layout (SIL white layout #12)
        slide = self.add_content_slide()
        logger.info("Added behaviors slide using SIL white layout")

        # Add header
        self._add_header(slide, team_name)

        # Generate insight text
        insight = self._generate_insight_text(merchant_ranker, team_short)

        # Add elements with 6.5" chart coordinated positioning
        self._add_insight_text(slide, insight)  # TOP left - large text
        self._add_chart_titles(slide, team_name)  # Titles for both sides
        self._add_community_chart(slide, chart_path)  # LEFT chart - 6.5" wide
        self._add_fan_wheel(slide, fan_wheel_path)  # RIGHT wheel - 5.5" diameter
        self._add_chart_explanation(slide)  # BOTTOM left explanation

        logger.info(f"Generated behaviors slide for {team_name}")
        return self.presentation

    def _create_fan_wheel(self, merchant_ranker: MerchantRanker,
                          team_config: Dict[str, Any]) -> Path:
        """Create fan wheel visualization with logo support"""
        # Get data
        wheel_data = merchant_ranker.get_fan_wheel_data(
            min_audience_pct=0.20,
            top_n_communities=10
        )

        if wheel_data.empty:
            raise ValueError("No fan wheel data available")

        # Create visualization with logo support enabled
        fan_wheel = FanWheel(team_config, enable_logos=True)

        # Generate logo report for debugging
        logo_report = fan_wheel.generate_logo_report(wheel_data)
        logger.info(f"Logo coverage: {logo_report['with_logos']}/{logo_report['total_merchants']} "
                    f"({logo_report['coverage_percentage']:.1f}%)")

        if logo_report['missing_list']:
            logger.debug(f"Missing logos for: {', '.join(logo_report['missing_list'])}")

        output_path = Path('temp_fan_wheel.png')
        fan_wheel.create(wheel_data, output_path)

        return output_path

    def _create_community_chart(self, merchant_ranker: MerchantRanker,
                                team_colors: Dict[str, str]) -> Path:
        """Create community index chart"""
        # Get data with COMPOSITE_INDEX
        communities_df = merchant_ranker.get_top_communities(
            min_audience_pct=0.20,
            top_n=10
        )

        # Rename columns
        data = communities_df.rename(columns={
            'COMMUNITY': 'Community',
            'PERC_AUDIENCE': 'Audience_Pct',
            'COMPOSITE_INDEX': 'Composite_Index'
        })

        # Create chart
        chart = CommunityIndexChart(team_colors)
        output_path = Path('temp_community_chart.png')
        chart.create(data, output_path)

        return output_path

    def _add_header(self, slide, team_name: str):
        """Add header with team name and slide title"""
        # Header background (adjusted for 16:9)
        header_rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(13.333), Inches(0.5)  # Full 16:9 width
        )
        header_rect.fill.solid()
        header_rect.fill.fore_color.rgb = RGBColor(240, 240, 240)
        header_rect.line.color.rgb = RGBColor(200, 200, 200)
        header_rect.line.width = Pt(0.5)

        # Team name (left)
        team_text = slide.shapes.add_textbox(
            Inches(0.2), Inches(0.1),
            Inches(3), Inches(0.3)
        )
        team_text.text_frame.text = team_name
        p = team_text.text_frame.paragraphs[0]
        p.font.name = self.default_font  # Red Hat Display
        p.font.size = Pt(14)
        p.font.bold = True

        # Slide indicator (right)
        slide_text = slide.shapes.add_textbox(
            Inches(6.5), Inches(0.1),  # Adjusted for 16:9
            Inches(6.633), Inches(0.3)  # Adjusted width
        )
        slide_text.text_frame.text = f"Fan Behaviors: How Are {team_name} Fans Unique"
        p = slide_text.text_frame.paragraphs[0]
        p.font.name = self.default_font  # Red Hat Display
        p.alignment = PP_ALIGN.RIGHT
        p.font.size = Pt(14)

    def _add_insight_text(self, slide, insight: str):
        """Add large insight text at TOP - CENTERED with 6.5" chart"""
        # Determine font size based on text length
        insight_length = len(insight)
        if insight_length <= 100:
            font_size = Pt(18)  # Large for short insights
        elif insight_length <= 140:
            font_size = Pt(16)  # Standard size
        elif insight_length <= 160:
            font_size = Pt(14)  # Smaller for longer insights
        else:
            font_size = Pt(12)  # Minimum size for very long insights

        logger.info(f"Insight length: {insight_length} chars, using font size: {font_size.pt}pt")

        text_box = slide.shapes.add_textbox(
            Inches(0.4), Inches(1.0),  # Moved up from 1.2" to 0.8"
            Inches(6.5), Inches(0.8)  # Reduced height to contain text
        )
        text_box.text_frame.text = insight
        text_box.text_frame.word_wrap = True

        # Format text - dynamic size, bold, and CENTER aligned
        p = text_box.text_frame.paragraphs[0]
        p.font.name = self.default_font  # Red Hat Display
        p.font.size = font_size  # Dynamic font size
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER  # CENTER aligned
        p.line_spacing = 1.0  # Tighter line spacing

    def _add_chart_titles(self, slide, team_name: str):
        """Add titles above both visualizations"""
        # Community chart title (LEFT side) - centered over 6.5" chart
        bar_chart_title = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.0),  # Match chart left position
            Inches(6.5), Inches(0.3)  # Match chart width
        )
        bar_chart_title.text_frame.text = f"Top Ten {team_name} Fan Communities"
        p = bar_chart_title.text_frame.paragraphs[0]
        p.font.name = self.default_font
        p.font.size = Pt(12)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Fan wheel title (RIGHT side) - ADJUSTED for new spacing
        title_center = 10.1665
        title_width = 5.5
        title_left = title_center - (title_width / 2)  # 7.4165"

        fan_wheel_title = slide.shapes.add_textbox(
            Inches(title_left), Inches(0.95),
            Inches(title_width), Inches(0.3)
        )
        fan_wheel_title.text_frame.text = "Top Community Fan Purchases"
        p = fan_wheel_title.text_frame.paragraphs[0]
        p.font.name = self.default_font
        p.font.size = Pt(12)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

    def _add_community_chart(self, slide, image_path: Path):
        """Add community index chart - LEFT side, 6.5" WIDTH"""
        left = Inches(0.4)  # Left margin
        top = Inches(2.4)  # Below title
        width = Inches(6.5)  # 6.5" width

        slide.shapes.add_picture(str(image_path), left, top, width=width)

    def _add_fan_wheel(self, slide, image_path: Path):
        """Add fan wheel - RIGHT side, 5.5" diameter"""
        width = Inches(5.5)  # 5.5" diameter
        left = Inches(7.4165)
        top = Inches(1.35)  # Vertical position

        slide.shapes.add_picture(str(image_path), left, top, width=width)

    def _add_chart_explanation(self, slide):
        """Add explanation text below community chart - CENTERED with 6.5" chart"""
        explanation_box = slide.shapes.add_textbox(
            Inches(0.3), Inches(6.4),  # Match chart left position
            Inches(7.2), Inches(0.8)  # Match chart width
        )
        explanation_text = (
            "The top ten fan communities are ranked according to a composite index score "
            "of likelihood to purchase, likelihood to make more purchases per fan versus "
            "the local general population, and likelihood to spend more per fan. The chart "
            "on the right illustrates the top brands (by % fans purchasing) within each of "
            "the top ten communities."
        )
        explanation_box.text_frame.text = explanation_text
        explanation_box.text_frame.word_wrap = True

        p = explanation_box.text_frame.paragraphs[0]
        p.font.name = self.default_font
        p.font.size = Pt(10)
        p.alignment = PP_ALIGN.CENTER  # CENTER aligned
        p.line_spacing = 1.2

    def _generate_insight_text(self, merchant_ranker: MerchantRanker,
                               team_short: str) -> str:
        """Generate insight text based on top communities using AI or fallback"""
        try:
            # Get top communities with their data
            communities = merchant_ranker.get_top_communities(
                min_audience_pct=0.20,
                top_n=10  # Get more for better AI context
            )

            if communities.empty:
                return self._generate_fallback_insight(team_short)

            # Check if AI insights are available
            if self.use_ai_insights and hasattr(self, 'openai_client'):
                return self._generate_ai_behavior_insight(communities, team_short)
            else:
                return self._generate_template_behavior_insight(communities, team_short)

        except Exception as e:
            logger.warning(f"Error generating insight text: {e}")
            return self._generate_fallback_insight(team_short)

    def _generate_ai_behavior_insight(self, communities_df: pd.DataFrame,
                                      team_short: str) -> str:
        """Generate AI-powered behavior insights using OpenAI"""
        try:
            # Prepare data summary for AI
            data_summary = []
            for idx, row in communities_df.head(10).iterrows():
                community = row['COMMUNITY']
                audience_pct = row['PERC_AUDIENCE'] * 100
                composite_idx = row['COMPOSITE_INDEX']

                # Add context about what makes this community special
                data_summary.append(
                    f"{community}: {audience_pct:.1f}% of fans, "
                    f"Composite Index: {composite_idx:.0f}"
                )

            # Create prompt with emphasis on brevity
            prompt = f"""Analyze the top fan communities for {team_short} fans and create a single, engaging sentence that captures their unique behavioral characteristics.

Top Fan Communities (ranked by composite index):
{chr(10).join(data_summary)}

Requirements:
- Write exactly ONE sentence (not a paragraph)
- Try to keep under 130 characters if possible (but quality is more important than length)
- Focus on the top 2-3 most distinctive communities
- Transform community names into engaging descriptive phrases
- Use natural, conversational language
- Make it exciting and actionable for sponsors
- Include an exclamation point for energy
- Avoid listing community names directly; instead describe the behaviors they represent

Example transformations:
- "Live Entertainment Streaming" → "entertainment enthusiasts"
- "Cost Conscious Families" → "value-seeking families"
- "Travel Adventures" → "adventure seekers"
- "Brand Conscious Shoppers" → "brand loyalists"

Good examples:
- "{team_short} fans are entertainment lovers who seek great deals and adventure!" 
- "{team_short} fans are tech-savvy adventurers who love live events!"
"""

            response = self.openai_client.chat.completions.create(
                model="gpt-4",
                messages=[
                    {"role": "system",
                     "content": "You are a marketing analyst who creates compelling behavioral profiles for sports fan bases. Transform data about fan communities into exciting, sponsor-friendly insights."},
                    {"role": "user", "content": prompt}
                ],
                max_tokens=75,
                temperature=0.3  # Lower temperature for more consistent output
            )

            ai_insight = response.choices[0].message.content.strip()

            # Validate and clean response
            if ai_insight and len(ai_insight) > 20:
                ai_insight = ai_insight.strip('"').strip("'")
                if not ai_insight.endswith('!'):
                    ai_insight = ai_insight.rstrip('.') + '!'

                logger.info(f"Generated AI behavior insight ({len(ai_insight)} chars)")
                return ai_insight
            else:
                logger.warning("AI insight too short, using template")
                return self._generate_template_behavior_insight(communities_df, team_short)

        except Exception as e:
            logger.error(f"Error generating AI behavior insights: {e}")
            return self._generate_template_behavior_insight(communities_df, team_short)

    def _generate_template_behavior_insight(self, communities_df: pd.DataFrame,
                                            team_short: str) -> str:
        """Generate template-based behavior insight (fallback)"""
        # Extract community names and convert to insights
        community_names = communities_df['COMMUNITY'].tolist()
        insights = []

        for community in community_names[:2]:  # Use top 2
            if 'entertainment' in community.lower() or 'live' in community.lower():
                insights.append("values-driven live entertainment seekers")
            elif 'cost' in community.lower() or 'conscious' in community.lower():
                insights.append("on the lookout for a deal")
            elif 'travel' in community.lower():
                insights.append("adventure-seeking travelers")
            elif 'beauty' in community.lower():
                insights.append("beauty enthusiasts")
            elif 'brand' in community.lower():
                insights.append("brand-conscious shoppers")
            elif 'movie' in community.lower() or 'film' in community.lower():
                insights.append("movie buffs")
            elif 'game' in community.lower():
                insights.append("gaming enthusiasts")
            elif 'sports' in community.lower():
                insights.append("sports fans")
            elif 'pet' in community.lower():
                insights.append("pet lovers")
            elif 'streaming' in community.lower():
                insights.append("streaming content consumers")

        # Create insight text
        if len(insights) >= 2:
            # Add movie reference for entertainment seekers
            if "entertainment" in insights[0]:
                return f"{team_short} fans are {insights[0]} who are {insights[1]} and a good movie!"
            else:
                return f"{team_short} fans are {insights[0]} who are {insights[1]}!"
        elif len(insights) == 1:
            if "entertainment" in insights[0]:
                return f"{team_short} fans are {insights[0]} who love a good movie!"
            else:
                return f"{team_short} fans are {insights[0]}!"

        return self._generate_fallback_insight(team_short)

    def _generate_fallback_insight(self, team_short: str) -> str:
        """Ultimate fallback insight"""
        return f"{team_short} fans have unique behaviors that set them apart from the general population!"


# Convenience function
def create_behaviors_slide(merchant_ranker: MerchantRanker,
                           team_config: Dict[str, Any],
                           presentation: Optional[Presentation] = None,
                           use_ai_insights: bool = True) -> Presentation:
    """
    Create a behaviors slide for the presentation

    Args:
        merchant_ranker: Data source
        team_config: Team configuration
        presentation: Existing presentation (creates new if None)
        use_ai_insights: Whether to use AI for insight generation

    Returns:
        Presentation with behaviors slide added
    """
    generator = BehaviorsSlide(presentation, use_ai_insights)
    return generator.generate(merchant_ranker, team_config)