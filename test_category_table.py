#!/usr/bin/env python3
"""
Simple test script to quickly iterate on category table formatting
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path


def create_test_slide():
    """Create a single slide with just the category tables for testing"""

    # Create presentation (16:9 format)
    pres = Presentation()
    pres.slide_width = Inches(13.333)
    pres.slide_height = Inches(7.5)

    # Add blank slide
    slide = pres.slides.add_slide(pres.slide_layouts[6])

    # Add a simple title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(5), Inches(0.5))
    title_box.text = "Travel Sponsor Analysis"
    p = title_box.text_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True

    # Create Category Table (top table)
    create_category_table(slide)

    # Create Subcategory Table (bottom table)
    create_subcategory_table(slide)

    # Save
    pres.save('test_category_tables.pptx')
    print("✅ Created test_category_tables.pptx")


def create_category_table(slide):
    """Create the category metrics table"""
    # Position and size
    left = Inches(5.8)
    top = Inches(1.5)
    width = Inches(6.8)
    height = Inches(0.8)

    # Create table
    table_shape = slide.shapes.add_table(2, 4, left, top, width, height)
    table = table_shape.table

    # Column widths
    table.columns[0].width = Inches(1.3)
    table.columns[1].width = Inches(1.6)
    table.columns[2].width = Inches(2.0)
    table.columns[3].width = Inches(1.9)

    # Headers
    headers = ['Category', 'Percent of Fans\nWho Spend', 'How likely fans are to\nspend vs. gen pop',
               'How many more purchases\nper fan v gen pop']

    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        format_header_cell(cell)

    # Data row
    table.cell(1, 0).text = "Travel"
    table.cell(1, 1).text = "52.7%"
    table.cell(1, 2).text = "97% More"
    table.cell(1, 3).text = "55% more"

    # Format data cells
    for i in range(4):
        format_data_cell(table.cell(1, i))


def create_subcategory_table(slide):
    """Create the subcategory table"""
    # Position and size
    left = Inches(5.8)
    top = Inches(2.7)
    width = Inches(6.8)
    rows = 5  # 4 subcategories + header

    # Create table
    table_shape = slide.shapes.add_table(rows, 4, left, top, width, Inches(0.3 * rows))
    table = table_shape.table

    # Column widths
    table.columns[0].width = Inches(1.6)
    table.columns[1].width = Inches(1.5)
    table.columns[2].width = Inches(1.9)
    table.columns[3].width = Inches(1.8)

    # Headers
    headers = ['Sub-Category', 'Percent of Fans\nWho Spend', 'How likely fans are to\nspend vs. gen pop',
               'How many more purchases\nper fan v gen pop']

    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        format_header_cell(cell)

    # Data rows
    data = [
        ('Airlines', '31.7%', '84% More', '73% more'),
        ('Travel Agencies', '18.2%', '82% More', '31% more'),
        ('Rental Cars', '15.3%', '104% More', '33% more'),
        ('Cruises', '4.0%', '65% More', '4% more')
    ]

    for row_idx, (subcat, pct, likely, purchases) in enumerate(data, 1):
        table.cell(row_idx, 0).text = subcat
        table.cell(row_idx, 1).text = pct
        table.cell(row_idx, 2).text = likely
        table.cell(row_idx, 3).text = purchases

        # Format cells
        for col in range(4):
            format_data_cell(table.cell(row_idx, col))


def format_header_cell(cell):
    """Format header cell"""
    # Background color
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(217, 217, 217)  # Light gray

    # Text formatting
    text_frame = cell.text_frame
    text_frame.margin_left = Inches(0.05)
    text_frame.margin_right = Inches(0.05)
    text_frame.margin_top = Inches(0.03)
    text_frame.margin_bottom = Inches(0.03)
    text_frame.word_wrap = True

    # Format all paragraphs
    for paragraph in text_frame.paragraphs:
        paragraph.font.name = "Arial"  # Change to "Red Hat Display" if available
        paragraph.font.size = Pt(10)
        paragraph.font.bold = True
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.line_spacing = 1.0


def format_data_cell(cell):
    """Format data cell"""
    text_frame = cell.text_frame
    text_frame.margin_left = Inches(0.05)
    text_frame.margin_right = Inches(0.05)
    text_frame.margin_top = Inches(0.03)
    text_frame.margin_bottom = Inches(0.03)
    text_frame.word_wrap = True

    # Format all paragraphs
    for paragraph in text_frame.paragraphs:
        paragraph.font.name = "Arial"  # Change to "Red Hat Display" if available
        paragraph.font.size = Pt(11)
        paragraph.alignment = PP_ALIGN.CENTER

        # Color coding
        if 'More' in cell.text:
            paragraph.font.color.rgb = RGBColor(0, 176, 80)  # Green
        elif 'Less' in cell.text:
            paragraph.font.color.rgb = RGBColor(255, 0, 0)  # Red


if __name__ == "__main__":
    # Quick edit zone - modify these values to test different settings

    # You can quickly test different positions here:
    # LEFT_POSITION = 5.8  # Change this to move tables left/right
    # TABLE_WIDTH = 6.8    # Change this to make tables wider/narrower

    create_test_slide()

    print("\nQuick tips for editing:")
    print("1. Change table positions in create_category_table() and create_subcategory_table()")
    print("2. Adjust column widths by modifying table.columns[x].width values")
    print("3. Change font sizes in format_header_cell() and format_data_cell()")
    print("4. Test different colors by modifying RGBColor values")